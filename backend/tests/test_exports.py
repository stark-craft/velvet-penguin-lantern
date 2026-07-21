from __future__ import annotations

import csv
import io
import json
import os
import unittest
import zipfile
from pathlib import Path
from tempfile import TemporaryDirectory

from docx import Document
from pptx import Presentation
from pydantic import ValidationError

from signalroom.models import (
    ArticleCreate,
    ArticleSourceCreate,
    DiscoveryMethod,
    ProfileId,
    make_stable_id,
)
from signalroom.services.exports import (
    ExportArticleNotFoundError,
    ExportFormat,
    ExportRequest,
    ExportService,
)
from signalroom.storage import Repository


class ExportServiceTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary_directory = TemporaryDirectory()
        self.repository = Repository(Path(self.temporary_directory.name) / "signalroom.db")
        self.service = ExportService(self.repository)
        self.retained = self._article(
            suffix="retained",
            profile=ProfileId.DEFAULT,
            title="=newsScrapper launch changes the market",
            retained=True,
        )
        self.rejected = self._article(
            suffix="rejected",
            profile=ProfileId.DEFAULT,
            title="A rejected signal",
            retained=False,
        )
        self.broadcast = self._article(
            suffix="broadcast",
            profile=ProfileId.BROADCAST,
            title="Broadcast standards update",
            retained=True,
        )

    def tearDown(self) -> None:
        self.temporary_directory.cleanup()

    def _article(self, *, suffix, profile, title, retained):
        url = f"https://news.example/{suffix}"
        source_url = f"https://wire.example/{suffix}?ref=signalroom"
        source = ArticleSourceCreate(
            stable_id=make_stable_id("source", profile.value, suffix),
            profile=profile,
            source_key="wire",
            publisher="Example Wire",
            url=source_url,
            canonical_url=url,
            discovery_method=DiscoveryMethod.RSS,
        )
        return self.repository.upsert_article(
            ArticleCreate(
                stable_id=make_stable_id("article", url),
                title=title,
                canonical_url=url,
                published_at="2026-07-20T07:15:00+05:30",
                summary=f"Summary for {suffix} with a clear strategic consequence.",
                intent="opportunity" if retained else "noise",
                top_image_url=f"https://images.example/{suffix}.jpg",
                region="India",
                category="AI & Technology",
                importance_score=0.91 if retained else 0.2,
                keywords=("agent governance", "enterprise AI"),
                profiles=(profile,),
                sources=(source,),
                model_metadata={
                    "gatekeeper": {
                        "decision": "keep" if retained else "drop",
                        "keep": retained,
                    }
                },
                metadata={
                    "author": "Maya Rao",
                    "insight": "Governance is now a product requirement.",
                    "retained": retained,
                },
            )
        )

    def _request(self, export_format: ExportFormat, **changes) -> ExportRequest:
        values = {
            "profile": ProfileId.DEFAULT,
            "article_ids": (self.retained.id,),
            "format": export_format,
            "filename": "../../Board Briefing.exe",
        }
        values.update(changes)
        return ExportRequest(**values)

    def test_json_contract_is_ordered_scoped_and_deterministic(self) -> None:
        request = self._request(
            ExportFormat.JSON,
            article_ids=(self.retained.id, self.rejected.id),
        )
        first = self.service.generate(request)
        second = self.service.generate(request)
        payload = json.loads(first.content)

        self.assertEqual(first.content, second.content)
        self.assertEqual(first.filename, "Board-Briefing.json")
        self.assertEqual(first.media_type, "application/json")
        self.assertEqual(first.article_count, 1)
        self.assertEqual(payload["profile"], "default")
        self.assertEqual(payload["article_count"], 1)
        self.assertEqual(payload["articles"][0]["article_id"], str(self.retained.id))
        self.assertEqual(payload["articles"][0]["image_url"], str(self.retained.top_image_url))
        self.assertEqual(
            payload["articles"][0]["sources"][0]["url"],
            str(self.retained.sources[0].url),
        )
        self.assertNotIn(str(self.rejected.id), first.content.decode("utf-8"))

    def test_inclusion_flags_remove_optional_sections_and_can_include_rejected(self) -> None:
        result = self.service.generate(
            self._request(
                ExportFormat.JSON,
                article_ids=(self.rejected.id,),
                include_rejected_items=True,
                include_images=False,
                include_summaries=False,
                include_source_links=False,
                include_opinions=False,
                include_metadata=False,
            )
        )
        article = json.loads(result.content)["articles"][0]
        self.assertEqual(result.article_count, 1)
        self.assertEqual(
            set(article),
            {"article_id", "stable_id", "title", "profile", "published_at"},
        )

    def test_csv_is_utf8_flat_and_guards_against_formula_injection(self) -> None:
        result = self.service.generate(self._request(ExportFormat.CSV))
        rows = list(csv.DictReader(io.StringIO(result.content.decode("utf-8-sig"))))
        self.assertEqual(result.media_type, "text/csv; charset=utf-8")
        self.assertEqual(rows[0]["title"], "'=newsScrapper launch changes the market")
        self.assertIn("https://wire.example/retained", rows[0]["source_urls"])
        self.assertEqual(rows[0]["keywords"], '["agent governance","enterprise AI"]')

    def test_office_exports_are_valid_in_memory_packages_with_text_urls(self) -> None:
        for export_format in (ExportFormat.XLSX, ExportFormat.DOCX, ExportFormat.PPTX):
            with self.subTest(export_format=export_format.value):
                request = self._request(export_format)
                first = self.service.generate(request)
                second = self.service.generate(request)
                self.assertTrue(first.content.startswith(b"PK"))
                self.assertEqual(first.content, second.content)
                with zipfile.ZipFile(io.BytesIO(first.content)) as package:
                    self.assertIn("[Content_Types].xml", package.namelist())

                if export_format == ExportFormat.XLSX:
                    with zipfile.ZipFile(io.BytesIO(first.content)) as package:
                        xml = b"\n".join(
                            package.read(name)
                            for name in package.namelist()
                            if name.endswith(".xml")
                        ).decode("utf-8", errors="ignore")
                    self.assertIn("newsScrapper", xml)
                    self.assertIn("https://wire.example/retained", xml)
                elif export_format == ExportFormat.DOCX:
                    document = Document(io.BytesIO(first.content))
                    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
                    self.assertIn(self.retained.title, text)
                    self.assertIn("https://wire.example/retained", text)
                    self.assertIn("Image URL", text)
                else:
                    presentation = Presentation(io.BytesIO(first.content))
                    text = "\n".join(
                        shape.text
                        for slide in presentation.slides
                        for shape in slide.shapes
                        if hasattr(shape, "text")
                    )
                    self.assertEqual(len(presentation.slides), 2)
                    self.assertIn(self.retained.title, text)
                    self.assertIn("https://wire.example/retained", text)
                    self.assertIn("IMAGE URL", text)

    def test_root_style_legacy_pptx_template_is_used_when_configured(self) -> None:
        template_path = Path(self.temporary_directory.name) / "template.pptx"
        template = Presentation()
        signature_slide = template.slides.add_slide(template.slide_layouts[6])
        signature_slide.shapes.add_textbox(0, 0, 100, 100).text = "LEGACY TEMPLATE SIGNATURE"
        template.save(template_path)
        previous = os.environ.get("SIGNALROOM_PPTX_TEMPLATE")
        os.environ["SIGNALROOM_PPTX_TEMPLATE"] = str(template_path)
        try:
            result = self.service.generate(self._request(ExportFormat.PPTX))
        finally:
            if previous is None:
                os.environ.pop("SIGNALROOM_PPTX_TEMPLATE", None)
            else:
                os.environ["SIGNALROOM_PPTX_TEMPLATE"] = previous
        presentation = Presentation(io.BytesIO(result.content))
        text = "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        self.assertEqual(len(presentation.slides), 3)
        self.assertIn("LEGACY TEMPLATE SIGNATURE", text)

    def test_profile_membership_is_enforced_for_every_explicit_id(self) -> None:
        with self.assertRaises(ExportArticleNotFoundError):
            self.service.generate(
                self._request(
                    ExportFormat.JSON,
                    article_ids=(self.broadcast.id,),
                )
            )

    def test_request_bounds_and_duplicate_ids_are_rejected(self) -> None:
        with self.assertRaises(ValidationError):
            self._request(
                ExportFormat.JSON,
                article_ids=(self.retained.id, self.retained.id),
            )
        with self.assertRaises(ValidationError):
            self._request(
                ExportFormat.JSON,
                article_ids=tuple(self.retained.id for _ in range(101)),
            )


if __name__ == "__main__":
    unittest.main()
