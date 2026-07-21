"""Deterministic, in-memory exports for an explicit set of articles.

The first API version deliberately accepts article UUIDs instead of a query or
an unbounded "all results" scope.  This keeps authorization and resource use
easy to reason about: the caller chooses at most 100 records, and every record
must belong to the requested profile.

No renderer performs network I/O.  Image and source URLs are exported as text,
so the generated files remain reproducible and cannot be used as an SSRF
primitive.
"""

from __future__ import annotations

import csv
import io
import json
import os
import re
import unicodedata
import zipfile
from dataclasses import dataclass
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Mapping, Optional, Sequence, Tuple
from uuid import UUID

import xlsxwriter
from docx import Document
from docx.enum.text import WD_BREAK
from docx.shared import Inches as DocxInches
from docx.shared import Pt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt as PptxPt
from pydantic import Field, StringConstraints, field_validator
from typing_extensions import Annotated

from signalroom.branding import EXPORT_TITLE, PRODUCT_NAME
from signalroom.models import ArticleRead, ProfileId, StrictModel
from signalroom.storage import RecordNotFoundError, SQLiteRepository


class ExportFormat(str, Enum):
    JSON = "json"
    CSV = "csv"
    XLSX = "xlsx"
    DOCX = "docx"
    PPTX = "pptx"


class ExportRequest(StrictModel):
    """Validated service input used directly by ``POST /api/v1/exports``."""

    profile: ProfileId
    article_ids: Annotated[Tuple[UUID, ...], Field(min_length=1, max_length=100)]
    format: ExportFormat
    filename: Optional[
        Annotated[str, StringConstraints(strip_whitespace=True, min_length=1, max_length=180)]
    ] = None
    include_images: bool = True
    include_summaries: bool = True
    include_source_links: bool = True
    include_opinions: bool = True
    include_metadata: bool = True
    include_rejected_items: bool = False

    @field_validator("article_ids")
    @classmethod
    def article_ids_are_unique(cls, values: Tuple[UUID, ...]) -> Tuple[UUID, ...]:
        if len(values) != len(set(values)):
            raise ValueError("article_ids must be unique")
        return values


@dataclass(frozen=True)
class ExportResult:
    content: bytes
    media_type: str
    filename: str
    article_count: int


class ExportError(RuntimeError):
    """Base error for export generation failures."""


class ExportArticleNotFoundError(ExportError):
    """An explicit article ID does not exist in the requested profile."""


_MEDIA_TYPES = {
    ExportFormat.JSON: "application/json",
    ExportFormat.CSV: "text/csv; charset=utf-8",
    ExportFormat.XLSX: "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ExportFormat.DOCX: "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ExportFormat.PPTX: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

_FIXED_DOCUMENT_TIME = datetime(2000, 1, 1, 0, 0, 0)
_ZIP_TIME = (1980, 1, 1, 0, 0, 0)
_WINDOWS_RESERVED_NAMES = {
    "CON",
    "PRN",
    "AUX",
    "NUL",
    *(f"COM{number}" for number in range(1, 10)),
    *(f"LPT{number}" for number in range(1, 10)),
}


def sanitize_filename(
    requested: Optional[str], *, profile: ProfileId, export_format: ExportFormat
) -> str:
    """Return a path-free filename with exactly the requested format extension."""

    raw = requested or f"signalroom-{profile.value}-export"
    raw = unicodedata.normalize("NFKC", raw).replace("\\", "/").rsplit("/", 1)[-1]
    if "." in raw:
        raw = raw.rsplit(".", 1)[0]
    raw = re.sub(r"[\x00-\x1f\x7f]", "", raw)
    raw = re.sub(r"[^A-Za-z0-9 _.-]+", "-", raw)
    raw = re.sub(r"[\s._-]+", "-", raw).strip("-.")
    stem = raw[:110].rstrip("-.") or f"signalroom-{profile.value}-export"
    if stem.upper() in _WINDOWS_RESERVED_NAMES:
        stem = f"signalroom-{stem}"
    return f"{stem}.{export_format.value}"


def _iso(value: Any) -> Optional[str]:
    if value is None:
        return None
    if hasattr(value, "isoformat"):
        return str(value.isoformat())
    return str(value)


def _is_rejected(article: ArticleRead) -> bool:
    retained = article.metadata.get("retained")
    if retained is False:
        return True
    gatekeeper = article.model_metadata.get("gatekeeper")
    if isinstance(gatekeeper, Mapping):
        if gatekeeper.get("keep") is False:
            return True
        if str(gatekeeper.get("decision") or "").casefold() in {"drop", "dropped", "reject"}:
            return True
    return str(article.metadata.get("status") or "").casefold() in {"rejected", "dropped"}


def _opinion(article: ArticleRead) -> Dict[str, Optional[str]]:
    insight = article.metadata.get("insight") or article.metadata.get("why_it_matters")
    return {
        "intent": article.intent,
        "insight": str(insight) if insight is not None else None,
    }


def _source_record(source: Any) -> Dict[str, Any]:
    return {
        "publisher": source.publisher,
        "source_key": source.source_key,
        "url": str(source.url),
        "canonical_url": str(source.canonical_url) if source.canonical_url else None,
        "published_at": _iso(source.published_at),
        "discovery_method": source.discovery_method.value,
    }


def _article_record(article: ArticleRead, request: ExportRequest) -> Dict[str, Any]:
    record: Dict[str, Any] = {
        "article_id": str(article.id),
        "stable_id": article.stable_id,
        "title": article.title,
        "profile": request.profile.value,
        "published_at": _iso(article.published_at),
    }
    if request.include_summaries:
        record["summary"] = article.summary
    if request.include_opinions:
        record.update(_opinion(article))
    if request.include_images:
        record["image_url"] = str(article.top_image_url) if article.top_image_url else None
    if request.include_source_links:
        record["article_url"] = str(article.canonical_url)
        record["sources"] = [_source_record(source) for source in article.sources]
    if request.include_metadata:
        record.update(
            {
                "region": article.region,
                "category": article.category,
                "language": article.language,
                "importance_score": article.importance_score,
                "keywords": list(article.keywords),
                "author": article.metadata.get("author"),
                "created_at": _iso(article.created_at),
                "model_metadata": article.model_metadata,
                "metadata": article.metadata,
            }
        )
    return record


def _inclusion_manifest(request: ExportRequest) -> Dict[str, bool]:
    return {
        "images": request.include_images,
        "summaries": request.include_summaries,
        "source_links": request.include_source_links,
        "opinions": request.include_opinions,
        "metadata": request.include_metadata,
        "rejected_items": request.include_rejected_items,
    }


def _json_text(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False, separators=(",", ":"), sort_keys=True)


def _flat_columns(request: ExportRequest) -> List[str]:
    columns = ["article_id", "stable_id", "title", "profile", "published_at"]
    if request.include_summaries:
        columns.append("summary")
    if request.include_opinions:
        columns.extend(["intent", "insight"])
    if request.include_images:
        columns.append("image_url")
    if request.include_source_links:
        columns.extend(["article_url", "source_publishers", "source_urls"])
    if request.include_metadata:
        columns.extend(
            [
                "region",
                "category",
                "language",
                "importance_score",
                "keywords",
                "author",
                "created_at",
                "model_metadata",
                "metadata",
            ]
        )
    return columns


def _flat_record(record: Mapping[str, Any], columns: Sequence[str]) -> Dict[str, Any]:
    flattened = {key: record.get(key) for key in columns}
    sources = record.get("sources") or []
    if "source_publishers" in columns:
        flattened["source_publishers"] = " | ".join(
            str(source.get("publisher") or "") for source in sources
        )
    if "source_urls" in columns:
        flattened["source_urls"] = " | ".join(
            str(source.get("url") or source.get("canonical_url") or "") for source in sources
        )
    for key, value in tuple(flattened.items()):
        if isinstance(value, (dict, list, tuple)):
            flattened[key] = _json_text(value)
        elif value is None:
            flattened[key] = ""
    return flattened


def _csv_safe(value: Any) -> Any:
    """Prevent CSV cells from becoming spreadsheet formulas when opened."""

    if isinstance(value, str) and value.startswith(("=", "+", "-", "@")):
        return "'" + value
    return value


def _render_json(
    request: ExportRequest, records: Sequence[Mapping[str, Any]]
) -> bytes:
    payload = {
        "schema_version": 1,
        "profile": request.profile.value,
        "article_count": len(records),
        "inclusions": _inclusion_manifest(request),
        "articles": list(records),
    }
    return (json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True) + "\n").encode(
        "utf-8"
    )


def _render_csv(request: ExportRequest, records: Sequence[Mapping[str, Any]]) -> bytes:
    columns = _flat_columns(request)
    stream = io.StringIO(newline="")
    writer = csv.DictWriter(stream, fieldnames=columns, lineterminator="\n")
    writer.writeheader()
    for record in records:
        writer.writerow(
            {key: _csv_safe(value) for key, value in _flat_record(record, columns).items()}
        )
    return stream.getvalue().encode("utf-8-sig")


def _xlsx_cell(value: Any) -> Any:
    if value is None:
        return ""
    if isinstance(value, (dict, list, tuple)):
        value = _json_text(value)
    if isinstance(value, str) and len(value) > 32_000:
        return value[:31_997] + "..."
    return value


def _render_xlsx(request: ExportRequest, records: Sequence[Mapping[str, Any]]) -> bytes:
    output = io.BytesIO()
    workbook = xlsxwriter.Workbook(
        output,
        {
            "in_memory": True,
            "strings_to_formulas": False,
            "strings_to_urls": False,
        },
    )
    workbook.set_properties(
        {
            "title": EXPORT_TITLE,
            "subject": f"{request.profile.value.title()} profile intelligence",
            "author": PRODUCT_NAME,
            "company": PRODUCT_NAME,
            "created": _FIXED_DOCUMENT_TIME,
        }
    )
    accent = "6957E8" if request.profile == ProfileId.DEFAULT else "D9475C"
    header = workbook.add_format(
        {
            "bold": True,
            "font_color": "FFFFFF",
            "bg_color": accent,
            "border": 0,
            "text_wrap": True,
            "valign": "vcenter",
        }
    )
    section = workbook.add_format({"bold": True, "font_color": accent})
    wrap = workbook.add_format({"text_wrap": True, "valign": "top"})
    number = workbook.add_format({"num_format": "0%", "valign": "top"})

    overview = workbook.add_worksheet("Overview")
    overview.hide_gridlines(2)
    overview.set_column("A:A", 24)
    overview.set_column("B:B", 38)
    overview.write("A1", EXPORT_TITLE, header)
    overview.write("A3", "Profile", section)
    overview.write("B3", request.profile.value.title())
    overview.write("A4", "Articles", section)
    overview.write_number("B4", len(records))
    overview.write("A6", "Included fields", section)
    for row_index, (name, enabled) in enumerate(_inclusion_manifest(request).items(), start=6):
        overview.write(row_index, 0, name.replace("_", " ").title())
        overview.write(row_index, 1, "Yes" if enabled else "No")

    columns = _flat_columns(request)
    articles = workbook.add_worksheet("Articles")
    articles.hide_gridlines(2)
    articles.freeze_panes(1, 0)
    articles.set_row(0, 28)
    articles.write_row(0, 0, [column.replace("_", " ").title() for column in columns], header)
    for row_index, record in enumerate(records, start=1):
        flattened = _flat_record(record, columns)
        for column_index, key in enumerate(columns):
            value = _xlsx_cell(flattened.get(key))
            cell_format = number if key == "importance_score" and isinstance(value, float) else wrap
            articles.write(row_index, column_index, value, cell_format)
    if records:
        articles.autofilter(0, 0, len(records), len(columns) - 1)
    for column_index, column in enumerate(columns):
        width = 44 if column in {"title", "summary", "metadata", "model_metadata"} else 22
        articles.set_column(column_index, column_index, width)

    if request.include_source_links:
        sources_sheet = workbook.add_worksheet("Sources")
        sources_sheet.hide_gridlines(2)
        source_columns = [
            "article_id",
            "article_title",
            "publisher",
            "source_key",
            "url",
            "canonical_url",
            "published_at",
            "discovery_method",
        ]
        sources_sheet.write_row(
            0, 0, [column.replace("_", " ").title() for column in source_columns], header
        )
        source_row = 1
        for record in records:
            for source in record.get("sources") or []:
                values = [record.get("article_id"), record.get("title")]
                values.extend(source.get(column) for column in source_columns[2:])
                sources_sheet.write_row(
                    source_row, 0, [_xlsx_cell(value) for value in values], wrap
                )
                source_row += 1
        sources_sheet.freeze_panes(1, 0)
        sources_sheet.set_column(0, 1, 38)
        sources_sheet.set_column(2, 3, 20)
        sources_sheet.set_column(4, 5, 56)
        sources_sheet.set_column(6, 7, 22)

    workbook.close()
    return _canonicalize_zip(output.getvalue())


def _set_docx_style(document: Document) -> None:
    normal = document.styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10)
    for section in document.sections:
        section.top_margin = DocxInches(0.65)
        section.bottom_margin = DocxInches(0.65)
        section.left_margin = DocxInches(0.75)
        section.right_margin = DocxInches(0.75)


def _docx_label(document: Document, label: str, value: Any) -> None:
    if value in (None, "", [], {}):
        return
    paragraph = document.add_paragraph()
    label_run = paragraph.add_run(f"{label}: ")
    label_run.bold = True
    paragraph.add_run(str(value))


def _render_docx(request: ExportRequest, records: Sequence[Mapping[str, Any]]) -> bytes:
    document = Document()
    _set_docx_style(document)
    document.core_properties.title = EXPORT_TITLE
    document.core_properties.subject = f"{request.profile.value.title()} profile intelligence"
    document.core_properties.author = PRODUCT_NAME
    document.core_properties.last_modified_by = PRODUCT_NAME
    document.core_properties.created = _FIXED_DOCUMENT_TIME
    document.core_properties.modified = _FIXED_DOCUMENT_TIME

    title = document.add_heading(EXPORT_TITLE, level=0)
    title.runs[0].font.color.rgb = DocxRGBColor(52, 43, 118)
    document.add_paragraph(
        f"{request.profile.value.title()} profile · {len(records)} article"
        f"{'s' if len(records) != 1 else ''}"
    )
    document.add_paragraph(
        "Images and source references are represented as URLs; this document contains no "
        "remotely fetched media."
    )

    if not records:
        document.add_paragraph("No articles matched the requested export scope.")
    for index, record in enumerate(records, start=1):
        if index > 1:
            document.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
        document.add_heading(f"{index}. {record['title']}", level=1)
        _docx_label(document, "Published", record.get("published_at") or "Not supplied")
        if request.include_summaries and record.get("summary"):
            document.add_heading("AI summary", level=2)
            document.add_paragraph(str(record["summary"]))
        if request.include_opinions:
            document.add_heading("Editorial interpretation", level=2)
            _docx_label(document, "Intent", record.get("intent"))
            _docx_label(document, "Why it matters", record.get("insight"))
        if request.include_images:
            _docx_label(document, "Image URL", record.get("image_url") or "Not supplied")
        if request.include_source_links:
            document.add_heading("Source links", level=2)
            _docx_label(document, "Canonical article", record.get("article_url"))
            for source in record.get("sources") or []:
                document.add_paragraph(
                    f"{source.get('publisher') or 'Unknown source'} — "
                    f"{source.get('url') or source.get('canonical_url') or ''}",
                    style="List Bullet",
                )
        if request.include_metadata:
            document.add_heading("Metadata", level=2)
            for label, key in (
                ("Article ID", "article_id"),
                ("Category", "category"),
                ("Region", "region"),
                ("Language", "language"),
                ("Importance", "importance_score"),
                ("Keywords", "keywords"),
                ("Author", "author"),
            ):
                value = record.get(key)
                if isinstance(value, list):
                    value = ", ".join(str(item) for item in value)
                _docx_label(document, label, value)

    output = io.BytesIO()
    document.save(output)
    return _canonicalize_zip(output.getvalue())


def _pptx_text(
    slide: Any,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: float,
    color: str,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = "Aptos"
    run.font.size = PptxPt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def _pptx_background(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def _pptx_template_path() -> Optional[Path]:
    """Resolve the optional legacy template without depending on cwd.

    Windows operators may place ``template.pptx`` in the repository root as
    the legacy application expected. An explicit environment override remains
    available for deployments that keep mutable assets outside the checkout.
    """

    configured = str(os.environ.get("SIGNALROOM_PPTX_TEMPLATE") or "").strip()
    project_root = Path(__file__).resolve().parents[3]
    candidates = (
        Path(configured).expanduser() if configured else None,
        project_root / "template.pptx",
        project_root / "backend" / "template.pptx",
    )
    for candidate in candidates:
        if candidate is not None and candidate.is_file():
            return candidate.resolve()
    return None


def _legacy_layout(presentation: Presentation, name: str, fallback_index: int) -> Any:
    return next(
        (layout for layout in presentation.slide_layouts if layout.name == name),
        presentation.slide_layouts[min(fallback_index, len(presentation.slide_layouts) - 1)],
    )


def _legacy_placeholder_map(layout: Any) -> Dict[str, int]:
    markers = {
        "#TITLE": "title",
        "#SUMMARY": "summary",
        "#LINK": "link",
        "#INSIGHT": "insight",
        "#DATE_HERE": "date",
        "#TARGATED_SRID_TEAM": "team",
        "#TARGETED_SRID_TEAM": "team",
    }
    result: Dict[str, int] = {}
    for shape in layout.placeholders:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type == PP_PLACEHOLDER.PICTURE:
            result["picture"] = shape.placeholder_format.idx
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text.strip().upper()
        for marker, key in markers.items():
            if marker in text:
                result[key] = shape.placeholder_format.idx
                break
    return result


def _legacy_set_text(shape: Any, value: str, *, size: Optional[int] = None) -> None:
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = value
    if size is not None:
        paragraph.font.size = PptxPt(size)


def _render_template_pptx(
    request: ExportRequest,
    records: Sequence[Mapping[str, Any]],
    template_path: Path,
) -> bytes:
    """Render the marker-based contract from the supplied legacy main.py."""

    presentation = Presentation(str(template_path))
    cover_layout = _legacy_layout(presentation, "CoverLayout", 0)
    news_layout = _legacy_layout(presentation, "NewsLayout", 1)
    placeholder_map = _legacy_placeholder_map(news_layout)
    month_label = datetime.now().strftime("%b'%y")

    cover = presentation.slides.add_slide(cover_layout)
    for shape in cover.shapes:
        if shape.has_text_frame and "#DATE_HERE" in shape.text.upper():
            _legacy_set_text(shape, month_label)

    for record in records:
        slide = presentation.slides.add_slide(news_layout)
        for shape in slide.placeholders:
            index = shape.placeholder_format.idx
            if index == placeholder_map.get("title"):
                _legacy_set_text(shape, str(record.get("title") or "Untitled signal"))
            elif index == placeholder_map.get("summary"):
                summary = str(record.get("summary") or "Summary unavailable.")
                frame = shape.text_frame
                frame.clear()
                sentences = [item.strip() for item in re.split(r"(?<=[.!?])\s+", summary) if item.strip()]
                for sentence_index, sentence in enumerate(sentences or [summary]):
                    paragraph = frame.paragraphs[0] if sentence_index == 0 else frame.add_paragraph()
                    paragraph.text = sentence
                    paragraph.level = 0
                    paragraph.font.name = "Calibri"
                    paragraph.font.size = PptxPt(18)
            elif index == placeholder_map.get("link"):
                _legacy_set_text(shape, str(record.get("article_url") or ""), size=10)
            elif index == placeholder_map.get("insight"):
                insight = str(record.get("insight") or record.get("intent") or "Unavailable")
                _legacy_set_text(shape, f"Insight : {insight}", size=14)
            elif index == placeholder_map.get("date"):
                _legacy_set_text(shape, month_label)
            elif index == placeholder_map.get("team"):
                metadata = record.get("metadata") if isinstance(record.get("metadata"), Mapping) else {}
                team = str(metadata.get("team") or "ALL")
                _legacy_set_text(shape, f"Targeted SRID TEAM : {team}")
            elif index == placeholder_map.get("picture") and request.include_images:
                # Remote image fetching remains deliberately disabled in the
                # export renderer. Preserve the template's picture frame and
                # attach the source URL as its click target when one exists.
                image_url = str(record.get("image_url") or "").strip()
                if image_url:
                    try:
                        shape.click_action.hyperlink.address = image_url
                    except (AttributeError, ValueError):
                        pass

    presentation.core_properties.title = EXPORT_TITLE
    presentation.core_properties.subject = f"{request.profile.value.title()} profile intelligence"
    presentation.core_properties.author = PRODUCT_NAME
    presentation.core_properties.last_modified_by = PRODUCT_NAME
    output = io.BytesIO()
    presentation.save(output)
    return _canonicalize_zip(output.getvalue())


def _render_pptx(request: ExportRequest, records: Sequence[Mapping[str, Any]]) -> bytes:
    template_path = _pptx_template_path()
    if template_path is not None:
        return _render_template_pptx(request, records, template_path)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = EXPORT_TITLE
    presentation.core_properties.subject = f"{request.profile.value.title()} profile intelligence"
    presentation.core_properties.author = PRODUCT_NAME
    presentation.core_properties.last_modified_by = PRODUCT_NAME
    presentation.core_properties.created = _FIXED_DOCUMENT_TIME
    presentation.core_properties.modified = _FIXED_DOCUMENT_TIME

    accent = "B9FF66" if request.profile == ProfileId.DEFAULT else "FFB35C"
    accent_dark = "6957E8" if request.profile == ProfileId.DEFAULT else "D9475C"
    ink = "101820"
    paper = "F6F7F9"
    muted = "68737D"

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _pptx_background(title_slide, ink)
    accent_bar = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(0.72), Inches(0.12), Inches(5.9)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor.from_string(accent)
    accent_bar.line.fill.background()
    _pptx_text(
        title_slide,
        f"{PRODUCT_NAME.upper()} · INTELLIGENCE EXPORT",
        left=1.15,
        top=1.1,
        width=10.8,
        height=0.45,
        size=14,
        color=accent,
        bold=True,
    )
    _pptx_text(
        title_slide,
        f"{request.profile.value.title()}\nbriefing",
        left=1.1,
        top=1.72,
        width=10.9,
        height=2.3,
        size=42,
        color="FFFFFF",
        bold=True,
    )
    _pptx_text(
        title_slide,
        f"{len(records)} curated article{'s' if len(records) != 1 else ''} · "
        "source and image URLs included as text",
        left=1.15,
        top=5.75,
        width=10.7,
        height=0.55,
        size=14,
        color="AAB4BC",
    )

    for index, record in enumerate(records, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _pptx_background(slide, paper)
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, Inches(0.16)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = RGBColor.from_string(accent_dark)
        top_bar.line.fill.background()
        context = " · ".join(
            value
            for value in (
                request.profile.value.upper(),
                str(record.get("category") or ""),
                str(record.get("region") or ""),
            )
            if value
        )
        _pptx_text(
            slide,
            context,
            left=0.7,
            top=0.48,
            width=10.8,
            height=0.35,
            size=11,
            color=accent_dark,
            bold=True,
        )
        _pptx_text(
            slide,
            str(record["title"]),
            left=0.68,
            top=0.9,
            width=11.7,
            height=1.18,
            size=26,
            color=ink,
            bold=True,
        )
        summary_lines: List[str] = []
        if request.include_summaries and record.get("summary"):
            summary_lines.append("AI SUMMARY\n" + str(record["summary"]))
        if request.include_opinions and (record.get("intent") or record.get("insight")):
            opinion = " · ".join(
                value
                for value in (
                    f"Intent: {record.get('intent')}" if record.get("intent") else "",
                    str(record.get("insight") or ""),
                )
                if value
            )
            summary_lines.append("EDITORIAL INTERPRETATION\n" + opinion)
        if not summary_lines:
            summary_lines.append("Article selected for intelligence review.")
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.68),
            Inches(2.25),
            Inches(8.15),
            Inches(3.8),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
        panel.line.color.rgb = RGBColor.from_string("E2E5E8")
        _pptx_text(
            slide,
            "\n\n".join(summary_lines),
            left=0.95,
            top=2.5,
            width=7.62,
            height=3.28,
            size=15,
            color=ink,
        )

        rail_lines = [
            f"PUBLISHED\n{record.get('published_at') or 'Not supplied'}",
        ]
        if request.include_metadata:
            rail_lines.extend(
                [
                    f"IMPORTANCE\n{round(float(record.get('importance_score') or 0) * 100)}%",
                    "KEYWORDS\n" + ", ".join(str(item) for item in record.get("keywords") or []),
                ]
            )
        if request.include_images:
            rail_lines.append(f"IMAGE URL\n{record.get('image_url') or 'Not supplied'}")
        if request.include_source_links:
            source_urls = [
                str(source.get("url") or source.get("canonical_url") or "")
                for source in record.get("sources") or []
            ]
            rail_lines.append(
                "SOURCE URLS\n"
                + ("\n".join(source_urls) or str(record.get("article_url") or "Not supplied"))
            )
        _pptx_text(
            slide,
            "\n\n".join(rail_lines),
            left=9.15,
            top=2.28,
            width=3.45,
            height=3.82,
            size=11,
            color=muted,
        )
        _pptx_text(
            slide,
            f"{index:02d} / {len(records):02d}    {record['article_id']}",
            left=0.7,
            top=6.72,
            width=11.8,
            height=0.28,
            size=9,
            color=muted,
            align=PP_ALIGN.RIGHT,
        )

    output = io.BytesIO()
    presentation.save(output)
    return _canonicalize_zip(output.getvalue())


def _canonicalize_zip(payload: bytes) -> bytes:
    """Normalize OOXML ZIP timestamps and member order for repeatable bytes."""

    source = io.BytesIO(payload)
    target = io.BytesIO()
    with zipfile.ZipFile(source, "r") as archive, zipfile.ZipFile(
        target, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9
    ) as output:
        for name in sorted(archive.namelist()):
            original = archive.getinfo(name)
            info = zipfile.ZipInfo(filename=name, date_time=_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.comment = original.comment
            info.extra = b""
            info.create_system = 0
            info.external_attr = original.external_attr
            output.writestr(info, archive.read(name))
    return target.getvalue()


class ExportService:
    """Load profile-scoped articles and render one in-memory export."""

    def __init__(self, repository: SQLiteRepository) -> None:
        self.repository = repository

    def _load_articles(self, request: ExportRequest) -> Tuple[ArticleRead, ...]:
        articles: List[ArticleRead] = []
        for article_id in request.article_ids:
            try:
                article = self.repository.get_article(article_id, profile=request.profile)
            except RecordNotFoundError as exc:
                raise ExportArticleNotFoundError(
                    f"article {article_id} was not found in profile {request.profile.value}"
                ) from exc
            if request.profile not in article.profiles:
                raise ExportArticleNotFoundError(
                    f"article {article_id} was not found in profile {request.profile.value}"
                )
            if request.include_rejected_items or not _is_rejected(article):
                articles.append(article)
        return tuple(articles)

    def generate(self, request: ExportRequest) -> ExportResult:
        articles = self._load_articles(request)
        records = tuple(_article_record(article, request) for article in articles)
        renderers = {
            ExportFormat.JSON: _render_json,
            ExportFormat.CSV: _render_csv,
            ExportFormat.XLSX: _render_xlsx,
            ExportFormat.DOCX: _render_docx,
            ExportFormat.PPTX: _render_pptx,
        }
        content = renderers[request.format](request, records)
        return ExportResult(
            content=content,
            media_type=_MEDIA_TYPES[request.format],
            filename=sanitize_filename(
                request.filename,
                profile=request.profile,
                export_format=request.format,
            ),
            article_count=len(records),
        )


__all__ = [
    "ExportArticleNotFoundError",
    "ExportError",
    "ExportFormat",
    "ExportRequest",
    "ExportResult",
    "ExportService",
    "sanitize_filename",
]
