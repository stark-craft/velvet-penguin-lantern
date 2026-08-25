from fastapi import FastAPI, Query, HTTPException, Body, Request, Response
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from contextlib import asynccontextmanager
from apscheduler.schedulers.background import BackgroundScheduler
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
# --- PPTX IMPORTS ---
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
# --- DOCX IMPORTS ---
from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
# --- EXCEL IMPORTS ---
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
# --- IMAGE & AI IMPORTS ---
from PIL import Image, ImageOps
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from typing import List, Optional
import subprocess
import os
import json
import sys
import datetime
import time
import glob
import pandas as pd
import io
import requests
import queue
from news_scrapper import learner
import re
import pickle
import numpy as np
import threading
import secrets
import platform
import hashlib
import ipaddress
import socket
from pathlib import Path
from urllib.parse import urljoin, urlsplit, urlunsplit
from dotenv import load_dotenv
from sentence_transformers import SentenceTransformer
from newspaper import Article
from fastapi import BackgroundTasks
from concurrent.futures import ThreadPoolExecutor
from threading import Semaphore
from core.secure_http import tls_verify
from core.storage import JsonStore
from news_scrapper.personalization import PersonalizationService
from news_scrapper.recommendation.identity import bind_viewer_request, set_viewer_cookie
from news_scrapper.source_catalog import (
    build_shadow_briefing,
    canonical_url as canonical_source_url,
    load_sites,
    normalize_site,
    write_shadow_catalog,
)
from core.profile import client_ip as resolve_client_ip
from core.profile import normalize_ip
from core.profile import resolve_profile
from core.settings import (
    FRONTEND_DIST,
    MODEL_ROOT,
    NEWS_CONFIG_DIR,
    NEWS_CRAWLER_DIR,
    NEWS_RUNTIME_DIR,
    PROJECT_ROOT,
    ensure_runtime_directories,
    migrate_legacy_news_runtime,
    migrate_unified_news_state,
    model_path as resolve_model_path,
)


# Load the local backend configuration before models, adapters, or security
# settings read environment variables.  Existing process environment values win
# over .env so production launchers can override this file safely.
BASE_DIR = PROJECT_ROOT
load_dotenv(PROJECT_ROOT / ".env", override=False)
ensure_runtime_directories()
migrate_legacy_news_runtime()
UNIFIED_MIGRATION_REPORT = migrate_unified_news_state()


def env_bool(name: str, default: bool = False) -> bool:
    value = os.environ.get(name)
    if value is None:
        return default
    return value.strip().lower() in {"1", "true", "yes", "on"}


def env_csv(name: str, default: str = "") -> set[str]:
    return {
        value.strip()
        for value in os.environ.get(name, default).split(",")
        if value.strip()
    }


def env_ip_set(name: str, default: str = "") -> set[str]:
    """Read a comma-separated IP allowlist using one canonical form."""

    return {normalize_ip(value) for value in env_csv(name, default)}


def env_ip_map(name: str, default: str = "") -> dict[str, str]:
    result = {}
    for entry in os.environ.get(name, default).split(","):
        if ":" not in entry:
            continue
        address, label = entry.split(":", 1)
        if address.strip() and label.strip():
            result[normalize_ip(address)] = label.strip()
    return result


try:
    from news_scrapper.adapters.samsung_web_search import (
        check_samsung_web_search,
        enrich_article_with_web_search,
    )
except Exception as adapter_error:
    check_samsung_web_search = None
    enrich_article_with_web_search = None
    print(f"[ADAPTER] Samsung Web Search unavailable: {adapter_error}", flush=True)

try:
    from news_scrapper.adapters.samsung_chat import (
        check_samsung_chat,
        summarize_article_with_chat,
    )
except Exception as adapter_error:
    check_samsung_chat = None
    summarize_article_with_chat = None
    print(f"[ADAPTER] Samsung Chat unavailable: {adapter_error}", flush=True)

try:
    from news_scrapper.adapters.article_metadata import enrich_article_image_metadata
except Exception as adapter_error:
    enrich_article_image_metadata = None
    print(f"[ADAPTER] Article metadata unavailable: {adapter_error}", flush=True)

# ==========================================
# AI GATEKEEPER MODEL LOADING
# ==========================================
print("Waking up the AI Gatekeeper...")

BOUNCER_MODEL_FILENAMES = {
    "default": "bouncer_model.pkl",
}

try:
    current_dir = str(NEWS_RUNTIME_DIR)
    model_folder = str(resolve_model_path("all-MiniLM-L6-v2", "local_miniLM_model"))
    bouncer_embedder = SentenceTransformer(model_folder)
    bouncer_models = {}

    for profile_name, model_filename in BOUNCER_MODEL_FILENAMES.items():
        model_path = os.path.join(current_dir, model_filename)
        if not os.path.exists(model_path):
            print(f"No {profile_name} bouncer found yet: {model_filename}")
            continue
        try:
            with open(model_path, "rb") as f:
                bouncer_models[profile_name] = pickle.load(f)
            print(f"Loaded {profile_name} bouncer: {model_filename}")
        except Exception as model_error:
            print(
                f"Could not load {profile_name} bouncer "
                f"({model_filename}): {model_error}"
            )

    bouncer_model = bouncer_models.get("default")
    if bouncer_models:
        print("AI Gatekeeper is awake with the unified model.")
    else:
        print("No bouncer models loaded. Scanning without filter.")
except Exception as e:
    print(f"Warning: Gatekeeper not found. Scanning without filter. Error: {e}")
    bouncer_embedder = None
    bouncer_models = {}
    bouncer_model = None

# ==========================================
# CONFIGURATION
# ==========================================
MORNING_KEYWORDS = (
    "OpenAI , Robot , Samsung , LG , Sony , Nvidia , TCL , OLED , QNED , "
    "Artificial Intelligence, chatGPT , Anthropic , Claude , Gemini , LED , "
    "Robotics , Television , TV , display , Grok , GPU , Processor , Jio , TPU"
)
BROADCAST_MORNING_KEYWORDS = (
    "DTH, Cable TV, IPTV, Broadcast, Digital terrestrial transmission, DTT, "
    "DVB S, DVB S2, DVB C, DVB C2, DVB T, DVB T2, conditional access system, "
    "digital rights management, FAST, OTT, Connected TV, Tuner, Set top box, "
    "Linear ad insertion, Linear ads, TRAI, MIB, broadcast regulation, HBB TV, "
    "DVB I, 5G broadcast, D2M"
)
DIRECTOR_KEY = os.environ.get("DIRECTOR_KEY", "1357")
ANALYTICS_KEY = os.environ.get("ANALYTICS_KEY", DIRECTOR_KEY)
GATEKEEPER_KEY = os.environ.get("GATEKEEPER_KEY", DIRECTOR_KEY)
APP_ENV = os.environ.get("NEWSSCRAPPER_ENV", "development").strip().lower()
ROOT_DIR = str(NEWS_RUNTIME_DIR)
HISTORY_DIR = os.path.join(ROOT_DIR, "history_archive")
MANUAL_LOG_FILE = os.path.join(ROOT_DIR, "manual_search_logs.xlsx")
WORKFLOW_FILE = os.path.join(ROOT_DIR, "workflow_store.json")
TRAINING_FILE = os.path.join(ROOT_DIR, "trainingData.json")
NOT_INTERESTED_FILE = os.path.join(ROOT_DIR, "not_interested_store.json")
NOT_INTERESTED_EXPIRY_HOURS = 22
USAGE_TRACKER_FILE = os.path.join(ROOT_DIR, "usage_tracker.json")
VIEWER_PROFILES_FILE = os.path.join(ROOT_DIR, "viewer_profiles.json")
VIEWER_HIDDEN_FILE = os.path.join(ROOT_DIR, "viewer_hidden_store.json")
VIEWER_SAVED_FILE = os.path.join(ROOT_DIR, "viewer_saved_store.json")
VIEWER_BRIEFING_FILE = os.path.join(ROOT_DIR, "viewer_url_briefings.json")
VIEWER_IDENTITY_CLAIMS_FILE = os.path.join(ROOT_DIR, "viewer_identity_claims.json")
VIEWER_PERSONALIZATION_FILE = os.path.join(ROOT_DIR, "viewer_personalization.json")
IP_HASH_SECRET = os.environ.get("NEWSSCRAPPER_IP_HASH_SECRET", "development-only-change-this-secret")
if APP_ENV in {"production", "prod"}:
    if IP_HASH_SECRET == "development-only-change-this-secret":
        raise RuntimeError("NEWSSCRAPPER_IP_HASH_SECRET must be changed in production.")
    if len(DIRECTOR_KEY) < 6 or len(ANALYTICS_KEY) < 6 or len(GATEKEEPER_KEY) < 6:
        raise RuntimeError("Production approval, analytics, and Gatekeeper keys must contain at least six characters.")
SCRAPY_ROBOTSTXT_OBEY = env_bool("SCRAPY_ROBOTSTXT_OBEY", True)
SCHEDULER_ENABLED = env_bool("SCHEDULER_ENABLED", True)
UNIFIED_CORPUS_SHADOW_ENABLED = env_bool("UNIFIED_CORPUS_SHADOW_ENABLED", False)
UNIFIED_CORPUS_ENABLED = env_bool("UNIFIED_CORPUS_ENABLED", True)
LEGACY_PROFILE_ROUTING_ENABLED = env_bool("LEGACY_PROFILE_ROUTING_ENABLED", False)
HISTORY_RETENTION_DAYS = max(1, int(os.environ.get("HISTORY_RETENTION_DAYS", "30")))
CRAWL_LOOKBACK_DAYS = max(1, int(os.environ.get("CRAWL_LOOKBACK_DAYS", "1")))
WEB_SEARCH_ENRICHMENT_ENABLED = env_bool("WEB_SEARCH_ENRICHMENT_ENABLED", False)
WEB_SEARCH_REQUIRE_SUCCESS = env_bool("WEB_SEARCH_REQUIRE_SUCCESS", False)
WEB_SEARCH_REQUIRE_KEYWORD_MATCH = env_bool(
    "WEB_SEARCH_REQUIRE_KEYWORD_MATCH", True
)
WEB_SEARCH_MAX_ENRICH_PER_RUN = max(0, int(os.environ.get("WEB_SEARCH_MAX_ENRICH_PER_RUN", "0")))
WEB_SEARCH_ENRICH_DELAY_SECONDS = max(0.0, float(os.environ.get("WEB_SEARCH_ENRICH_DELAY_SECONDS", "0")))
FINAL_CHAT_SUMMARY_ENABLED = env_bool("FINAL_CHAT_SUMMARY_ENABLED", False)
FINAL_CHAT_SUMMARY_DELAY_SECONDS = max(0.0, float(os.environ.get("FINAL_CHAT_SUMMARY_DELAY_SECONDS", "0")))
FINAL_CHAT_SUMMARY_MAX_ARTICLES = max(0, int(os.environ.get("FINAL_CHAT_SUMMARY_MAX_ARTICLES", "0")))
SAMSUNG_PIPELINE_ENABLED = env_bool(
    "SAMSUNG_PIPELINE_ENABLED",
    WEB_SEARCH_ENRICHMENT_ENABLED and FINAL_CHAT_SUMMARY_ENABLED,
)
if SAMSUNG_PIPELINE_ENABLED:
    WEB_SEARCH_ENRICHMENT_ENABLED = True
    FINAL_CHAT_SUMMARY_ENABLED = True
    WEB_SEARCH_REQUIRE_SUCCESS = True
SAMSUNG_DISCOVERY_ONLY = SAMSUNG_PIPELINE_ENABLED and env_bool(
    "SAMSUNG_DISCOVERY_ONLY", True
)
SAMSUNG_PIPELINE_CREDENTIALS_READY = all(
    os.environ.get(name, "").strip()
    for name in (
        "SAMSUNG_WEB_SEARCH_CLIENT",
        "SAMSUNG_WEB_SEARCH_TOKEN",
        "SAMSUNG_CHAT_CLIENT",
        "SAMSUNG_CHAT_TOKEN",
        "SAMSUNG_CHAT_MODEL_ID",
    )
)
if SAMSUNG_PIPELINE_ENABLED and not SAMSUNG_PIPELINE_CREDENTIALS_READY:
    print(
        "[PIPELINE] Samsung pipeline is enabled but one or more Web Search/"
        "Chat credentials are missing. Existing briefings remain available, "
        "but new Samsung extraction cannot complete until configuration is fixed.",
        flush=True,
    )
SAMSUNG_CACHE_MAX_ITEMS = max(
    100, int(os.environ.get("SAMSUNG_CACHE_MAX_ITEMS", "2500"))
)
SAMSUNG_HEALTH_CACHE_SECONDS = max(
    0, int(os.environ.get("SAMSUNG_HEALTH_CACHE_SECONDS", "900"))
)
SAMSUNG_PIPELINE_CACHE_DIR = NEWS_RUNTIME_DIR / "samsung_pipeline_cache"
WEB_SEARCH_CACHE = JsonStore(
    SAMSUNG_PIPELINE_CACHE_DIR / "web_search_success.json", dict
)
CHAT_SUMMARY_CACHE = JsonStore(
    SAMSUNG_PIPELINE_CACHE_DIR / "chat_summary_success.json", dict
)
pipeline_health_lock = threading.Lock()
pipeline_health_cache = {"checked_at": 0.0, "result": None}


def resolve_pipeline_capabilities(force=False):
    """Choose external or local stages once, before a crawler is launched."""

    if not SAMSUNG_PIPELINE_ENABLED:
        return {
            "web_search": False,
            "chat": False,
            "discovery_only": False,
            "summary_engine": "local_bart",
            "mode": "local_scrapy_bart",
        }

    now = time.monotonic()
    with pipeline_health_lock:
        cached = pipeline_health_cache.get("result")
        age = now - float(pipeline_health_cache.get("checked_at") or 0)
        if (
            not force
            and isinstance(cached, dict)
            and age <= SAMSUNG_HEALTH_CACHE_SECONDS
        ):
            print(
                f"[PIPELINE:PRECHECK] Reusing {round(age)}s-old capability "
                f"result: {cached['mode']}.",
                flush=True,
            )
            return dict(cached)

        print(
            "[PIPELINE:PRECHECK] Testing Samsung Web Search and Samsung Chat "
            "before crawling starts...",
            flush=True,
        )
        web_result = (
            check_samsung_web_search()
            if check_samsung_web_search is not None
            else {"available": False, "error": "Web Search adapter import failed"}
        )
        chat_result = (
            check_samsung_chat()
            if check_samsung_chat is not None
            else {"available": False, "error": "Chat adapter import failed"}
        )
        web_available = bool(web_result.get("available"))
        chat_available = bool(chat_result.get("available"))

        if web_available:
            print(
                f"[PIPELINE:PRECHECK] PASS Samsung Web Search "
                f"({web_result.get('latency_ms', '?')} ms). Scrapy will discover "
                "URLs only; exact Web Search references and targeted completion "
                "will supply article text.",
                flush=True,
            )
        else:
            print(
                "[PIPELINE:PRECHECK] FAIL Samsung Web Search: "
                f"{web_result.get('error', 'unknown failure')}. FALLBACK: Scrapy "
                "will crawl and extract complete article text in this run.",
                flush=True,
            )
        if chat_available:
            print(
                f"[PIPELINE:PRECHECK] PASS Samsung Chat "
                f"({chat_result.get('latency_ms', '?')} ms). Chat will generate "
                "the lead, bullets, intent, and Why This Matters.",
                flush=True,
            )
        else:
            print(
                "[PIPELINE:PRECHECK] FAIL Samsung Chat: "
                f"{chat_result.get('error', 'unknown failure')}. FALLBACK: local "
                "BART will summarize and local FLAN-T5 will serve strategic insight.",
                flush=True,
            )

        if web_available and chat_available:
            mode = "samsung_web_search_and_chat"
        elif web_available:
            mode = "samsung_web_search_local_bart"
        elif chat_available:
            mode = "scrapy_extraction_samsung_chat"
        else:
            mode = "local_scrapy_bart"
        result = {
            "web_search": web_available,
            "chat": chat_available,
            "discovery_only": web_available and SAMSUNG_DISCOVERY_ONLY,
            "summary_engine": "samsung_chat" if chat_available else "local_bart",
            "mode": mode,
            "web_search_check": web_result,
            "chat_check": chat_result,
        }
        pipeline_health_cache.update({"checked_at": now, "result": dict(result)})
        print(f"[PIPELINE:PRECHECK] Selected run mode: {mode}.", flush=True)
        return result

# ==========================================
# PROFILE ROUTING AND STORAGE
# ==========================================
DEFAULT_PROFILE = "default"
BROADCAST_PROFILE = "broadcast"
UNIFIED_PROFILE = DEFAULT_PROFILE
BROADCAST_SPECIAL_IPS = env_ip_set(
    "BROADCAST_SPECIAL_IPS",
    "",
)
ANALYTICS_ALLOWED_IPS = env_ip_set(
    "ANALYTICS_ALLOWED_IPS",
    "127.0.0.1,::1",
)
GATEKEEPER_ALLOWED_IPS = env_ip_set(
    "GATEKEEPER_ALLOWED_IPS",
    "127.0.0.1,::1",
)
PROFILE_SETTINGS_ALLOWED_IPS = env_ip_set(
    "PROFILE_SETTINGS_ALLOWED_IPS",
    "127.0.0.1,::1",
)
SYSTEM_STATUS_ALLOWED_IPS = env_ip_set(
    "SYSTEM_STATUS_ALLOWED_IPS",
    "127.0.0.1,::1",
)
TRUSTED_PROXY_IPS = env_ip_set("TRUSTED_PROXY_IPS", "127.0.0.1,::1")

DEFAULT_SITES_FILE = str(NEWS_CONFIG_DIR / "sites.json")
BROADCAST_SITES_FILE = str(NEWS_CONFIG_DIR / "sites_broadcast.json")
INTELLIGENCE_STORE_DIR = os.path.join(ROOT_DIR, "intelligence_store")
UNIFIED_HISTORY_DIR = os.path.join(INTELLIGENCE_STORE_DIR, "unified", "history")


def merge_keyword_packs(*packs):
    values = []
    for pack in packs:
        values.extend(part.strip() for part in str(pack or "").split(","))
    return ", ".join(dict.fromkeys(value for value in values if value))


UNIFIED_MORNING_KEYWORDS = merge_keyword_packs(
    MORNING_KEYWORDS,
    BROADCAST_MORNING_KEYWORDS,
)

PROFILE_CONFIGS = {
    DEFAULT_PROFILE: {
        "label": "Unified Intelligence",
        "keywords": UNIFIED_MORNING_KEYWORDS,
        "sites_file": DEFAULT_SITES_FILE,
        "history_dir": UNIFIED_HISTORY_DIR,
        "use_bouncer": True,
    },
    # Rollback metadata only. Serving functions normalize this legacy name to
    # DEFAULT_PROFILE while unified mode is active.
    BROADCAST_PROFILE: {
        "label": "Legacy Broadcast Rollback",
        "keywords": BROADCAST_MORNING_KEYWORDS,
        "sites_file": BROADCAST_SITES_FILE,
        "history_dir": os.path.join(INTELLIGENCE_STORE_DIR, BROADCAST_PROFILE, "history"),
        "use_bouncer": True,
    },
}
WORKFLOW_FILES = {
    DEFAULT_PROFILE: WORKFLOW_FILE,
    BROADCAST_PROFILE: WORKFLOW_FILE,
}
NOT_INTERESTED_FILES = {
    DEFAULT_PROFILE: NOT_INTERESTED_FILE,
    BROADCAST_PROFILE: NOT_INTERESTED_FILE,
}
TRAINING_FILES = {
    DEFAULT_PROFILE: TRAINING_FILE,
    BROADCAST_PROFILE: TRAINING_FILE,
}
BOUNCER_MODEL_FILES = {
    DEFAULT_PROFILE: os.path.join(ROOT_DIR, "bouncer_model.pkl"),
    BROADCAST_PROFILE: os.path.join(ROOT_DIR, "bouncer_model.pkl"),
}
REGION_LEARNING_FILES = {
    DEFAULT_PROFILE: os.path.join(ROOT_DIR, "region_learning.json"),
    BROADCAST_PROFILE: os.path.join(ROOT_DIR, "region_learning.json"),
}

# Personnel/network mappings are deployment data, never source-code defaults.
# Configure them only in the untracked .env file on the internal server.
TEAM_IP_MAP = env_ip_map("TEAM_IP_MAP")


def get_client_ip(request: Request = None):
    if not request:
        return "unknown"
    return resolve_client_ip(
        request.client.host if request.client else "unknown",
        request.headers,
        TRUSTED_PROXY_IPS,
    )


def get_profile_for_request(request: Request = None):
    if UNIFIED_CORPUS_ENABLED and not LEGACY_PROFILE_ROUTING_ENABLED:
        return UNIFIED_PROFILE
    client_ip = get_client_ip(request)
    if request:
        requested_profile = (
            request.headers.get("x-sense-profile")
            or request.query_params.get("profile")
            or ""
        ).strip().lower()
        # Explicit switching is a developer/operations feature. Ordinary users
        # are always assigned from their real client IP.
        if (
            requested_profile in PROFILE_CONFIGS
            and client_ip in PROFILE_SETTINGS_ALLOWED_IPS
        ):
            return requested_profile

    return resolve_profile(
        client_ip,
        BROADCAST_SPECIAL_IPS,
        requested_profile=requested_profile if request else "",
        switch_allowed_ips=PROFILE_SETTINGS_ALLOWED_IPS,
    )


def get_active_profile_name(request: Request = None):
    return get_profile_for_request(request)


def public_profile_name(profile: str) -> str:
    return "unified" if UNIFIED_CORPUS_ENABLED else profile


def get_profile_config(profile: str):
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    return PROFILE_CONFIGS.get(profile, PROFILE_CONFIGS[DEFAULT_PROFILE])


def get_sites_file_for_profile(profile: str):
    return DEFAULT_SITES_FILE if UNIFIED_CORPUS_ENABLED else get_profile_config(profile)["sites_file"]


def get_profile_history_dir(profile: str):
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    history_dir = get_profile_config(profile)["history_dir"]
    os.makedirs(history_dir, exist_ok=True)
    return history_dir


def ensure_profile_storage():
    profiles = [UNIFIED_PROFILE] if UNIFIED_CORPUS_ENABLED else list(PROFILE_CONFIGS)
    for profile in profiles:
        os.makedirs(get_profile_history_dir(profile), exist_ok=True)


def get_profile_history_files(profile: str, include_legacy_default: bool = True):
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    files = glob.glob(os.path.join(get_profile_history_dir(profile), "*.json"))
    if include_legacy_default and profile == DEFAULT_PROFILE and not UNIFIED_CORPUS_ENABLED:
        files.extend(glob.glob(os.path.join(HISTORY_DIR, "*.json")))
    return list(dict.fromkeys(os.path.abspath(file_path) for file_path in files))


def resolve_profile_history_file(filename: str, profile: str):
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    safe_name = Path(filename).name
    profile_path = os.path.join(get_profile_history_dir(profile), safe_name)
    if os.path.exists(profile_path):
        return profile_path
    legacy_path = os.path.join(HISTORY_DIR, safe_name)
    if not UNIFIED_CORPUS_ENABLED and profile == DEFAULT_PROFILE and os.path.exists(legacy_path):
        return legacy_path
    return None


def get_latest_briefing_file_for_profile(profile: str):
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    files = glob.glob(os.path.join(get_profile_history_dir(profile), "briefing_*.json"))
    if not UNIFIED_CORPUS_ENABLED and profile == DEFAULT_PROFILE and not files:
        files = glob.glob(os.path.join(HISTORY_DIR, "briefing_*.json"))
    valid_files = []
    for file_path in files:
        try:
            with open(file_path, "r", encoding="utf-8") as file_obj:
                if json.load(file_obj):
                    valid_files.append(file_path)
        except Exception:
            continue
    return max(valid_files, key=os.path.getmtime) if valid_files else None


def purge_expired_history(profile: str, keep_days: int = HISTORY_RETENTION_DAYS):
    """Delete profile history older than the retention window.

    Workflow, selected, approved, and viewer data live in separate JSON stores,
    so retained user decisions are not removed with an expired briefing file.
    """

    cutoff = datetime.datetime.now().timestamp() - (keep_days * 86400)
    removed = 0
    for file_path in get_profile_history_files(profile, include_legacy_default=False):
        try:
            if os.path.getmtime(file_path) < cutoff:
                os.remove(file_path)
                removed += 1
        except OSError as error:
            print(f"[RETENTION:{profile}] Could not remove {file_path}: {error}", flush=True)
    if removed:
        print(f"[RETENTION:{profile}] Removed {removed} briefing file(s) older than {keep_days} days.", flush=True)
    return removed


def get_workflow_file_for_request(request: Request = None):
    profile = get_profile_for_request(request)
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    return WORKFLOW_FILES.get(profile, WORKFLOW_FILE)


def get_not_interested_file_for_profile(profile: str):
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    return NOT_INTERESTED_FILES.get(profile, NOT_INTERESTED_FILE)


def get_training_file_for_profile(profile: str):
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    return TRAINING_FILES.get(profile, TRAINING_FILE)


def get_bouncer_model_file_for_profile(profile: str):
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    return BOUNCER_MODEL_FILES.get(profile, BOUNCER_MODEL_FILES[DEFAULT_PROFILE])


def get_bouncer_model_for_profile(profile: str):
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    return bouncer_models.get(profile) if bouncer_embedder is not None else None


def get_team_owner_for_ip(ip: str):
    return TEAM_IP_MAP.get(str(ip or "").strip())


def is_analytics_allowed_ip(ip: str) -> bool:
    return str(ip or "").strip() in ANALYTICS_ALLOWED_IPS


def require_analytics_access(request: Request, key: str = None):
    ip = get_client_ip(request)

    if not is_analytics_allowed_ip(ip):
        raise HTTPException(status_code=403, detail="Analytics is not enabled for this network.")

    if key != ANALYTICS_KEY:
        raise HTTPException(status_code=403, detail="Invalid analytics key.")

    return ip


def get_profile_debug_info(profile: str):
    config = get_profile_config(profile)
    latest_file = get_latest_briefing_file_for_profile(profile)
    return {
        "profile": profile,
        "label": config["label"],
        "sites_file": config["sites_file"],
        "sites_file_exists": os.path.exists(config["sites_file"]),
        "history_dir": config["history_dir"],
        "history_dir_exists": os.path.exists(config["history_dir"]),
        "latest_briefing": os.path.basename(latest_file) if latest_file else None,
        "bouncer_model_file": get_bouncer_model_file_for_profile(profile),
        "bouncer_model_exists": os.path.exists(get_bouncer_model_file_for_profile(profile)),
        "training_file": get_training_file_for_profile(profile),
        "training_file_exists": os.path.exists(get_training_file_for_profile(profile)),
        "not_interested_file": get_not_interested_file_for_profile(profile),
        "not_interested_file_exists": os.path.exists(get_not_interested_file_for_profile(profile)),
    }


BOUNCER_LOW_PRIORITY_THRESHOLD = 0.45
BOUNCER_HARD_DROP_THRESHOLD = 0.60

# ==========================================
# --- THREAD POOL FOR ML INFERENCE ---
# ==========================================
ml_executor = ThreadPoolExecutor(max_workers=2, thread_name_prefix="ml_worker")
gatekeeper_executor = ThreadPoolExecutor(max_workers=1, thread_name_prefix="gatekeeper_restore")
personal_briefing_executor = ThreadPoolExecutor(
    max_workers=2,
    thread_name_prefix="personal_briefing",
)

# Feedback can arrive faster than a local bouncer model can retrain.  A single
# worker consumes one request per profile at a time, while coalescing any votes
# that arrive during the current run into exactly one follow-up run.  This
# avoids both concurrent pickle writes and the old "already training; skip"
# data-loss behavior.
training_queue = queue.Queue()
training_queue_lock = threading.Lock()
training_queued_profiles = set()
training_running_profiles = set()
training_dirty_profiles = set()
training_worker = None
training_worker_stop = threading.Event()

# ==========================================
# --- MULTI-USER JOB TRACKING ---
# ==========================================
crawl_semaphore = Semaphore(3)
active_jobs = {}

SCHEDULER_STATUS = {
    "is_active": False,
    "message": "System Ready.",
    "mode": "idle",
}
SCHEDULER_STATE_STORE = JsonStore(
    NEWS_RUNTIME_DIR / "scheduler_state.json",
    lambda: {
        "schema_version": 1,
        "status": "idle",
        "stage": "idle",
        "completed_partitions": [],
        "failures": [],
    },
)


def update_durable_scheduler_state(**changes):
    """Persist scheduler progress atomically for diagnostics and restart audits."""

    def updater(current):
        state = current if isinstance(current, dict) else {}
        return {"schema_version": 1, **state, **changes}

    return SCHEDULER_STATE_STORE.update(updater)


def _load_briefing_articles(path):
    if not path:
        return []
    try:
        with open(path, "r", encoding="utf-8") as handle:
            payload = json.load(handle)
        return payload if isinstance(payload, list) else []
    except (OSError, json.JSONDecodeError):
        return []


def publish_unified_shadow(run_id):
    """Create non-serving parity artifacts after both legacy partitions succeed."""

    catalog_report = write_shadow_catalog(
        Path(DEFAULT_SITES_FILE),
        Path(BROADCAST_SITES_FILE),
        NEWS_RUNTIME_DIR,
    )
    default_items = _load_briefing_articles(
        get_latest_briefing_file_for_profile(DEFAULT_PROFILE)
    )
    broadcast_items = _load_briefing_articles(
        get_latest_briefing_file_for_profile(BROADCAST_PROFILE)
    )
    destination = (
        NEWS_RUNTIME_DIR
        / "unified_shadow"
        / "history"
        / f"briefing_{run_id}.json"
    )
    briefing_report = build_shadow_briefing(
        default_items,
        broadcast_items,
        destination,
        semantic=True,
    )
    report = {
        "schema_version": 1,
        "run_id": run_id,
        "catalog": catalog_report,
        "briefing": briefing_report,
        "serving_enabled": False,
        "generated_at": datetime.datetime.now(datetime.timezone.utc).isoformat(
            timespec="seconds"
        ),
    }
    JsonStore(
        NEWS_RUNTIME_DIR / "unified_shadow" / "latest_parity_report.json",
        dict,
    ).write(report)
    return report

scheduler_lock = threading.Lock()
scheduler_process_lock = threading.Lock()
scheduler_shutdown_event = threading.Event()
scheduler_processes = {}
scheduler_pending_run = False
scheduler_retry_scheduled = False
file_lock = threading.Lock()
train_lock = threading.Lock()
not_interested_lock = threading.Lock()
briefing_lock = threading.RLock()
tracker_lock = threading.Lock()
workflow_lock = threading.RLock()
voc_lock = threading.Lock()
sites_lock = threading.Lock()
viewer_hidden_lock = threading.Lock()
viewer_saved_lock = threading.RLock()
viewer_briefing_lock = threading.RLock()
viewer_identity_lock = threading.RLock()
region_learning_lock = threading.Lock()
dropped_lock = threading.Lock()
opinion_lock = threading.Lock()
insight_cache_lock = threading.Lock()
gatekeeper_queue_lock = threading.Lock()
insight_cache = {}

CLOSE_FDS = platform.system() != "Windows"

if not os.path.exists(HISTORY_DIR):
    os.makedirs(HISTORY_DIR)

# ==========================================
# --- INITIALIZE HUGGING FACE OPINION ENGINE ---
# ==========================================
print("Initializing Hugging Face Opinion Engine (Local)...")
if SAMSUNG_PIPELINE_ENABLED:
    # Samsung Chat supplies both the structured summary and strategic
    # implication. Avoid holding FLAN-T5 in memory when that production path is
    # active.
    tokenizer = None
    model = None
    print("Opinion Engine delegated to Samsung Chat; local FLAN-T5 not loaded.")
else:
    try:
        # This is a local folder, not a Hugging Face Hub model name. The folder
        # contains config/tokenizer files plus model.safetensors.
        opinion_model_path = str(resolve_model_path("flan-t5-local", "flan-t5-local"))

        # local_files_only=True tells Transformers to load only files already
        # present on disk. It should not try to download missing artifacts.
        tokenizer = AutoTokenizer.from_pretrained(opinion_model_path, local_files_only=True)
        model = AutoModelForSeq2SeqLM.from_pretrained(opinion_model_path, local_files_only=True)
        print("Opinion Engine Ready.")
    except Exception as e:
        print(f"Opinion Engine Failed to Load: {e}")
        tokenizer = None
        model = None

opinion_model_load_attempted = model is not None


def ensure_local_opinion_model():
    """Lazily load FLAN-T5 when Samsung Chat is unavailable at runtime."""

    global tokenizer, model, opinion_model_load_attempted
    if model is not None and tokenizer is not None:
        return True
    with opinion_lock:
        if model is not None and tokenizer is not None:
            return True
        if opinion_model_load_attempted:
            return False
        opinion_model_load_attempted = True
        try:
            opinion_model_path = str(
                resolve_model_path("flan-t5-local", "flan-t5-local")
            )
            print(
                "[PIPELINE:FALLBACK] Loading local FLAN-T5 for Why This "
                "Matters because Samsung Chat is unavailable.",
                flush=True,
            )
            tokenizer = AutoTokenizer.from_pretrained(
                opinion_model_path, local_files_only=True
            )
            model = AutoModelForSeq2SeqLM.from_pretrained(
                opinion_model_path, local_files_only=True
            )
            print("[PIPELINE:FALLBACK] Local FLAN-T5 ready.", flush=True)
            return True
        except Exception as error:
            print(
                "[PIPELINE:FALLBACK] Local FLAN-T5 failed to load: "
                f"{type(error).__name__}: {error}. Deterministic insight "
                "fallback remains active.",
                flush=True,
            )
            tokenizer = None
            model = None
            return False


def generate_opinion(text):
    """Generate a short executive opinion using local FLAN-T5 when available."""

    # If the local model failed to load, export generation degrades gracefully
    # instead of calling an external API.
    if not text or not ensure_local_opinion_model():
        return "Insight generation unavailable."

    try:
        # Only the first 500 chars are used to keep generation fast and bounded.
        prompt = f"Briefly analyze this news and give a one-sentence professional opinion: {text[:500]}"

        # The model object is shared across requests. The lock prevents two
        # threads from calling generate() on the same local model at once.
        with opinion_lock:
            inputs = tokenizer(
                prompt, return_tensors="pt", max_length=512, truncation=True
            )
            outputs = model.generate(**inputs, max_new_tokens=40, do_sample=False)

        return tokenizer.decode(outputs[0], skip_special_tokens=True)

    except Exception as e:
        print(f"Insight Gen Error: {e}")
        return "Could not compute insight."


def fallback_why_it_matters(item):
    source_count = int(item.get("source_count", 1) or 1)
    category = str(item.get("category", "technology intelligence") or "technology intelligence")
    category_lower = category.lower()
    if "broadcast" in category_lower:
        consequence = "It could reshape distribution reach, rights economics, and regulatory priorities for broadcast operators."
    elif "ai" in category_lower:
        consequence = "It may shift product capability, compute demand, and competitive positioning for AI teams."
    elif any(term in category_lower for term in ["display", "television", "device"]):
        consequence = "It may influence product roadmaps, supplier choices, and near-term competitive differentiation."
    else:
        consequence = "It may change competitive priorities, investment choices, or execution risk for decision-makers."
    return (
        f"{consequence} The signal is supported by {source_count} "
        f"source{'s' if source_count != 1 else ''}."
    )


def is_weak_generated_insight(insight, title, summary=""):
    normalized_insight = " ".join(
        re.findall(r"[a-z0-9]+", str(insight or "").lower())
    )
    normalized_summary = " ".join(
        re.findall(r"[a-z0-9]+", str(summary or "").lower())
    )
    generated_words = normalized_insight.split()
    title_words = set(re.findall(r"[a-z0-9]+", str(title or "").lower()))
    if len(generated_words) < 9:
        return True
    # Local generation can occasionally echo the first summary sentence even
    # when that sentence differs from the (sometimes generic) stored title.
    if (
        len(normalized_insight) >= 24
        and normalized_summary
        and normalized_insight in normalized_summary
    ):
        return True
    overlap = sum(1 for word in generated_words if word in title_words)
    return overlap / max(len(generated_words), 1) > 0.7


def generate_why_it_matters(item, profile=DEFAULT_PROFILE):
    title = str(item.get("title", "") or "").strip()
    summary = str(
        item.get("master_summary")
        or item.get("summary")
        or item.get("snippet")
        or ""
    ).strip()
    cache_key = hashlib.sha256(
        f"{profile}|{title}|{summary[:1000]}".encode("utf-8")
    ).hexdigest()

    with insight_cache_lock:
        cached = insight_cache.get(cache_key)
    if cached:
        cached_insight, cached_source = cached
        if not is_weak_generated_insight(cached_insight, title, summary):
            return cached_insight, f"{cached_source}-cache"
        # Older application versions cached occasional headline echoes. Do not
        # let those bypass the current quality gate for the lifetime of the
        # server process.
        with insight_cache_lock:
            insight_cache.pop(cache_key, None)

    if not summary or not ensure_local_opinion_model():
        insight = fallback_why_it_matters(item)
        source = "fallback"
    else:
        prompt = (
            "You are preparing an executive intelligence briefing. "
            "Complete the sentence 'This matters because...' in one concise sentence, "
            "describing strategic impact, risk, opportunity, or market consequence, "
            "not repeating the headline. "
            f"Title: {title}. Summary: {summary[:800]}"
        )
        try:
            with opinion_lock:
                inputs = tokenizer(
                    prompt,
                    return_tensors="pt",
                    max_length=512,
                    truncation=True,
                )
                outputs = model.generate(
                    **inputs,
                    max_new_tokens=64,
                    do_sample=False,
                )
            insight = tokenizer.decode(outputs[0], skip_special_tokens=True).strip()
            source = "flan-t5-local"
            if is_weak_generated_insight(insight, title, summary):
                insight = fallback_why_it_matters(item)
                source = "fallback-after-flan"
        except Exception as e:
            print(f"Why It Matters Gen Error: {e}")
            insight = fallback_why_it_matters(item)
            source = "fallback"

    with insight_cache_lock:
        if len(insight_cache) >= 1000:
            insight_cache.pop(next(iter(insight_cache)))
        insight_cache[cache_key] = (insight, source)
    return insight, source


# ==========================================
# --- TEAM ROUTING LOGIC ---
# ==========================================
def determine_target_team(title, summary):
    text = (title + " " + summary).lower()
    if any(k in text for k in ["cloud", "server", "aws", "azure"]):
        return "Cloud Team"
    elif any(k in text for k in ["hardware", "chip", "semiconductor", "nvidia"]):
        return "Hardware Team"
    elif any(k in text for k in ["robot", "robotics", "automation"]):
        return "Robotics Team"
    elif any(k in text for k in ["tv", "oled", "display", "tizen", "streaming"]):
        return "TV & Display Team"
    elif any(
        k in text
        for k in ["broadcast", "dth", "iptv", "ott", "dvb", "set top box", "trai", "mib"]
    ):
        return "Broadcast Team"
    else:
        return "ALL"


# ==========================================
# --- NEWS CATEGORIZATION ENGINE ---
# ==========================================
CATEGORY_MATRIX = {
    "AI Models": [
        "llm", "gpt", "gemini", "claude", "llama", "foundation model",
        "parameters", "openai", "anthropic", "neural network",
    ],
    "AI Agents": [
        "agent", "autonomous agent", "copilot", "ai assistant",
        "digital assistant", "virtual assistant",
    ],
    "Smart Features": [
        "smart feature", "intelligent feature", "auto-",
        "smart tracking", "adaptive", "predictive",
    ],
    "Form Factor": [
        "form factor", "foldable", "rollable", "wearable", "design",
        "chassis", "slimmer", "hinge",
    ],
    "New Product": [
        "launch", "unveil", "release", "debut", "announced",
        "new lineup", "introducing",
    ],
    "Robotics": [
        "robot", "humanoid", "boston dynamics", "automation",
        "bipedal", "robotic", "drone",
    ],
    "Services": [
        "subscription", "service", "cloud service",
        "platform as a service", "saas", "streaming",
    ],
    "Security": [
        "security", "privacy", "cybersecurity", "hack", "breach",
        "encryption", "knox", "malware",
    ],
    "Smart Home": [
        "smart home", "iot", "thermostat", "fridge", "appliance",
        "smartthings", "matter",
    ],
    "Display Tech": [
        "oled", "microled", "display", "screen", "monitor", "tv",
        "resolution", "nits", "panel",
    ],
    "Partnership": [
        "partnership", "collaboration", "team up", "teaming up",
        "joint venture", "partnered",
    ],
    "Research": [
        "research", "study", "paper", "breakthrough", "scientists",
        "developed a new", "laboratory",
    ],
    "Patent": [
        "patent", "trademark", "intellectual property", "uspto",
        "filed a patent",
    ],
    "Broadcasting": [
        "broadcast", "dth", "cable tv", "iptv", "dvb", "ott", "fast",
        "connected tv", "set top box", "tuner", "linear ad insertion",
        "trai", "mib", "broadcast regulation", "hbbtv", "5g broadcast", "d2m",
        "digital terrestrial transmission", "conditional access",
        "digital rights management",
    ],
    "AI Features": [
        "ai-powered", "generative ai", "genai", "ai capability",
        "ai tool", "ai update",
    ],
}


def assign_category(title, summary):
    text = (str(title) + " " + str(summary)).lower()
    scores = {cat: 0 for cat in CATEGORY_MATRIX}
    for category, keywords in CATEGORY_MATRIX.items():
        for kw in keywords:
            if kw in text:
                scores[category] += text.count(kw)
    best_category = max(scores, key=scores.get)
    if scores[best_category] == 0:
        return "Tech News"
    return best_category


# ==========================================
# --- USER-TAUGHT REGION CLASSIFICATION ---
# ==========================================
def normalize_region_label(value):
    candidate = str(value or "").strip().lower()
    if candidate == "local":
        return "Local"
    if candidate == "global":
        return "Global"
    return None


def normalize_region_keywords(value):
    if isinstance(value, list):
        raw_values = value
    else:
        raw_values = re.split(r"[,;\n]", str(value or ""))
    keywords = []
    for raw in raw_values:
        keyword = re.sub(r"\s+", " ", str(raw).strip().lower())
        if len(keyword) >= 2 and keyword not in keywords:
            keywords.append(keyword)
    return keywords[:20]


def get_region_learning_file_for_profile(profile: str):
    return REGION_LEARNING_FILES.get(profile, REGION_LEARNING_FILES[DEFAULT_PROFILE])


def load_region_learning(profile=DEFAULT_PROFILE):
    blank = {"Local": [], "Global": [], "corrections": []}
    learning_file = get_region_learning_file_for_profile(profile)
    if not os.path.exists(learning_file):
        return blank
    try:
        with open(learning_file, "r", encoding="utf-8") as f:
            stored = json.load(f)
        for region in ("Local", "Global"):
            blank[region] = normalize_region_keywords(stored.get(region, []))
        blank["corrections"] = stored.get("corrections", [])[-500:]
    except (json.JSONDecodeError, OSError, TypeError):
        pass
    return blank


def save_region_learning(data, profile=DEFAULT_PROFILE):
    destination = get_region_learning_file_for_profile(profile)
    temp_file = f"{destination}.{secrets.token_hex(6)}.tmp"
    with open(temp_file, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)
        f.flush()
        os.fsync(f.fileno())
    os.replace(temp_file, destination)


def region_text_for_article(item):
    fields = [
        item.get("title", ""),
        item.get("master_summary", ""),
        item.get("summary", ""),
        item.get("snippet", ""),
        item.get("full_contents", ""),
        item.get("source", ""),
        " ".join(item.get("keywords_found", []) or []),
    ]
    return " ".join(str(field) for field in fields if field).lower()


def apply_learned_region(item, profile=DEFAULT_PROFILE):
    if not isinstance(item, dict):
        return item
    learned = load_region_learning(profile)
    title = str(item.get("title", "")).strip().lower()
    for correction in reversed(learned["corrections"]):
        if title and title == str(correction.get("title", "")).strip().lower():
            next_item = dict(item)
            next_item["region"] = correction["region"]
            next_item["region_basis"] = "User corrected"
            return next_item

    text = region_text_for_article(item)
    for region in ("Local", "Global"):
        matches = [keyword for keyword in learned[region] if keyword in text]
        if matches:
            next_item = dict(item)
            next_item["region"] = region
            next_item["region_basis"] = f"Learned keyword: {matches[0]}"
            return next_item
    return item


def apply_learned_regions(items, profile=DEFAULT_PROFILE):
    return [apply_learned_region(item, profile) for item in (items or [])]


# ==========================================
# --- HELPER: Robust Image Downloader ---
# ==========================================
def download_image_for_export(url, add_border=False):
    if not url:
        return None
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(
            url,
            headers=headers,
            timeout=(5, 10),
            verify=tls_verify("ARTICLE_IMAGE_METADATA"),
        )
        if response.status_code == 200:
            img = Image.open(io.BytesIO(response.content))
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            if add_border:
                img = ImageOps.expand(img, border=15, fill="black")
            out_stream = io.BytesIO()
            img.save(out_stream, format="PNG")
            out_stream.seek(0)
            return out_stream
    except Exception as e:
        print(f"   [IMG] Exception: {e}")
        return None
    return None


def sanitize_filename(filename: str) -> str:
    return Path(filename).name


def cleanup_job_files(*files):
    for f in files:
        try:
            if os.path.exists(f):
                os.remove(f)
        except Exception as e:
            print(f"Cleanup error for {f}: {e}")


# ==========================================
# BOUNCER HELPERS
# ==========================================
def normalize_bouncer_keywords(keywords_found):
    if isinstance(keywords_found, list):
        return ", ".join(str(keyword).strip() for keyword in keywords_found if str(keyword).strip())
    return str(keywords_found or "").strip()


def build_bouncer_text(title, summary, keywords_found=None):
    return (
        f"Title: {str(title or '').strip()}\n"
        f"Keywords: {normalize_bouncer_keywords(keywords_found)}\n"
        f"Summary: {str(summary or '').strip()}"
    )


def get_bouncer_not_interested_score(
    title,
    summary,
    keywords_found=None,
    profile=DEFAULT_PROFILE,
):
    model_for_profile = get_bouncer_model_for_profile(profile)
    if bouncer_embedder is None or model_for_profile is None:
        return None
    try:
        check_text = build_bouncer_text(title, summary, keywords_found)
        vector = bouncer_embedder.encode([check_text])
        if hasattr(model_for_profile, "predict_proba"):
            probabilities = model_for_profile.predict_proba(vector)[0]
            classes = list(getattr(model_for_profile, "classes_", []))
            for candidate in [0, "0", "not_interested", "not_intrested", "irrelevant", "drop", "dislike"]:
                if candidate in classes:
                    return float(probabilities[classes.index(candidate)])
            print(f"[BOUNCER:{profile}] Could not identify not_interested class: {classes}")
            return None
        prediction = model_for_profile.predict(vector)[0]
        return 1.0 if prediction in [0, "0", "not_interested", "not_intrested", "irrelevant", "drop", "dislike"] else 0.0
    except Exception as e:
        print(f"[BOUNCER:{profile}] Bouncer error: {e}")
        return None


def bouncer_decision(title, summary, keywords_found=None, profile=DEFAULT_PROFILE):
    score = get_bouncer_not_interested_score(title, summary, keywords_found, profile)
    if score is None:
        return {"keep": True, "decision": "keep", "score": None, "reason": f"bouncer_unavailable_{profile}"}
    if score >= BOUNCER_HARD_DROP_THRESHOLD:
        return {"keep": False, "decision": "drop", "score": round(score, 4), "reason": f"high_confidence_not_interested_{profile}"}
    if score >= BOUNCER_LOW_PRIORITY_THRESHOLD:
        return {"keep": True, "decision": "low_priority", "score": round(score, 4), "reason": f"medium_confidence_not_interested_{profile}"}
    return {"keep": True, "decision": "keep", "score": round(score, 4), "reason": f"likely_interesting_{profile}"}


def bouncer_check(title, summary, keywords_found=None, profile=DEFAULT_PROFILE):
    return bouncer_decision(title, summary, keywords_found, profile)["keep"]


def attach_bouncer_metadata(item, decision):
    item["bouncer_decision"] = decision.get("decision", "keep")
    item["bouncer_score"] = decision.get("score")
    item["bouncer_reason"] = decision.get("reason", "")
    return item


# ==========================================
# --- DATA MODELS ---
# ==========================================
class Source(BaseModel):
    name: str


class NewsItem(BaseModel):
    title: str
    master_summary: str
    summary_lead: Optional[str] = ""
    summary_points: List[str] = []
    ppt_summary: Optional[str] = ""
    snippet: Optional[str] = ""
    date: str
    link: str
    top_image: Optional[str] = None
    sources: List[Source] = []
    importance_score: int
    keywords_found: List[str] = []
    region: Optional[str] = "Global"
    full_contents: Optional[str] = ""
    selected_by: Optional[str] = None
    category: Optional[str] = "Tech News"
    why_it_matters: Optional[str] = ""
    article_intent: Optional[str] = ""
    summarized_by: Optional[str] = ""


class ExportRequest(BaseModel):
    items: List[NewsItem]
    filename: str = "SENSE_Brief.pptx"


class VotePayload(BaseModel):
    keywords: list
    summary: str
    vote: str
    title: Optional[str] = ""


# ==========================================
# --- NOT INTERESTED STORE ---
# ==========================================
def load_not_interested_store(request: Request = None):
    profile = get_active_profile_name(request)
    not_interested_file = get_not_interested_file_for_profile(profile)
    if not os.path.exists(not_interested_file):
        return []
    try:
        with open(not_interested_file, "r", encoding="utf-8") as f:
            store = json.load(f)
    except (json.JSONDecodeError, Exception):
        return []

    now = datetime.datetime.now()
    active = []
    expired_count = 0
    for item in store:
        try:
            rejected_at = datetime.datetime.strptime(
                item.get("rejected_at", ""), "%Y-%m-%d %H:%M:%S"
            )
            age_hours = (now - rejected_at).total_seconds() / 3600
            if age_hours <= NOT_INTERESTED_EXPIRY_HOURS:
                active.append(item)
            else:
                expired_count += 1
        except (ValueError, TypeError):
            active.append(item)

    if expired_count > 0:
        try:
            temp_file = f"{not_interested_file}.{secrets.token_hex(6)}.tmp"
            with open(temp_file, "w", encoding="utf-8") as f:
                json.dump(active, f, indent=4, ensure_ascii=False)
                f.flush()
                os.fsync(f.fileno())
            os.replace(temp_file, not_interested_file)
            print(f"[NOT-INTERESTED:{profile}] Expired {expired_count} entries (>{NOT_INTERESTED_EXPIRY_HOURS}h old)")
        except Exception:
            pass

    return active


def save_not_interested_store(store, request: Request = None):
    profile = get_active_profile_name(request)
    destination = get_not_interested_file_for_profile(profile)
    temp_file = f"{destination}.{secrets.token_hex(6)}.tmp"
    with open(temp_file, "w", encoding="utf-8") as f:
        json.dump(store, f, indent=4, ensure_ascii=False)
        f.flush()
        os.fsync(f.fileno())
    os.replace(temp_file, destination)


def is_already_rejected(title, store):
    normalized_title = title.strip().lower()
    for item in store:
        if item.get("title", "").strip().lower() == normalized_title:
            return True
    return False


# ==========================================
# --- WORKFLOW HELPERS ---
# ==========================================
def load_workflow_store(request: Request = None):
    workflow_file = get_workflow_file_for_request(request)
    if not os.path.exists(workflow_file):
        return {"selected": [], "approved": []}
    try:
        with open(workflow_file, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, dict):
            return {"selected": [], "approved": []}
        data.setdefault("selected", [])
        data.setdefault("approved", [])
        return data
    except Exception:
        return {"selected": [], "approved": []}


def save_workflow_store(data, request: Request = None):
    destination = get_workflow_file_for_request(request)
    temp_file = f"{destination}.{secrets.token_hex(6)}.tmp"
    with open(temp_file, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)
        f.flush()
        os.fsync(f.fileno())
    os.replace(temp_file, destination)


def resolve_workflow_identities(store):
    """Render editable viewer names from stable protected viewer identifiers."""
    profiles = load_viewer_profiles()
    resolved = {"selected": [], "approved": []}
    for list_name in ("selected", "approved"):
        for raw_item in store.get(list_name, []):
            item = dict(raw_item)
            selected_key = str(item.get("selected_by_id", "")).strip()
            approved_key = str(item.get("approved_by_id", "")).strip()
            if selected_key and selected_key in profiles:
                item["selected_by"] = profiles[selected_key].get("display_name") or item.get("selected_by")
            if approved_key and approved_key in profiles:
                item["approved_by"] = profiles[approved_key].get("display_name") or item.get("approved_by")
            resolved[list_name].append(item)
    return resolved


def refresh_workflow_identity(viewer_key, previous_name, display_name):
    """Keep legacy copied names compatible while stable IDs become canonical."""
    with workflow_lock:
        for workflow_file in set(WORKFLOW_FILES.values()):
            if not os.path.exists(workflow_file):
                continue
            try:
                with open(workflow_file, "r", encoding="utf-8") as handle:
                    store = json.load(handle)
                changed = False
                for list_name in ("selected", "approved"):
                    for item in store.get(list_name, []):
                        selected_matches = (
                            item.get("selected_by_id") == viewer_key
                            or (
                                previous_name
                                and str(item.get("selected_by", "")).casefold()
                                == str(previous_name).casefold()
                            )
                        )
                        approved_matches = (
                            item.get("approved_by_id") == viewer_key
                            or (
                                previous_name
                                and str(item.get("approved_by", "")).casefold()
                                == str(previous_name).casefold()
                            )
                        )
                        if selected_matches:
                            item["selected_by_id"] = viewer_key
                            item["selected_by"] = display_name
                            changed = True
                        if approved_matches:
                            item["approved_by_id"] = viewer_key
                            item["approved_by"] = display_name
                            changed = True
                if changed:
                    temp_file = f"{workflow_file}.{secrets.token_hex(6)}.tmp"
                    with open(temp_file, "w", encoding="utf-8") as handle:
                        json.dump(store, handle, indent=4, ensure_ascii=False)
                        handle.flush()
                        os.fsync(handle.fileno())
                    os.replace(temp_file, workflow_file)
            except Exception as error:
                print(f"[IDENTITY] Could not refresh {workflow_file}: {error}", flush=True)


# ==========================================
# --- DROPPED ARTICLES DATA ---
# ==========================================
def log_dropped_article(title, summary, keywords_found, bouncer_info=None, profile=DEFAULT_PROFILE, article=None):
    dropped_file = os.path.join(ROOT_DIR, "dropped_articles.json")
    bouncer_info = bouncer_info or {}
    article = dict(article or {})
    timestamp = datetime.datetime.now().isoformat(timespec="seconds")
    new_entry = {
        "id": secrets.token_hex(12),
        "timestamp": timestamp,
        "updated_at": timestamp,
        "status": "dropped",
        "profile": profile,
        "title": title,
        "keyword": keywords_found if keywords_found else [title[:50]],
        "keywords_found": keywords_found if keywords_found else [],
        "summary": summary,
        "source": article.get("source", "Unknown"),
        "link": article.get("link") or article.get("url") or "",
        "date": article.get("date", ""),
        "top_image": article.get("top_image", ""),
        "full_contents": article.get("full_contents") or article.get("full_content") or summary,
        "sources": article.get("sources", []),
        "label": "not_interested",
        "bouncer_decision": bouncer_info.get("decision", "drop"),
        "bouncer_score": bouncer_info.get("score"),
        "bouncer_reason": bouncer_info.get("reason", ""),
        "bouncer_stage": article.get("bouncer_stage", "unknown"),
        "restore_eligible": bool(article.get("link") or article.get("url")),
        "stage_history": [{
            "timestamp": timestamp,
            "stage": article.get("bouncer_stage", "unknown"),
            "reason": bouncer_info.get("reason", ""),
        }],
    }
    with dropped_lock:
        dropped = []
        if os.path.exists(dropped_file):
            try:
                with open(dropped_file, "r", encoding="utf-8") as f:
                    dropped = json.load(f)
            except Exception:
                dropped = []
        normalized_link = str(new_entry.get("link", "")).strip().lower()
        normalized_title = str(title or "").strip().casefold()
        existing = next(
            (
                entry for entry in dropped
                if entry.get("profile", DEFAULT_PROFILE) == profile
                and (
                    (normalized_link and str(entry.get("link", "")).strip().lower() == normalized_link)
                    or str(entry.get("title", "")).strip().casefold() == normalized_title
                )
            ),
            None,
        )
        if existing:
            history = list(existing.get("stage_history", []))
            history.extend(new_entry["stage_history"])
            preserved_id = existing.get("id") or new_entry["id"]
            existing.update(new_entry)
            existing["id"] = preserved_id
            existing["stage_history"] = history[-25:]
            existing["status"] = "dropped"
        else:
            dropped.append(new_entry)
        dropped = dropped[-500:]
        with open(dropped_file, "w", encoding="utf-8") as f:
            json.dump(dropped, f, indent=4, ensure_ascii=False)


# ==========================================
# TRAINING AND FILTER HELPERS
# ==========================================
def get_bouncer_summary_from_item(item):
    parts = []
    for key in ["master_summary", "summary", "snippet", "full_content", "full_contents"]:
        text = str(item.get(key, "") or "").strip()
        if text and text not in parts:
            parts.append(text)
    return (" ".join(parts).strip() or str(item.get("title", "")).strip())[:2500]


def pipeline_cache_key(item, *, include_sources=False):
    links = []
    if include_sources:
        for source in item.get("sources", []) or []:
            if isinstance(source, dict):
                links.append(str(source.get("link") or source.get("url") or ""))
    identity = "|".join(
        [
            str(item.get("canonical_link") or item.get("link") or "").strip().lower(),
            str(item.get("title") or "").strip().casefold(),
            str(item.get("date") or "").strip(),
            "|".join(sorted(link.strip().lower() for link in links if link.strip())),
            hashlib.sha256(
                str(
                    item.get("full_contents")
                    or item.get("summary_input")
                    or ""
                )[:20000].encode("utf-8")
            ).hexdigest()
            if include_sources
            else "",
        ]
    )
    return hashlib.sha256(identity.encode("utf-8")).hexdigest()


def cache_success(store, key, item):
    def update(values):
        values = values if isinstance(values, dict) else {}
        values[key] = dict(item)
        while len(values) > SAMSUNG_CACHE_MAX_ITEMS:
            values.pop(next(iter(values)))
        return values

    store.update(update)


CHAT_SUMMARY_CACHE_FIELDS = {
    "title",
    "summary",
    "summary_lead",
    "summary_points",
    "key_points",
    "master_summary",
    "ppt_summary",
    "why_it_matters",
    "why_matters",
    "attention_hook",
    "what_changed",
    "why_now",
    "watch_next",
    "hook_type",
    "hook_source",
    "hook_grounded",
    "article_intent",
    "category",
    "region",
    "importance_score",
    "chat_summary_status",
    "summarized_by",
    "summary_format",
    "chat_model_id",
}


def chat_summary_cache_payload(item):
    """Cache generated fields only, never a private or source article body."""

    return {
        key: value
        for key, value in dict(item or {}).items()
        if key in CHAT_SUMMARY_CACHE_FIELDS
    }


def keyword_terms(keywords):
    if isinstance(keywords, str):
        values = keywords.split(",")
    else:
        values = keywords or []
    return [
        str(value).strip()
        for value in values
        if str(value).strip()
    ]


def matched_pipeline_keywords(item, keywords):
    text = " ".join(
        str(item.get(key) or "")
        for key in (
            "title",
            "full_contents",
            "web_search_content",
            "summary_input",
            "summary",
            "snippet",
        )
    )
    matches = []
    for keyword in keyword_terms(keywords):
        if re.search(
            r"(?<!\w)" + re.escape(keyword) + r"(?!\w)",
            text,
            flags=re.IGNORECASE,
        ):
            matches.append(keyword)
    return matches


class WebSearchRuntimeFailure(RuntimeError):
    """Signal that discovery-only candidates require a full Scrapy retry."""


def run_bouncer_filter_on_items(items, profile=DEFAULT_PROFILE, stage="raw"):
    filtered_items = []
    dropped_count = 0
    low_priority_count = 0
    for item in items if isinstance(items, list) else []:
        decision = bouncer_decision(
            item.get("title", ""),
            get_bouncer_summary_from_item(item),
            item.get("keywords_found", []),
            profile,
        )
        if decision["keep"]:
            item["profile"] = profile
            item["bouncer_profile"] = profile
            item["bouncer_stage"] = stage
            item["bouncer_decision"] = decision["decision"]
            item["bouncer_score"] = decision["score"]
            item["bouncer_reason"] = decision["reason"]
            filtered_items.append(item)
            if decision["decision"] == "low_priority":
                low_priority_count += 1
        else:
            dropped_count += 1
            item["bouncer_profile"] = profile
            item["bouncer_stage"] = stage
            log_dropped_article(
                item.get("title", ""),
                get_bouncer_summary_from_item(item),
                item.get("keywords_found", []),
                decision,
                profile,
                item,
            )
    return filtered_items, dropped_count, low_priority_count


def enrich_raw_articles(
    items,
    keywords,
    profile=DEFAULT_PROFILE,
    use_web_search=None,
    raise_on_service_failure=False,
):
    """Use Samsung Web Search as the article extraction stage when configured."""

    source_items = list(items or [])
    enabled = (
        WEB_SEARCH_ENRICHMENT_ENABLED
        if use_web_search is None
        else bool(use_web_search)
    )
    if not enabled:
        return source_items
    if enrich_article_with_web_search is None:
        message = "Samsung Web Search adapter is unavailable"
        if WEB_SEARCH_REQUIRE_SUCCESS:
            raise RuntimeError(message)
        print(f"[PIPELINE:{profile}] {message}; keeping crawler content.", flush=True)
        return source_items

    limit = WEB_SEARCH_MAX_ENRICH_PER_RUN or len(source_items)
    output = []
    successful = 0
    rejected_for_keywords = 0
    service_failures = 0
    strict = WEB_SEARCH_REQUIRE_SUCCESS or SAMSUNG_PIPELINE_ENABLED
    print(f"[PIPELINE:{profile}] Web Search enrichment: {min(limit, len(source_items))} article(s).", flush=True)
    for index, item in enumerate(source_items):
        if index >= limit:
            if not strict:
                output.append(item)
            continue
        cache_key = pipeline_cache_key(item)
        cached_values = WEB_SEARCH_CACHE.read()
        cached = (
            cached_values.get(cache_key)
            if isinstance(cached_values, dict)
            else None
        )
        if isinstance(cached, dict):
            enriched = dict(cached)
            enriched["enrichment_cache"] = "hit"
        else:
            try:
                enriched = enrich_article_with_web_search(item, keywords=keywords)
            except Exception as error:
                enriched = {
                    **item,
                    "enrichment_status": "failed",
                    "enrichment_error": f"{type(error).__name__}: {error}"[:500],
                }
            if enriched.get("enrichment_status") == "success":
                cache_success(WEB_SEARCH_CACHE, cache_key, enriched)
                enriched["enrichment_cache"] = "miss"
        if enriched.get("enrichment_status") == "success":
            matches = matched_pipeline_keywords(enriched, keywords)
            enriched["keywords_found"] = matches
            if WEB_SEARCH_REQUIRE_KEYWORD_MATCH and not matches:
                rejected_for_keywords += 1
                enriched["enrichment_status"] = "rejected_keyword_mismatch"
                print(
                    f"[PIPELINE:{profile}] Web Search content did not match "
                    f"profile keywords: {item.get('title', '')[:70]}",
                    flush=True,
                )
                if not strict:
                    output.append(enriched)
            else:
                successful += 1
                output.append(enriched)
        elif strict:
            service_failures += 1
            print(f"[PIPELINE:{profile}] Enrichment rejected: {item.get('title', '')[:70]}", flush=True)
        else:
            output.append(enriched)
        if index + 1 < min(limit, len(source_items)) and WEB_SEARCH_ENRICH_DELAY_SECONDS:
            time.sleep(WEB_SEARCH_ENRICH_DELAY_SECONDS)
    print(
        f"[PIPELINE:{profile}] Web Search enriched {successful}/"
        f"{min(limit, len(source_items))}; keyword mismatches "
        f"{rejected_for_keywords}.",
        flush=True,
    )
    if raise_on_service_failure and service_failures:
        raise WebSearchRuntimeFailure(
            "Samsung Web Search failed for "
            f"{service_failures}/{min(limit, len(source_items))} candidate(s) "
            "after preflight; full Scrapy extraction is required."
        )
    return output


def structure_summary_for_dossier(item, summarized_by="local_bart"):
    """Normalize any local summary into the same lead-plus-bullets UI contract."""

    output = dict(item or {})
    summary = str(
        output.get("master_summary")
        or output.get("summary")
        or output.get("snippet")
        or ""
    ).strip()
    content = str(
        output.get("full_contents")
        or output.get("full_content")
        or output.get("summary_input")
        or summary
    ).strip()
    summary_sentences = [
        part.strip()
        for part in re.split(r"(?<=[.!?])\s+", summary)
        if len(part.strip()) >= 12
    ]
    lead_parts = summary_sentences[:2] or ([summary] if summary else [])
    lead = " ".join(lead_parts).strip()
    candidates = summary_sentences[2:]
    candidates.extend(
        part.strip()
        for part in re.split(r"(?<=[.!?])\s+", content)
        if len(part.strip()) >= 24
    )
    points = []
    seen = {part.casefold() for part in lead_parts}
    for candidate in candidates:
        normalized = re.sub(r"\s+", " ", candidate).strip()
        key = normalized.casefold()
        if key and key not in seen:
            seen.add(key)
            points.append(normalized)
        if len(points) >= 5:
            break
    if not points and lead:
        points = [lead]
    output["summary"] = lead or summary
    output["summary_lead"] = lead or summary
    output["summary_points"] = points
    output["key_points"] = points
    output["master_summary"] = " ".join(
        [lead or summary, *[f"• {point}" for point in points]]
    ).strip()
    output["summary_format"] = "lead_and_bullets"
    output["summarized_by"] = summarized_by
    return output


def apply_local_bart_fallback(items, profile=DEFAULT_PROFILE):
    """Generate local BART summaries after Chat preflight or request failure."""

    source_items = [dict(item) for item in (items or [])]
    if not source_items:
        return []
    print(
        f"[PIPELINE:{profile}] LOCAL FALLBACK: loading BART to summarize "
        f"{len(source_items)} article(s).",
        flush=True,
    )
    try:
        from news_scrapper.semantic_clustering import MinimalSemanticEngine

        engine = MinimalSemanticEngine(load_summarizer=True)
        for index, item in enumerate(source_items, 1):
            content = str(
                item.get("full_contents")
                or item.get("full_content")
                or item.get("summary_input")
                or item.get("master_summary")
                or item.get("snippet")
                or ""
            ).strip()
            if content:
                bart_summary = engine.generate_dynamic_summary([content[:12000]])
                if bart_summary:
                    item["master_summary"] = bart_summary
                    item["summary"] = bart_summary
                    item["ppt_summary"] = engine.generate_ppt_summary([content[:12000]])
            source_items[index - 1] = structure_summary_for_dossier(
                item, "local_bart"
            )
            print(
                f"[PIPELINE:{profile}] BART fallback {index}/"
                f"{len(source_items)}: {item.get('title', '')[:70]}",
                flush=True,
            )
    except Exception as error:
        print(
            f"[PIPELINE:{profile}] BART fallback failed: "
            f"{type(error).__name__}: {error}. Using deterministic structured "
            "extractive summaries so the dossier remains readable.",
            flush=True,
        )
        source_items = [
            structure_summary_for_dossier(item, "local_extractive")
            for item in source_items
        ]
    return source_items


def enrich_final_articles(items, profile=DEFAULT_PROFILE, use_chat=None):
    """Add secure image metadata and optional Samsung Chat final summaries."""

    output = []
    chat_enabled = (
        FINAL_CHAT_SUMMARY_ENABLED if use_chat is None else bool(use_chat)
    )
    chat_limit = FINAL_CHAT_SUMMARY_MAX_ARTICLES or len(items or [])
    chat_failures = []
    for index, original in enumerate(items or []):
        item = dict(original)
        if enrich_article_image_metadata is not None:
            item = enrich_article_image_metadata(item)
        if chat_enabled and summarize_article_with_chat is not None and index < chat_limit:
            cache_key = pipeline_cache_key(item, include_sources=True)
            cached_values = CHAT_SUMMARY_CACHE.read()
            cached = (
                cached_values.get(cache_key)
                if isinstance(cached_values, dict)
                else None
            )
            if isinstance(cached, dict):
                # Preserve this request's identity, source, full text and
                # private-scope metadata. The cache only supplies generated
                # summary fields for matching article content.
                item.update(chat_summary_cache_payload(cached))
                item["chat_summary_cache"] = "hit"
            else:
                summarized = summarize_article_with_chat(item)
                if summarized.get("chat_summary_status") == "success":
                    cache_success(
                        CHAT_SUMMARY_CACHE,
                        cache_key,
                        chat_summary_cache_payload(summarized),
                    )
                item = summarized
                item["chat_summary_cache"] = "miss"
            if item.get("chat_summary_status") != "success":
                print(
                    f"[PIPELINE:{profile}] Samsung Chat failed for "
                    f"'{item.get('title', '')[:70]}': "
                    f"{item.get('chat_summary_error', 'unknown failure')}. "
                    "Queued for local BART fallback.",
                    flush=True,
                )
                chat_failures.append(len(output))
            if index + 1 < min(chat_limit, len(items)) and FINAL_CHAT_SUMMARY_DELAY_SECONDS:
                time.sleep(FINAL_CHAT_SUMMARY_DELAY_SECONDS)
        output.append(item)
    if not chat_enabled:
        return [
            structure_summary_for_dossier(item, "local_bart")
            for item in output
        ]
    if chat_failures:
        fallback_items = apply_local_bart_fallback(
            [output[index] for index in chat_failures],
            profile,
        )
        for index, fallback in zip(chat_failures, fallback_items):
            fallback["chat_summary_status"] = "fallback"
            output[index] = fallback
    return output


PERSONAL_BRIEFING_STORE = JsonStore(Path(VIEWER_BRIEFING_FILE), dict)
PRIVATE_VIEWER_CLAIMS = JsonStore(Path(VIEWER_IDENTITY_CLAIMS_FILE), dict)
PERSONALIZATION_SERVICE = PersonalizationService(
    Path(VIEWER_PERSONALIZATION_FILE),
    window_days=30,
)
PERSONAL_BRIEFING_MAX_URLS = max(
    1, min(50, int(os.environ.get("PERSONAL_BRIEFING_MAX_URLS", "20")))
)
PERSONAL_BRIEFING_FETCH_TIMEOUT = max(
    5, int(os.environ.get("PERSONAL_BRIEFING_FETCH_TIMEOUT", "30"))
)
PERSONAL_BRIEFING_MAX_BYTES = max(
    250_000, int(os.environ.get("PERSONAL_BRIEFING_MAX_BYTES", "5000000"))
)


def canonical_personal_url(value):
    try:
        parsed = urlsplit(str(value or "").strip())
    except ValueError:
        return ""
    if (
        parsed.scheme.lower() not in {"http", "https"}
        or not parsed.hostname
        or parsed.username
        or parsed.password
    ):
        return ""
    port = parsed.port
    if port and port not in {80, 443}:
        return ""
    netloc = parsed.hostname.lower().removeprefix("www.")
    if port:
        netloc = f"{netloc}:{port}"
    return urlunsplit(
        (
            parsed.scheme.lower(),
            netloc,
            parsed.path or "/",
            parsed.query,
            "",
        )
    )


def assert_public_article_url(value):
    canonical = canonical_personal_url(value)
    if not canonical:
        raise ValueError("Use a valid public HTTP or HTTPS article URL.")
    hostname = urlsplit(canonical).hostname
    if hostname in {"localhost"} or hostname.endswith((".local", ".internal")):
        raise ValueError("Private or local network URLs are not allowed.")
    try:
        addresses = {
            result[4][0]
            for result in socket.getaddrinfo(
                hostname,
                urlsplit(canonical).port or 443,
                type=socket.SOCK_STREAM,
            )
        }
    except socket.gaierror as error:
        raise ValueError("The article hostname could not be resolved.") from error
    for address in addresses:
        parsed_ip = ipaddress.ip_address(address)
        if not parsed_ip.is_global:
            raise ValueError("Private or local network URLs are not allowed.")
    return canonical


def fetch_personal_article(url):
    current_url = assert_public_article_url(url)
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 Chrome/124 Safari/537.36"
        )
    }
    for _ in range(5):
        response = requests.get(
            current_url,
            headers=headers,
            timeout=PERSONAL_BRIEFING_FETCH_TIMEOUT,
            verify=tls_verify("PERSONAL_BRIEFING"),
            allow_redirects=False,
            stream=True,
        )
        if response.status_code in {301, 302, 303, 307, 308}:
            location = response.headers.get("location", "")
            if not location:
                raise RuntimeError("Article redirect did not provide a destination.")
            current_url = assert_public_article_url(urljoin(current_url, location))
            continue
        if response.status_code >= 400:
            raise RuntimeError(f"Article server returned HTTP {response.status_code}.")
        content_type = str(response.headers.get("content-type", "")).lower()
        if content_type and "html" not in content_type:
            raise RuntimeError("The submitted URL is not an HTML article.")
        body = bytearray()
        for chunk in response.iter_content(chunk_size=64 * 1024):
            body.extend(chunk)
            if len(body) > PERSONAL_BRIEFING_MAX_BYTES:
                raise RuntimeError("Article response exceeded the safe download limit.")
        html = bytes(body).decode(response.encoding or "utf-8", errors="replace")
        article = Article(url=current_url)
        article.set_html(html)
        article.parse()
        text = re.sub(r"\s+", " ", str(article.text or "")).strip()
        if len(text) < 200:
            raise RuntimeError("Could not extract enough readable article text.")
        published = article.publish_date
        return {
            "title": str(article.title or urlsplit(current_url).hostname).strip(),
            "link": current_url,
            "canonical_link": current_url,
            "source": urlsplit(current_url).hostname.removeprefix("www."),
            "date": (
                published.isoformat()
                if hasattr(published, "isoformat")
                else datetime.date.today().isoformat()
            ),
            "top_image": str(article.top_image or "").strip(),
            "full_contents": text,
            "summary_input": text,
            "snippet": text[:1000],
            "master_summary": text[:1500],
            "sources": [
                {
                    "name": urlsplit(current_url).hostname.removeprefix("www."),
                    "link": current_url,
                }
            ],
            "source_count": 1,
            "origin": "personal_briefing",
        }
    raise RuntimeError("Article redirected too many times.")


def update_personal_briefing_job(viewer_key, profile, job_id, **changes):
    def updater(store):
        viewer = store.setdefault(viewer_key, {})
        jobs = viewer.setdefault(profile, [])
        for job in jobs:
            if job.get("id") == job_id:
                job.update(changes)
                job["updated_at"] = datetime.datetime.now().isoformat(
                    timespec="seconds"
                )
                break
        return store

    with viewer_briefing_lock:
        return PERSONAL_BRIEFING_STORE.update(updater)


def cluster_personal_briefing_article(viewer_key, profile, item):
    """Attach private, viewer-scoped semantic relationship metadata.

    Personal submissions are never merged into the shared briefing. Related
    submissions are only grouped inside the current viewer/profile namespace.
    If MiniLM is unavailable, every article remains a safe singleton.
    """

    store = PERSONAL_BRIEFING_STORE.read()
    jobs = store.get(viewer_key, {}).get(profile, [])
    existing = [
        dict(job["article"])
        for job in jobs
        if job.get("status") == "complete" and isinstance(job.get("article"), dict)
    ]
    articles = [*existing, dict(item)]
    if len(articles) < 2:
        item["personal_cluster_id"] = f"personal-cluster-{item['id']}"
        item["related_private_count"] = 0
        return item

    try:
        from news_scrapper.semantic_clustering import MinimalSemanticEngine

        clusters = MinimalSemanticEngine().semantic_cluster(articles)
    except Exception as error:
        print(
            f"[PERSONAL:{profile}] Semantic grouping unavailable: "
            f"{type(error).__name__}: {error}",
            flush=True,
        )
        clusters = [[article] for article in articles]

    cluster_updates = {}
    for cluster in clusters:
        member_ids = sorted(str(article.get("id") or "") for article in cluster)
        cluster_id = "personal-cluster-" + hashlib.sha256(
            "|".join(member_ids).encode("utf-8")
        ).hexdigest()[:16]
        for article in cluster:
            article_id = str(article.get("id") or "")
            cluster_updates[article_id] = {
                "personal_cluster_id": cluster_id,
                "related_private_count": max(0, len(cluster) - 1),
            }

    current_id = str(item.get("id") or "")
    item.update(cluster_updates.get(current_id, {}))

    def updater(current_store):
        viewer_jobs = current_store.setdefault(viewer_key, {}).setdefault(profile, [])
        for job in viewer_jobs:
            article = job.get("article")
            if isinstance(article, dict):
                article.update(cluster_updates.get(str(article.get("id") or ""), {}))
        return current_store

    with viewer_briefing_lock:
        PERSONAL_BRIEFING_STORE.update(updater)
    return item


def process_personal_briefing_job(viewer_key, profile, job_id, url):
    try:
        update_personal_briefing_job(
            viewer_key,
            profile,
            job_id,
            status="processing",
            stage="extracting",
            progress=18,
            message="Opening the article and extracting its story.",
        )
        seed = {
            "title": url,
            "link": url,
            "canonical_link": url,
            "source": urlsplit(url).hostname.removeprefix("www."),
            "origin": "personal_briefing",
        }
        capabilities = resolve_pipeline_capabilities()
        item = None
        extraction_engine = "targeted_scrapy_fallback"
        if capabilities.get("web_search") and enrich_article_with_web_search:
            update_personal_briefing_job(
                viewer_key,
                profile,
                job_id,
                stage="web_search",
                progress=30,
                message="Samsung Web Search is reading the exact article.",
            )
            enriched = enrich_article_with_web_search(seed, keywords=[])
            if enriched.get("enrichment_status") == "success":
                item = enriched
                extraction_engine = "samsung_web_search"
        if item is None:
            update_personal_briefing_job(
                viewer_key,
                profile,
                job_id,
                stage="local_extraction",
                progress=42,
                message="Using secure targeted extraction for this article.",
            )
            item = fetch_personal_article(url)

        item["category"] = assign_category(
            item.get("title", ""),
            item.get("full_contents") or item.get("summary", ""),
        )
        item["region"] = apply_learned_region(item, profile).get(
            "region", "Global"
        )
        item["extracted_by"] = extraction_engine
        item["private_scope"] = "current_viewer_only"
        item["origin"] = "personal_briefing"
        item["submitted_url"] = url
        item["id"] = f"personal-{job_id}"

        update_personal_briefing_job(
            viewer_key,
            profile,
            job_id,
            stage="summarizing",
            progress=68,
            message=(
                "Samsung Chat is shaping the summary and strategic context."
                if capabilities.get("chat")
                else "Local AI is shaping the summary and strategic context."
            ),
        )
        if capabilities.get("chat"):
            item = enrich_final_articles(
                [item],
                profile,
                use_chat=True,
            )[0]
        else:
            if enrich_article_image_metadata is not None:
                item = enrich_article_image_metadata(item)
            item = apply_local_bart_fallback([item], profile)[0]
        item["why_matters"], item["insight_source"] = generate_why_it_matters(
            item,
            profile,
        )
        item["created_at"] = datetime.datetime.now().isoformat(
            timespec="seconds"
        )
        item = cluster_personal_briefing_article(
            viewer_key,
            profile,
            item,
        )

        update_personal_briefing_job(
            viewer_key,
            profile,
            job_id,
            status="complete",
            stage="complete",
            progress=100,
            message="Your private briefing is ready.",
            article=item,
            error=None,
        )
    except Exception as error:
        print(
            f"[PERSONAL:{profile}] {url} failed: "
            f"{type(error).__name__}: {error}",
            flush=True,
        )
        update_personal_briefing_job(
            viewer_key,
            profile,
            job_id,
            status="failed",
            stage="failed",
            progress=100,
            message="This article could not be prepared.",
            error=f"{type(error).__name__}: {error}"[:500],
        )


def resume_personal_briefing_jobs():
    store = PERSONAL_BRIEFING_STORE.read()
    for viewer_key, profiles in store.items():
        if not isinstance(profiles, dict):
            continue
        for profile, jobs in profiles.items():
            for job in jobs if isinstance(jobs, list) else []:
                if job.get("status") in {"queued", "processing"}:
                    update_personal_briefing_job(
                        viewer_key,
                        profile,
                        job.get("id"),
                        status="queued",
                        stage="queued",
                        progress=5,
                        message="Resumed after server restart.",
                    )
                    try:
                        personal_briefing_executor.submit(
                            process_personal_briefing_job,
                            viewer_key,
                            profile,
                            job.get("id"),
                            job.get("url"),
                        )
                    except Exception as error:
                        update_personal_briefing_job(
                            viewer_key,
                            profile,
                            job.get("id"),
                            status="failed",
                            stage="dispatch_failed",
                            progress=100,
                            message=(
                                "The private briefing worker could not restart. "
                                "Use Retry when the service is ready."
                            ),
                            error=f"{type(error).__name__}: {error}"[:500],
                        )
                        print(
                            f"[PERSONAL:{profile}] Could not resume job "
                            f"{job.get('id')}: {type(error).__name__}: {error}",
                            flush=True,
                        )


def save_training_vote(
    keywords,
    summary,
    vote,
    title="",
    profile=DEFAULT_PROFILE,
    *,
    consensus_article_id="",
    consensus_meta=None,
):
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    training_file = get_training_file_for_profile(profile)
    new_row = {
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "title": str(title or "").strip(),
        "keyword": keywords,
        "summary": str(summary or "").strip(),
        "label": vote,
        "profile": profile,
    }
    if consensus_article_id:
        new_row["origin"] = "reaction_consensus"
        new_row["consensus_article_id"] = str(consensus_article_id)[:128]
        new_row["consensus"] = dict(consensus_meta or {})
    with file_lock:
        memory = []
        if os.path.exists(training_file):
            try:
                with open(training_file, "r", encoding="utf-8") as f:
                    memory = json.load(f)
            except json.JSONDecodeError:
                memory = []

        # A mature aggregate is one authoritative row per article. If later
        # votes flip the consensus, replace the previous aggregate instead of
        # teaching the Bouncer contradictory labels for the same story.
        if consensus_article_id:
            memory = [
                row for row in memory
                if str(row.get("consensus_article_id") or "") != str(consensus_article_id)
            ]

        normalized_title = str(title or "").strip().lower()[:150]
        normalized_new = str(summary or "").strip().lower()[:200]
        normalized_vote = str(vote or "").strip().lower()
        is_duplicate = False
        for existing in memory:
            existing_title = existing.get("title", "").strip().lower()[:150]
            existing_summary = existing.get("summary", "").strip().lower()[:200]
            existing_vote = existing.get("label", "").strip().lower()
            if (
                existing_title == normalized_title
                and existing_summary == normalized_new
                and existing_vote == normalized_vote
            ):
                is_duplicate = True
                break

        if not is_duplicate:
            memory.append(new_row)
            temp_file = f"{training_file}.{secrets.token_hex(6)}.tmp"
            with open(temp_file, "w", encoding="utf-8") as f:
                json.dump(memory, f, indent=4, ensure_ascii=False)
                f.flush()
                os.fsync(f.fileno())
            os.replace(temp_file, training_file)
        else:
            print(f"Dedup: Skipped duplicate {vote} vote for: {summary[:50]}...")

    return len(memory)


def reload_bouncer_model_for_profile(profile=DEFAULT_PROFILE):
    global bouncer_model
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    model_file = get_bouncer_model_file_for_profile(profile)
    if not os.path.exists(model_file):
        return False
    try:
        with open(model_file, "rb") as f:
            bouncer_models[profile] = pickle.load(f)
        if profile == DEFAULT_PROFILE:
            bouncer_model = bouncer_models.get(profile)
        return True
    except Exception as e:
        print(f"[BOUNCER:{profile}] Failed to reload model: {e}")
        return False


def retrain_and_reload(profile=DEFAULT_PROFILE):
    """Train one profile while serializing model replacement.

    The request-facing code calls :func:`enqueue_bouncer_retrain`; this
    function remains safe to call directly from a maintenance script as well.
    It deliberately waits for the lock instead of dropping a training event.
    """

    profile = profile if profile in PROFILE_CONFIGS else DEFAULT_PROFILE
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    with train_lock:
        try:
            print(f"\n[BOUNCER:{profile}] Retraining with new data...", flush=True)
            subprocess.run(
                [
                    sys.executable,
                    "-m",
                    "news_scrapper.train_bouncer",
                    "--profile",
                    profile,
                ],
                check=True,
                cwd=str(PROJECT_ROOT),
            )
            reloaded = reload_bouncer_model_for_profile(profile)
            print(
                f"[BOUNCER:{profile}] Brain successfully upgraded and reloaded "
                f"(model_loaded={reloaded}).\n",
                flush=True,
            )
            return True
        except Exception as e:
            print(f"[BOUNCER:{profile}] Failed to retrain Bouncer: {e}", flush=True)
            return False


def _training_worker_loop():
    """Consume profile retraining requests until application shutdown."""

    while not training_worker_stop.is_set():
        try:
            profile = training_queue.get(timeout=0.5)
        except queue.Empty:
            continue
        if profile is None:
            training_queue.task_done()
            break

        with training_queue_lock:
            training_queued_profiles.discard(profile)
            training_running_profiles.add(profile)
        try:
            retrain_and_reload(profile)
        finally:
            with training_queue_lock:
                training_running_profiles.discard(profile)
                # Votes that arrived while this profile was training are
                # represented by one follow-up run over the complete JSON file.
                if profile in training_dirty_profiles and not training_worker_stop.is_set():
                    training_dirty_profiles.discard(profile)
                    training_queued_profiles.add(profile)
                    training_queue.put(profile)
            training_queue.task_done()


def _ensure_training_worker():
    global training_worker
    with training_queue_lock:
        if training_worker and training_worker.is_alive():
            return
        training_worker_stop.clear()
        training_worker = threading.Thread(
            target=_training_worker_loop,
            name="bouncer-training-worker",
            daemon=True,
        )
        training_worker.start()


def enqueue_bouncer_retrain(profile=DEFAULT_PROFILE):
    """Queue a profile retrain without losing rapid successive votes."""

    profile = profile if profile in PROFILE_CONFIGS else DEFAULT_PROFILE
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    _ensure_training_worker()
    with training_queue_lock:
        if profile in training_running_profiles:
            training_dirty_profiles.add(profile)
            return {"status": "queued", "profile": profile, "coalesced": True}
        if profile in training_queued_profiles:
            return {"status": "queued", "profile": profile, "coalesced": True}
        training_queued_profiles.add(profile)
        training_queue.put(profile)
    return {"status": "queued", "profile": profile, "coalesced": False}


# ==========================================
# --- USAGE TRACKING HELPERS ---
# ==========================================
def load_tracker():
    if not os.path.exists(USAGE_TRACKER_FILE):
        return {}
    try:
        with open(USAGE_TRACKER_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except:
        return {}


def save_tracker(data):
    temp_file = f"{USAGE_TRACKER_FILE}.tmp"
    with open(temp_file, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4, ensure_ascii=False)
    os.replace(temp_file, USAGE_TRACKER_FILE)


def get_device_id(ip, fingerprint):
    raw = f"{IP_HASH_SECRET}:{ip}:{fingerprint}"
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:16]


def get_viewer_key(ip):
    return hashlib.sha256(f"{IP_HASH_SECRET}:{ip}".encode("utf-8")).hexdigest()


def get_private_viewer_key(request: Request):
    """Use the signed browser identity, with IP hash as a legacy fallback."""

    return str(
        getattr(request.state, "private_viewer_key", "")
        or get_viewer_key(get_client_ip(request))
    )


def _private_item_identity(item):
    if not isinstance(item, dict):
        return str(item)
    return str(
        item.get("id")
        or item.get("article_key")
        or item.get("canonical_link")
        or item.get("link")
        or item.get("url")
        or item.get("title")
        or ""
    ).strip().casefold()


def _merge_private_buckets(current, legacy):
    """Merge legacy profile lists without duplicating saved items or jobs."""

    merged = dict(current) if isinstance(current, dict) else {}
    for profile, legacy_items in (legacy.items() if isinstance(legacy, dict) else []):
        destination = merged.setdefault(profile, [])
        if not isinstance(destination, list) or not isinstance(legacy_items, list):
            continue
        seen = {_private_item_identity(item) for item in destination}
        for item in legacy_items:
            identity = _private_item_identity(item)
            if identity and identity not in seen:
                destination.append(item)
                seen.add(identity)
    return merged


def claim_legacy_private_bucket(store: dict, request: Request):
    """Move an IP-keyed desk once to the first signed browser at that IP."""

    private_key = get_private_viewer_key(request)
    legacy_key = get_viewer_key(get_client_ip(request))
    if private_key == legacy_key:
        return private_key, False

    with viewer_identity_lock:
        claims = PRIVATE_VIEWER_CLAIMS.read()
        claims = claims if isinstance(claims, dict) else {}
        owner = claims.get(legacy_key)
        if owner is None:
            # Only a newly issued browser identity may reserve old data for
            # this IP. An existing identity that merely roamed onto a new IP
            # must never inherit that network's legacy private records.
            if not bool(getattr(request.state, "private_viewer_created", False)):
                return private_key, False
            claims[legacy_key] = private_key
            PRIVATE_VIEWER_CLAIMS.write(claims)
            owner = private_key

    if owner != private_key or legacy_key not in store:
        return private_key, False
    store[private_key] = _merge_private_buckets(
        store.get(private_key, {}),
        store.get(legacy_key, {}),
    )
    store.pop(legacy_key, None)
    return private_key, True


def migrate_tracker_privacy():
    """Replace legacy raw-IP analytics records with keyed hashes in place."""

    with tracker_lock:
        tracker = load_tracker()
        migrated = {}
        changed = False
        for old_device_id, device in tracker.items():
            record = dict(device or {})
            raw_ip = str(record.pop("ip", "") or "").strip()
            if raw_ip:
                record["ip_hash"] = get_viewer_key(raw_ip)
                record["owner"] = record.get("display_name") or get_team_owner_for_ip(raw_ip) or record.get("owner", "Unknown")
                record["known_team_member"] = bool(get_team_owner_for_ip(raw_ip))
                new_device_id = get_device_id(raw_ip, record.get("fingerprint", "unknown"))
                changed = True
            else:
                new_device_id = old_device_id
            migrated[new_device_id] = record
        if changed:
            save_tracker(migrated)
            print("[PRIVACY] Migrated legacy analytics records from raw IPs to keyed hashes.", flush=True)
        return changed


def load_viewer_profiles():
    if not os.path.exists(VIEWER_PROFILES_FILE):
        return {}
    try:
        with open(VIEWER_PROFILES_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
            return data if isinstance(data, dict) else {}
    except Exception:
        return {}


def save_viewer_profiles(data):
    temp_file = f"{VIEWER_PROFILES_FILE}.tmp"
    with open(temp_file, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    os.replace(temp_file, VIEWER_PROFILES_FILE)


def get_viewer_profile(ip):
    return load_viewer_profiles().get(get_viewer_key(ip), {})


def _article_identity(article):
    """Return a stable, profile-independent identity for a stored article."""

    if not isinstance(article, dict):
        return ""
    link = str(
        article.get("canonical_link")
        or article.get("link")
        or article.get("url")
        or ""
    ).strip().casefold()
    title = str(article.get("title") or "").strip().casefold()
    raw = link or title
    return hashlib.sha256(raw.encode("utf-8")).hexdigest() if raw else ""


def load_viewer_hidden_store():
    if not os.path.exists(VIEWER_HIDDEN_FILE):
        return {}
    try:
        with open(VIEWER_HIDDEN_FILE, "r", encoding="utf-8") as file_obj:
            data = json.load(file_obj)
            return data if isinstance(data, dict) else {}
    except Exception:
        return {}


def save_viewer_hidden_store(data):
    temp_file = f"{VIEWER_HIDDEN_FILE}.tmp"
    with open(temp_file, "w", encoding="utf-8") as file_obj:
        json.dump(data, file_obj, indent=2, ensure_ascii=False)
    os.replace(temp_file, VIEWER_HIDDEN_FILE)


def load_viewer_saved_store():
    data = JsonStore(Path(VIEWER_SAVED_FILE), dict).read()
    return data if isinstance(data, dict) else {}


def save_viewer_saved_store(data):
    return JsonStore(Path(VIEWER_SAVED_FILE), dict).write(data)


def get_viewer_saved_items(request, profile=None):
    profile_name = profile or get_profile_for_request(request)
    with viewer_saved_lock:
        store = load_viewer_saved_store()
        viewer_key, migrated = claim_legacy_private_bucket(store, request)
        if migrated:
            save_viewer_saved_store(store)
        viewer_store = store.get(viewer_key, {})
        items = viewer_store.get(profile_name, [])
    return items if isinstance(items, list) else []


def get_viewer_hidden_items(request, profile=None):
    profile_name = profile or get_profile_for_request(request)
    viewer_key = get_viewer_key(get_client_ip(request))
    store = load_viewer_hidden_store()
    viewer_store = store.get(viewer_key, {})
    items = viewer_store.get(profile_name, [])
    return items if isinstance(items, list) else []


def filter_viewer_hidden(items, request, profile=None):
    hidden_keys = {
        str(item.get("article_key") or _article_identity(item))
        for item in get_viewer_hidden_items(request, profile)
    }
    return [
        item for item in (items or [])
        if _article_identity(item) not in hidden_keys
    ]


def record_usage_best_effort(ip, profile, fingerprint, action, detail=""):
    """Never turn a committed private desk action into an apparent failure."""

    try:
        return record_usage_activity(
            ip, profile, fingerprint, action, detail
        )
    except Exception as error:
        print(
            f"[USAGE] {action} tracking failed after the desk action committed: "
            f"{type(error).__name__}: {error}",
            flush=True,
        )
        return False


def record_recommendation_best_effort(
    request,
    response,
    action,
    detail,
    *,
    event_id="",
    occurred_at="",
    active_ms=0,
    visible_ratio=0.0,
):
    """Bridge a validated shared-surface action without breaking its UI task."""

    try:
        from news_scrapper.recommendation.router import record_shared_briefing_event

        return record_shared_briefing_event(
            request,
            response,
            track_action=action,
            detail=detail if isinstance(detail, dict) else {},
            event_id=event_id,
            occurred_at=occurred_at,
            active_ms=active_ms,
            visible_ratio=visible_ratio,
        )
    except Exception as error:
        print(
            f"[FOR YOU BRIDGE] Could not record {action}: "
            f"{type(error).__name__}: {error}",
            flush=True,
        )
        return {"accepted": 0, "duplicates": 0, "rejected": 1, "ignored": True}


def get_today():
    return datetime.datetime.now().strftime("%Y-%m-%d")


def get_empty_day():
    return {
        "page_loads": 0,
        "searches": [],
        "articles_clicked": 0,
        "votes_interested": 0,
        "votes_not_interested": 0,
        "saved_for_later": 0,
        "removed_from_saved": 0,
        "exports": [],
        "briefing_views": 0,
        "heartbeats": 0,
        "voc_feedback": [],
        "selections": 0,
        "approvals": 0,
        "personal_hides": 0,
        "workflow_removals": 0,
        "sources_added": 0,
        "action_counts": {},
        "events": [],
    }


def purge_old_entries(device_data, keep_days=30):
    cutoff = (
        datetime.datetime.now() - datetime.timedelta(days=keep_days)
    ).strftime("%Y-%m-%d")
    activity = device_data.get("activity", {})
    device_data["activity"] = {
        date: data for date, data in activity.items()
        if date >= cutoff
    }
    return device_data


# ==========================================
# Legacy default-only scheduler retained for archive compatibility; the
# profile-aware scheduler definition below is the runtime entry point.
# ==========================================
def _legacy_run_morning_briefing():
    global SCHEDULER_STATUS

    with scheduler_lock:
        if SCHEDULER_STATUS["is_active"]:
            print("[SCHEDULER] Already running. Skipping.")
            return

        running_manual = sum(
            1 for j in active_jobs.values() if j.get("status") == "running"
        )
        if running_manual > 0:
            print(
                f"[SCHEDULER] {running_manual} manual scan(s) in progress. "
                f"Deferring autonomous run by 10 minutes."
            )
            threading.Timer(600, _legacy_run_morning_briefing).start()
            return

        SCHEDULER_STATUS["is_active"] = True
        SCHEDULER_STATUS["message"] = "Autonomous Engine Scanning..."
        SCHEDULER_STATUS["mode"] = "autonomous"

    scheduler_job_id = (
        f"scheduler_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    )

    try:
        today = datetime.datetime.now()
        yesterday = today - datetime.timedelta(days=1)
        from_date = yesterday.strftime("%Y-%m-%d")
        to_date = today.strftime("%Y-%m-%d")

        root_dir = ROOT_DIR
        output_file = os.path.join(
            root_dir, f"ui_results_{scheduler_job_id}.json"
        )
        cluster_file = os.path.join(
            root_dir, f"clustered_results_{scheduler_job_id}.json"
        )

        for pattern in [
            "ui_results_scheduler_*.json",
            "clustered_results_scheduler_*.json",
        ]:
            for old_file in glob.glob(os.path.join(root_dir, pattern)):
                try:
                    os.remove(old_file)
                except:
                    pass

        print(f"[SCHEDULER] Starting scan: {scheduler_job_id}")
        print(
            f"   -> Deploying Spiders for keywords: {MORNING_KEYWORDS[:50]}..."
        )

        cmd = [
            sys.executable, "-m", "scrapy", "crawl", "news_spider",
            "-a", f"keyword={MORNING_KEYWORDS}",
            "-a", f"from_date={from_date}",
            "-a", f"to_date={to_date}",
            "-a", "target_sites=All",
            "-a", f"discovery_only={'true' if SAMSUNG_DISCOVERY_ONLY else 'false'}",
            "-s", f"ROBOTSTXT_OBEY={'True' if SCRAPY_ROBOTSTXT_OBEY else 'False'}",
            "-s", "USER_AGENT=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "-s", "TWISTED_REACTOR=twisted.internet.asyncioreactor.AsyncioSelectorReactor",
            "-O", output_file,
        ]
        spider_cwd = str(NEWS_CRAWLER_DIR)
        subprocess.run(cmd, cwd=spider_cwd, timeout=3600)

        if os.path.exists(output_file):
            try:
                print("   -> [SCHEDULER] Running AI Gatekeeper...")
                with open(output_file, "r", encoding="utf-8") as f:
                    raw_data = json.load(f)

                filtered_data = []
                dropped_count = 0

                for item in raw_data:
                    title = item.get("title", "")
                    summary = item.get("snippet", item.get("summary", ""))
                    keywords_found = item.get("keywords_found", [])

                    if bouncer_check(title, summary, keywords_found):
                        filtered_data.append(item)
                    else:
                        dropped_count += 1
                        print(f"   Dropped: {title}", flush=True)
                        log_dropped_article(title, summary, keywords_found)

                with open(output_file, "w", encoding="utf-8") as f:
                    json.dump(filtered_data, f, indent=4)

                print(
                    f"   -> [SCHEDULER] Gatekeeper removed {dropped_count} articles.",
                    flush=True,
                )
            except Exception as e:
                print(
                    f"   -> [SCHEDULER] Bouncer error, skipping filter: {e}"
                )

        print("   -> Activating Semantic Fusion Engine (Fast Mode)...")
        subprocess.run(
            [
                sys.executable, "-m", "news_scrapper.semantic_clustering",
                "--job-id", scheduler_job_id,
                "--fast-mode",
            ],
            cwd=str(PROJECT_ROOT),
            timeout=3600,
        )

        if os.path.exists(cluster_file):
            with open(cluster_file, "r", encoding="utf-8") as f:
                results = json.load(f)

            if results:
                for r in results:
                    r["category"] = assign_category(
                        r.get("title", ""),
                        r.get("master_summary", "") or r.get("snippet", ""),
                    )
                    r.update(apply_learned_region(r))

                print(f"   -> Archiving {len(results)} intelligence items...")
                learner.log_search_data(MORNING_KEYWORDS, results)

                timestamp = today.strftime("%Y-%m-%d_%H-%M-%S")
                history_path = os.path.join(
                    HISTORY_DIR, f"briefing_{timestamp}.json"
                )
                with open(history_path, "w", encoding="utf-8") as f:
                    json.dump(results, f, indent=4)

                try:
                    df = pd.DataFrame(results)
                    df["search_timestamp"] = today
                    df["search_keywords"] = "MORNING_BRIEFING"
                    df["global_week"] = today.isocalendar()[1]

                    if os.path.exists(MANUAL_LOG_FILE):
                        with pd.ExcelWriter(
                            MANUAL_LOG_FILE,
                            mode="a",
                            if_sheet_exists="overlay",
                            engine="openpyxl",
                        ) as writer:
                            pd.concat(
                                [pd.read_excel(MANUAL_LOG_FILE), df],
                                ignore_index=True,
                            ).to_excel(writer, index=False)
                    else:
                        df.to_excel(MANUAL_LOG_FILE, index=False)
                except Exception as e:
                    print(f"Excel Logging Error: {e}")

                print("[SCHEDULER] Morning Briefing Complete & Ready.")
            else:
                print("[SCHEDULER] Briefing finished but no news was found.")
        else:
            print("[SCHEDULER] Failed to generate cluster file.")

    except subprocess.TimeoutExpired:
        print("\n[SCHEDULER ERROR] Process timed out and was terminated.")
    except Exception as e:
        print(f"\n[SCHEDULER ERROR] Critical failure: {e}")
    finally:
        with scheduler_lock:
            SCHEDULER_STATUS["is_active"] = False
            SCHEDULER_STATUS["message"] = "Morning Briefing Complete."
            SCHEDULER_STATUS["mode"] = "idle"
        cleanup_job_files(output_file, cluster_file)
        print(f"[SCHEDULER] Cleaned up temp files for {scheduler_job_id}")


# ==========================================
# PROFILE-AWARE AUTONOMOUS SCHEDULER
# ==========================================
def run_scheduler_for_profile(profile: str, capabilities=None):
    if UNIFIED_CORPUS_ENABLED and profile != UNIFIED_PROFILE:
        print(
            f"[SCHEDULER] Ignoring legacy '{profile}' partition; the unified "
            "catalog is scanned once per cycle.",
            flush=True,
        )
        return True
    if UNIFIED_CORPUS_ENABLED:
        profile = UNIFIED_PROFILE
    config = get_profile_config(profile)
    capabilities = capabilities or resolve_pipeline_capabilities()
    today = datetime.datetime.now()
    run_label = "unified" if UNIFIED_CORPUS_ENABLED else profile
    scheduler_job_id = f"scheduler_{run_label}_{today.strftime('%Y%m%d_%H%M%S')}"
    output_file = os.path.join(ROOT_DIR, f"ui_results_{scheduler_job_id}.json")
    cluster_file = os.path.join(ROOT_DIR, f"clustered_results_{scheduler_job_id}.json")
    history_dir = get_profile_history_dir(profile)

    if not os.path.exists(config["sites_file"]):
        print(f"[SCHEDULER:{profile}] Missing sites file: {config['sites_file']}")
        return False

    try:
        from_date = (today - datetime.timedelta(days=CRAWL_LOOKBACK_DAYS - 1)).strftime("%Y-%m-%d")
        to_date = today.strftime("%Y-%m-%d")
        cmd = [
            sys.executable, "-m", "scrapy", "crawl", "news_spider",
            "-a", f"keyword={config['keywords']}",
            "-a", f"from_date={from_date}",
            "-a", f"to_date={to_date}",
            "-a", "target_sites=All",
            "-a", f"sites_file={config['sites_file']}",
            "-a", f"discovery_only={'true' if capabilities['discovery_only'] else 'false'}",
            "-s", f"ROBOTSTXT_OBEY={'True' if SCRAPY_ROBOTSTXT_OBEY else 'False'}",
            "-s", "USER_AGENT=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "-s", "TWISTED_REACTOR=twisted.internet.asyncioreactor.AsyncioSelectorReactor",
            "-O", output_file,
        ]
        print(
            f"[SCHEDULER:{profile}] Starting scan: {scheduler_job_id} "
            f"(pipeline={capabilities['mode']}, "
            f"scrapy={'discovery-only' if capabilities['discovery_only'] else 'full extraction'})",
            flush=True,
        )
        if scheduler_shutdown_event.is_set():
            print(f"[SCHEDULER:{profile}] Shutdown requested; scan skipped.", flush=True)
            return False
        def execute_crawl(command):
            process = subprocess.Popen(
                command,
                cwd=str(NEWS_CRAWLER_DIR),
            )
            try:
                with scheduler_process_lock:
                    scheduler_processes[profile] = process
                try:
                    return_code = process.wait(timeout=3600)
                except subprocess.TimeoutExpired:
                    process.terminate()
                    try:
                        process.wait(timeout=10)
                    except subprocess.TimeoutExpired:
                        process.kill()
                        process.wait(timeout=5)
                    raise
            finally:
                with scheduler_process_lock:
                    scheduler_processes.pop(profile, None)
            if return_code != 0:
                raise RuntimeError(f"Scrapy exited with code {return_code}")

        execute_crawl(cmd)

        if not os.path.exists(output_file):
            raise RuntimeError("Scrapy completed without creating an output file")

        with open(output_file, "r", encoding="utf-8") as f:
            raw_data = json.load(f)
        if not isinstance(raw_data, list):
            raise RuntimeError("Scrapy output must be a JSON list")

        web_search_used = bool(capabilities["web_search"])
        if web_search_used:
            print(
                f"[SCHEDULER:{profile}] Sending discovered URLs to Samsung "
                "Web Search before Gatekeeper scoring.",
                flush=True,
            )
            try:
                raw_data = enrich_raw_articles(
                    raw_data,
                    config["keywords"],
                    profile,
                    use_web_search=True,
                    raise_on_service_failure=True,
                )
            except WebSearchRuntimeFailure as error:
                print(
                    f"[SCHEDULER:{profile}] WEB SEARCH RUNTIME FAILURE: {error} "
                    "FALLBACK: restarting this profile with full Scrapy "
                    "article extraction.",
                    flush=True,
                )
                fallback_cmd = [
                    (
                        "discovery_only=false"
                        if str(argument).startswith("discovery_only=")
                        else argument
                    )
                    for argument in cmd
                ]
                execute_crawl(fallback_cmd)
                if not os.path.exists(output_file):
                    raise RuntimeError(
                        "Full Scrapy fallback completed without creating an "
                        "output file"
                    )
                with open(output_file, "r", encoding="utf-8") as f:
                    raw_data = json.load(f)
                if not isinstance(raw_data, list):
                    raise RuntimeError(
                        "Full Scrapy fallback output must be a JSON list"
                    )
                web_search_used = False

        if web_search_used:
            if config["use_bouncer"]:
                raw_data, dropped, low_priority = run_bouncer_filter_on_items(
                    raw_data, profile, "scheduler_enriched"
                )
                print(
                    f"[SCHEDULER:{profile}] Dropped {dropped}; "
                    f"low priority kept {low_priority}.",
                    flush=True,
                )
        else:
            if config["use_bouncer"]:
                raw_data, dropped, low_priority = run_bouncer_filter_on_items(
                    raw_data, profile, "scheduler_raw"
                )
                print(
                    f"[SCHEDULER:{profile}] Dropped {dropped}; "
                    f"low priority kept {low_priority}.",
                    flush=True,
                )
            raw_data = enrich_raw_articles(
                raw_data,
                config["keywords"],
                profile,
                use_web_search=False,
            )
        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(raw_data, f, indent=4, ensure_ascii=False)

        clustering_command = [
            sys.executable,
            "-m",
            "news_scrapper.semantic_clustering",
            "--job-id",
            scheduler_job_id,
        ]
        if capabilities["chat"]:
            clustering_command.append("--fast-mode")
            print(
                f"[SCHEDULER:{profile}] MiniLM clustering in fast mode; "
                "Samsung Chat owns final summarization.",
                flush=True,
            )
        else:
            print(
                f"[SCHEDULER:{profile}] MiniLM clustering with local BART "
                "summarization fallback.",
                flush=True,
            )
        clustering_result = subprocess.run(
            clustering_command,
            cwd=str(PROJECT_ROOT),
            timeout=3600,
            check=False,
        )
        if clustering_result.returncode != 0:
            raise RuntimeError(
                "Semantic clustering exited with code "
                f"{clustering_result.returncode}"
            )
        if not os.path.exists(cluster_file):
            print(f"[SCHEDULER:{profile}] Failed to generate cluster file.")
            return False
        with open(cluster_file, "r", encoding="utf-8") as f:
            results = json.load(f)
        if results and config["use_bouncer"]:
            results, dropped, low_priority = run_bouncer_filter_on_items(
                results, profile, "scheduler_final"
            )
        for item in results:
            verticals = item.get("verticals") or []
            if isinstance(verticals, str):
                verticals = [verticals]
            vertical = str(item.get("vertical") or (verticals[0] if verticals else "")).strip().lower()
            legacy_profile = str(item.get("legacy_profile") or item.get("keyword_pack") or "").strip().lower()
            if vertical not in {"technology", "broadcast"}:
                vertical = "broadcast" if legacy_profile == BROADCAST_PROFILE else "technology"
            item["profile"] = "unified" if UNIFIED_CORPUS_ENABLED else profile
            item["legacy_profile"] = legacy_profile or (
                BROADCAST_PROFILE if vertical == "broadcast" else DEFAULT_PROFILE
            )
            item["vertical"] = vertical
            item["verticals"] = list(dict.fromkeys([*verticals, vertical]))
            item["audiences"] = item.get("audiences") or ["all"]
            item["category"] = assign_category(
                item.get("title", ""),
                item.get("master_summary", "") or item.get("snippet", ""),
            )
            item.update(apply_learned_region(item, profile))
        results = enrich_final_articles(
            results,
            profile,
            use_chat=capabilities["chat"],
        )
        if not results:
            print(f"[SCHEDULER:{profile}] Finished but no news was found.")
            return True
        learner.log_search_data(config["keywords"], results)
        timestamp = today.strftime("%Y-%m-%d_%H-%M-%S")
        history_path = os.path.join(history_dir, f"briefing_{timestamp}.json")
        with open(history_path, "w", encoding="utf-8") as f:
            json.dump(results, f, indent=4, ensure_ascii=False)
        purge_expired_history(profile)
        print(
            f"[SCHEDULER:{run_label}] Unified briefing complete: {history_path}",
            flush=True,
        )
        return True
    except subprocess.TimeoutExpired:
        print(f"[SCHEDULER:{profile}] Process timed out and was terminated.")
    except Exception as e:
        print(f"[SCHEDULER:{profile}] Critical failure: {e}")
    finally:
        cleanup_job_files(output_file, cluster_file)
    return False


SCHEDULER_RETRY_DELAY_SECONDS = max(
    30,
    int(os.environ.get("SCHEDULER_RETRY_DELAY_SECONDS", "600")),
)


def _schedule_scheduler_retry(delay_seconds: int = SCHEDULER_RETRY_DELAY_SECONDS):
    """Schedule one deferred retry, never an unbounded timer chain."""

    global scheduler_retry_scheduled
    with scheduler_lock:
        if scheduler_shutdown_event.is_set() or scheduler_retry_scheduled:
            return
        scheduler_retry_scheduled = True

    def retry():
        global scheduler_retry_scheduled
        with scheduler_lock:
            scheduler_retry_scheduled = False
        run_morning_briefing()

    timer = threading.Timer(delay_seconds, retry)
    timer.daemon = True
    timer.start()


def run_morning_briefing():
    global SCHEDULER_STATUS, scheduler_pending_run
    run_id = f"cycle_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    scheduled_at = datetime.datetime.now(datetime.timezone.utc).isoformat(
        timespec="seconds"
    )
    defer_for_manual = False
    with scheduler_lock:
        if scheduler_shutdown_event.is_set():
            return
        if SCHEDULER_STATUS["is_active"]:
            # The APScheduler tick must not disappear while a long scan runs.
            scheduler_pending_run = True
            return
        if any(
            job.get("status") in {"queued", "running"}
            for job in active_jobs.values()
        ):
            scheduler_pending_run = True
            defer_for_manual = True
        else:
            SCHEDULER_STATUS.update(
                {
                    "is_active": True,
                    "message": "Autonomous Engine Scanning...",
                    "mode": "autonomous",
                    "last_started_at": datetime.datetime.now().isoformat(timespec="seconds"),
                    "last_error": None,
                }
            )
    if defer_for_manual:
        print("[SCHEDULER] Manual scan admitted first. Deferring autonomous scan.", flush=True)
        update_durable_scheduler_state(
            run_id=run_id,
            scheduled_at=scheduled_at,
            status="deferred",
            stage="waiting_for_manual_scan",
            retry_state="scheduled",
            completed_partitions=[],
            failures=[],
        )
        _schedule_scheduler_retry()
        return

    failed_profiles = []
    completed_profiles = []
    update_durable_scheduler_state(
        run_id=run_id,
        scheduled_at=scheduled_at,
        status="running",
        stage="capability_preflight",
        completed_partitions=[],
        failures=[],
        retry_state="none",
        publish_timestamp=None,
    )
    try:
        ensure_profile_storage()
        capabilities = resolve_pipeline_capabilities()
        serving_profiles = (
            [UNIFIED_PROFILE]
            if UNIFIED_CORPUS_ENABLED
            else [DEFAULT_PROFILE, BROADCAST_PROFILE]
        )
        for profile in serving_profiles:
            if scheduler_shutdown_event.is_set():
                break
            with scheduler_lock:
                SCHEDULER_STATUS["message"] = f"Autonomous Engine Scanning {profile} profile..."
            update_durable_scheduler_state(
                stage=f"processing_{profile}",
                active_partition=profile,
                completed_partitions=completed_profiles,
            )
            if not run_scheduler_for_profile(profile, capabilities):
                failed_profiles.append(profile)
                update_durable_scheduler_state(
                    failures=[
                        *failed_profiles,
                    ],
                    stage=f"failed_{profile}",
                )
            else:
                completed_profiles.append(profile)
                update_durable_scheduler_state(
                    completed_partitions=completed_profiles,
                    stage=f"completed_{profile}",
                )
        try:
            from news_scrapper.recommendation.router import process_reaction_consensus
            consensus = process_reaction_consensus()
            print(f"[FOR YOU] Reaction consensus batch: {consensus}", flush=True)
        except Exception as consensus_error:
            print(f"[FOR YOU] Reaction consensus batch failed safely: {consensus_error}", flush=True)
        if (
            not UNIFIED_CORPUS_ENABLED
            and
            UNIFIED_CORPUS_SHADOW_ENABLED
            and not failed_profiles
            and set(completed_profiles) == {DEFAULT_PROFILE, BROADCAST_PROFILE}
        ):
            update_durable_scheduler_state(stage="publishing_unified_shadow")
            shadow_report = publish_unified_shadow(run_id)
            update_durable_scheduler_state(
                stage="unified_shadow_published",
                shadow_report=shadow_report,
                publish_timestamp=datetime.datetime.now(
                    datetime.timezone.utc
                ).isoformat(timespec="seconds"),
            )
            print(
                "[SCHEDULER] Unified shadow briefing and parity report published; "
                "legacy serving remains active.",
                flush=True,
            )
    except Exception as error:
        failed_profiles = failed_profiles or ["orchestration"]
        update_durable_scheduler_state(
            status="failed",
            stage="orchestration_failed",
            failures=failed_profiles,
            last_error=str(error),
        )
        print(f"[SCHEDULER] Orchestration failure: {error}", flush=True)
    finally:
        with scheduler_lock:
            pending = scheduler_pending_run or bool(failed_profiles)
            scheduler_pending_run = False
            completed_at = datetime.datetime.now()
            if failed_profiles:
                next_due = completed_at + datetime.timedelta(
                    seconds=SCHEDULER_RETRY_DELAY_SECONDS
                )
            elif pending:
                next_due = completed_at + datetime.timedelta(seconds=1)
            else:
                started_value = SCHEDULER_STATUS.get("last_started_at")
                try:
                    next_due = datetime.datetime.fromisoformat(started_value)
                except (TypeError, ValueError):
                    next_due = completed_at
                next_due += datetime.timedelta(hours=4)
                while next_due <= completed_at:
                    next_due += datetime.timedelta(hours=4)
            SCHEDULER_STATUS.update(
                {
                    "is_active": False,
                    "message": "Morning Briefing Complete.",
                    "mode": "idle",
                    "last_completed_at": completed_at.isoformat(timespec="seconds"),
                    "next_run": next_due.isoformat(timespec="seconds"),
                    "last_profiles": (
                        ["unified"]
                        if UNIFIED_CORPUS_ENABLED
                        else [DEFAULT_PROFILE, BROADCAST_PROFILE]
                    ),
                    "last_failed_profiles": failed_profiles,
                    "last_error": (
                        f"Retry scheduled for: {', '.join(failed_profiles)}"
                        if failed_profiles
                        else None
                    ),
                }
            )
        update_durable_scheduler_state(
            status="failed" if failed_profiles else "complete",
            stage="retry_wait" if failed_profiles else "complete",
            active_partition=None,
            completed_partitions=completed_profiles,
            failures=failed_profiles,
            retry_state="scheduled" if pending else "none",
            completed_at=completed_at.astimezone(datetime.timezone.utc).isoformat(
                timespec="seconds"
            ),
            next_due_at=next_due.astimezone().isoformat(timespec="seconds"),
        )
        if pending and not scheduler_shutdown_event.is_set():
            _schedule_scheduler_retry(1 if not failed_profiles else SCHEDULER_RETRY_DELAY_SECONDS)


# ==========================================
# --- LIFECYCLE ---
# ==========================================
@asynccontextmanager
async def lifespan(app: FastAPI):
    ensure_runtime_directories()
    ensure_profile_storage()
    migrate_tracker_privacy()
    resume_personal_briefing_jobs()
    if UNIFIED_MIGRATION_REPORT.get("model_needs_retrain"):
        print(
            "[BOUNCER] Legacy feedback was merged; queuing one unified model retrain.",
            flush=True,
        )
        enqueue_bouncer_retrain(UNIFIED_PROFILE)
    scheduler_shutdown_event.clear()
    scheduler = BackgroundScheduler()
    next_run = datetime.datetime.now() + datetime.timedelta(minutes=2)

    if SCHEDULER_ENABLED:
        now = datetime.datetime.now()
        due_times = []
        due_profiles = (
            [UNIFIED_PROFILE]
            if UNIFIED_CORPUS_ENABLED
            else [DEFAULT_PROFILE, BROADCAST_PROFILE]
        )
        for profile in due_profiles:
            latest_file = get_latest_briefing_file_for_profile(profile)
            if not latest_file:
                print(f"[SCHEDULER:{profile}] No briefing exists; profile is due after startup.", flush=True)
                due_times.append(now + datetime.timedelta(minutes=2))
                continue
            last_run = datetime.datetime.fromtimestamp(os.path.getmtime(latest_file))
            due_times.append(max(now + datetime.timedelta(seconds=5), last_run + datetime.timedelta(hours=4)))
        if due_times:
            # Unified mode has one authoritative briefing timestamp. The
            # legacy loop remains only for an explicit rollback configuration.
            next_run = min(due_times)
        SCHEDULER_STATUS["next_run"] = next_run.isoformat(timespec="seconds")
        print(f"[SCHEDULER] Next unified scan: {next_run.strftime('%I:%M %p')}.", flush=True)
        scheduler.add_job(
            run_morning_briefing,
            "interval",
            hours=4,
            next_run_time=next_run,
            id="unified_briefing",
            max_instances=1,
            # If a Windows host sleeps through several four-hour ticks, run one
            # catch-up scan when it wakes instead of replaying every missed tick.
            coalesce=True,
            misfire_grace_time=24 * 60 * 60,
        )
        scheduler.start()
        print("SYSTEM: Autonomous Intelligence Engine online.")
    else:
        SCHEDULER_STATUS["next_run"] = None
        print("SYSTEM: Autonomous scheduler disabled by configuration.")
    yield
    scheduler_shutdown_event.set()
    with scheduler_process_lock:
        running_processes = list(scheduler_processes.values())
    for process in running_processes:
        if process.poll() is None:
            process.terminate()
    if scheduler.running:
        scheduler.shutdown(wait=False)
    training_worker_stop.set()
    training_queue.put(None)
    ml_executor.shutdown(wait=False)
    gatekeeper_executor.shutdown(wait=False)
    personal_briefing_executor.shutdown(wait=False)


# ==========================================
# --- APP INIT ---
# ==========================================
app = FastAPI(lifespan=lifespan)
cors_allowed_origins = sorted(env_csv("NEWSSCRAPPER_CORS_ALLOWED_ORIGINS"))
if cors_allowed_origins:
    app.add_middleware(
        CORSMiddleware,
        allow_origins=cors_allowed_origins,
        allow_credentials=True,
        allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
        allow_headers=["Content-Type", "X-Sense-Profile"],
    )


@app.middleware("http")
async def bind_private_viewer(request: Request, call_next):
    """Give every private viewer API a signed, browser-scoped identity."""

    identity = None
    if request.url.path == "/viewer" or request.url.path.startswith("/viewer/"):
        request.state.private_viewer_cookie_middleware = True
        identity = bind_viewer_request(request)
        if identity[1]:
            # Reserve any IP-keyed legacy desk for the first browser identity
            # created on that IP. The individual Saved/Briefing stores move
            # their bucket lazily when their endpoint is opened.
            claim_legacy_private_bucket({}, request)
    response = await call_next(request)
    if identity and identity[1]:
        set_viewer_cookie(request, response, identity[2])
    return response

# ==========================================
# --- STATIC FILES ---
# ==========================================
abs_frontend_path = str(FRONTEND_DIST)
if os.path.exists(os.path.join(abs_frontend_path, "assets")):
    app.mount(
        "/assets",
        StaticFiles(directory=os.path.join(abs_frontend_path, "assets")),
        name="assets",
    )


# ==========================================
# PROFILE ENDPOINT
# ==========================================
@app.get("/profile")
def get_profile(request: Request):
    profile = get_active_profile_name(request)
    config = get_profile_config(profile)
    response = {
        "status": "success",
        "ip": get_client_ip(request),
        "profile": public_profile_name(profile),
        "label": config["label"],
        "is_broadcast": False if UNIFIED_CORPUS_ENABLED else profile == BROADCAST_PROFILE,
        "corpus": "unified" if UNIFIED_CORPUS_ENABLED else "partitioned",
    }
    # Local storage paths and cross-profile configuration are operational
    # diagnostics, not normal viewer data.
    if get_client_ip(request) in SYSTEM_STATUS_ALLOWED_IPS:
        response["paths"] = {
            "sites_file": config["sites_file"],
            "history_dir": config["history_dir"],
            "workflow_file": get_workflow_file_for_request(request),
            "not_interested_file": get_not_interested_file_for_profile(profile),
            "training_file": get_training_file_for_profile(profile),
            "bouncer_model_file": get_bouncer_model_file_for_profile(profile),
        }
        response["all_profiles"] = (
            {"unified": get_profile_debug_info(UNIFIED_PROFILE)}
            if UNIFIED_CORPUS_ENABLED
            else {
                DEFAULT_PROFILE: get_profile_debug_info(DEFAULT_PROFILE),
                BROADCAST_PROFILE: get_profile_debug_info(BROADCAST_PROFILE),
            }
        )
    return response


@app.get("/trends/access")
def get_profile_settings_access(request: Request):
    ip = get_client_ip(request)
    return {
        "allowed": ip in PROFILE_SETTINGS_ALLOWED_IPS,
        "ip": ip,
        "active_profile": get_active_profile_name(request),
        "owner": get_viewer_profile(ip).get("display_name") or get_team_owner_for_ip(ip) or "Unknown",
    }


# ==========================================
# --- EXPORT: POWERPOINT ---
# ==========================================
@app.post("/export-ppt")
async def export_ppt(http_request: Request, request: ExportRequest):
    if get_active_profile_name(http_request) == BROADCAST_PROFILE:
        raise HTTPException(
            status_code=403,
            detail="PowerPoint export is disabled for Broadcast profile.",
        )
    safe_filename = sanitize_filename(request.filename)
    TEMPLATE_PATH = "template.pptx"
    if not os.path.exists(TEMPLATE_PATH):
        raise HTTPException(status_code=404, detail="template.pptx not found")

    prs = Presentation(TEMPLATE_PATH)
    cover_layout = None
    news_layout = None
    for layout in prs.slide_layouts:
        if layout.name == "CoverLayout":
            cover_layout = layout
        if layout.name == "NewsLayout":
            news_layout = layout

    if not cover_layout:
        cover_layout = prs.slide_layouts[0]
    if not news_layout:
        news_layout = prs.slide_layouts[0] if len(prs.slide_layouts) == 1 else prs.slide_layouts[1]

    ph_map = {"title": -1, "summary": -1, "link": -1, "insight": -1, "picture": -1, "date": -1, "team": -1}
    for shape in news_layout.placeholders:
        idx = shape.placeholder_format.idx
        if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            ph_map["picture"] = idx
            continue
        if shape.has_text_frame:
            text = shape.text.strip()
            if "#TITLE" in text: ph_map["title"] = idx
            elif "#SUMMARY" in text: ph_map["summary"] = idx
            elif "#LINK" in text: ph_map["link"] = idx
            elif "#INSIGHT" in text: ph_map["insight"] = idx
            elif "#DATE_HERE" in text: ph_map["date"] = idx
            elif "#Targated_SRID_Team" in text: ph_map["team"] = idx

    slide = prs.slides.add_slide(cover_layout)
    for shape in slide.shapes:
        if shape.has_text_frame and "#DATE_HERE" in shape.text:
            shape.text = datetime.datetime.now().strftime("%b'%y")

    for item in request.items:
        slide = prs.slides.add_slide(news_layout)
        ai_opinion = generate_opinion(item.master_summary)
        target_team = determine_target_team(item.title, item.master_summary)

        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == ph_map["title"]:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                kw = item.keywords_found[0] if item.keywords_found else ""
                if kw:
                    parts = re.split(f"({re.escape(kw)})", item.title, flags=re.IGNORECASE)
                    for part in parts:
                        run = p.add_run()
                        run.text = part
                        if part.lower() == kw.lower():
                            run.font.bold = True
                            run.font.underline = True
                            run.font.color.rgb = RGBColor(0, 112, 192)
                        else:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0, 0, 0)
                else:
                    run = p.add_run()
                    run.text = item.title
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 0, 0)
            elif idx == ph_map["summary"]:
                tf = shape.text_frame
                tf.clear()
                target_summary = item.ppt_summary if item.ppt_summary else item.master_summary
                sentences = [s.strip() + "." for s in target_summary.split(". ") if s.strip()]
                if not sentences: sentences = [target_summary]
                for i, sentence in enumerate(sentences):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = f"• {sentence}"
                    p.font.name = "Calibri"
                    p.font.size = Pt(18)
            elif idx == ph_map["link"]:
                shape.text = item.link
                try: shape.text_frame.paragraphs[0].font.size = Pt(10)
                except: pass
            elif idx == ph_map["insight"]:
                shape.text = f"Insight : {ai_opinion}"
                try: shape.text_frame.paragraphs[0].font.size = Pt(14)
                except: pass
            elif idx == ph_map["date"]:
                shape.text = datetime.datetime.now().strftime("%b'%y")
            elif idx == ph_map["team"]:
                shape.text = f"Targeted SRID TEAM : {target_team}"
            elif idx == ph_map["picture"] or shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                img_stream = download_image_for_export(item.top_image, add_border=False)
                if img_stream:
                    try:
                        left, top = shape.left, shape.top
                        width, height = shape.width, shape.height
                        pic = slide.shapes.add_picture(img_stream, left, top, width, height)
                        try: pic.click_action.hyperlink.address = item.link
                        except: pass
                    except Exception as e:
                        print(f"   [PPT] Insert Error: {e}")

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={safe_filename}"},
    )


# ==========================================
# --- EXPORT: EXCEL ---
# ==========================================
@app.post("/export-excel")
async def export_excel(request: ExportRequest):
    safe_filename = sanitize_filename(request.filename)
    wb = Workbook()
    ws = wb.active
    ws.title = "SENSE Report"
    ws["A2"] = "Weekly Report"
    ws["A2"].font = Font(size=20, bold=True)

    headers = ["Sr. No.", "WK", "Date", "Name of the Initiator", "Category", "Keyword", "News Highlight", "URL"]
    header_fill = PatternFill(start_color="ADD8E6", end_color="ADD8E6", fill_type="solid")
    for col_num, header_title in enumerate(headers, 1):
        cell = ws.cell(row=4, column=col_num, value=header_title)
        cell.font = Font(bold=True)
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")

    for idx, item in enumerate(request.items, 1):
        row_num = 4 + idx
        try:
            dt = datetime.datetime.strptime(item.date, "%Y-%m-%d")
            wk = dt.isocalendar()[1]
        except:
            wk = datetime.datetime.now().isocalendar()[1]

        initiator = getattr(item, "selected_by", "")
        category_val = getattr(item, "category", "Tech News")
        kw = ", ".join(item.keywords_found) if item.keywords_found else ""
        summary_sentences = [s.strip() + "." for s in item.master_summary.split(". ") if s.strip()]
        bullets = "\n".join([f"• {s}" for s in summary_sentences])
        highlight_text = f"{item.title.upper()}\n\n{bullets}"

        ws.cell(row=row_num, column=1, value=idx).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=2, value=wk).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=3, value=item.date).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=4, value=initiator).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=5, value=category_val.upper()).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=6, value=kw).alignment = Alignment(horizontal="center", vertical="top")
        ws.cell(row=row_num, column=7, value=highlight_text).alignment = Alignment(wrap_text=True, vertical="top")
        ws.cell(row=row_num, column=8, value=item.link).alignment = Alignment(vertical="top")

    ws.column_dimensions["A"].width = 8
    ws.column_dimensions["B"].width = 6
    ws.column_dimensions["C"].width = 12
    ws.column_dimensions["D"].width = 22
    ws.column_dimensions["E"].width = 18
    ws.column_dimensions["F"].width = 20
    ws.column_dimensions["G"].width = 80
    ws.column_dimensions["H"].width = 40

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={safe_filename}"},
    )


# ==========================================
# --- EXPORT: WORD ---
# ==========================================
@app.post("/export-word")
async def export_word(request: ExportRequest):
    safe_filename = sanitize_filename(request.filename)
    doc = Document()
    section = doc.sections[0]
    section.page_width = DocxInches(8.5)
    section.page_height = DocxInches(11)
    section.left_margin = DocxInches(1)
    section.right_margin = DocxInches(1)
    section.top_margin = DocxInches(1)
    section.bottom_margin = DocxInches(1)
    section.different_first_page_header_footer = True

    p_brand = doc.add_paragraph()
    run_brand = p_brand.add_run("Samsung General")
    run_brand.font.name = "Calibri"
    run_brand.font.size = DocxPt(7)
    run_brand.font.color.rgb = DocxRGBColor(0, 160, 70)
    p_brand.paragraph_format.space_before = DocxPt(0)
    p_brand.paragraph_format.space_after = DocxPt(0)

    header = section.header
    header_para = header.paragraphs[0]
    header_para.clear()
    run_h_left = header_para.add_run("SENSE Intelligence Brief")
    run_h_left.font.name = "Calibri"
    run_h_left.font.size = DocxPt(9)
    run_h_left.font.color.rgb = DocxRGBColor(120, 120, 120)
    header_para.add_run("\t")
    run_h_right = header_para.add_run(datetime.datetime.now().strftime("%d %B %Y"))
    run_h_right.font.name = "Calibri"
    run_h_right.font.size = DocxPt(9)
    run_h_right.font.color.rgb = DocxRGBColor(120, 120, 120)

    pPr_h = header_para._p.get_or_add_pPr()
    tabs_h = OxmlElement("w:tabs")
    tab_h = OxmlElement("w:tab")
    tab_h.set(qn("w:val"), "right")
    tab_h.set(qn("w:pos"), "9360")
    tabs_h.append(tab_h)
    pPr_h.append(tabs_h)
    pBdr_h = OxmlElement("w:pBdr")
    bottom_h = OxmlElement("w:bottom")
    bottom_h.set(qn("w:val"), "single")
    bottom_h.set(qn("w:sz"), "4")
    bottom_h.set(qn("w:space"), "1")
    bottom_h.set(qn("w:color"), "CCCCCC")
    pBdr_h.append(bottom_h)
    pPr_h.append(pBdr_h)

    footer = section.footer
    footer_para = footer.paragraphs[0]
    footer_para.clear()
    run_conf = footer_para.add_run("CONFIDENTIAL - For Internal Use Only")
    run_conf.font.name = "Calibri"
    run_conf.font.size = DocxPt(8)
    run_conf.font.color.rgb = DocxRGBColor(150, 150, 150)
    footer_para.add_run("\t")
    run_pg = footer_para.add_run()
    fldChar1 = OxmlElement("w:fldChar")
    fldChar1.set(qn("w:fldCharType"), "begin")
    instrText = OxmlElement("w:instrText")
    instrText.text = "PAGE"
    fldChar2 = OxmlElement("w:fldChar")
    fldChar2.set(qn("w:fldCharType"), "end")
    run_pg._r.append(fldChar1)
    run_pg._r.append(instrText)
    run_pg._r.append(fldChar2)
    run_pg.font.name = "Calibri"
    run_pg.font.size = DocxPt(8)
    run_pg.font.color.rgb = DocxRGBColor(150, 150, 150)

    pPr_f = footer_para._p.get_or_add_pPr()
    tabs_f = OxmlElement("w:tabs")
    tab_f = OxmlElement("w:tab")
    tab_f.set(qn("w:val"), "right")
    tab_f.set(qn("w:pos"), "9360")
    tabs_f.append(tab_f)
    pPr_f.append(tabs_f)
    pBdr_f = OxmlElement("w:pBdr")
    top_f = OxmlElement("w:top")
    top_f.set(qn("w:val"), "single")
    top_f.set(qn("w:sz"), "4")
    top_f.set(qn("w:space"), "1")
    top_f.set(qn("w:color"), "CCCCCC")
    pBdr_f.append(top_f)
    pPr_f.append(pBdr_f)

    def add_section_label(doc, text):
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.font.name = "Calibri"
        run.font.size = DocxPt(11)
        run.font.bold = True
        run.font.color.rgb = DocxRGBColor(0, 51, 102)
        run.font.all_caps = True
        p.paragraph_format.space_before = DocxPt(12)
        p.paragraph_format.space_after = DocxPt(4)
        return p

    def add_rule(doc):
        p = doc.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        pBdr = OxmlElement("w:pBdr")
        bottom = OxmlElement("w:bottom")
        bottom.set(qn("w:val"), "single")
        bottom.set(qn("w:sz"), "4")
        bottom.set(qn("w:space"), "1")
        bottom.set(qn("w:color"), "CCCCCC")
        pBdr.append(bottom)
        pPr.append(pBdr)
        p.paragraph_format.space_before = DocxPt(0)
        p.paragraph_format.space_after = DocxPt(10)

    for item in request.items:
        doc.add_page_break()
        p_title = doc.add_paragraph()
        p_title.paragraph_format.space_after = DocxPt(4)
        kw = item.keywords_found[0] if item.keywords_found else ""
        if kw:
            parts = re.split(f"({re.escape(kw)})", item.title, flags=re.IGNORECASE)
            for part in parts:
                run = p_title.add_run(part)
                run.font.name = "Calibri"
                run.font.size = DocxPt(16)
                run.font.bold = True
                run.font.color.rgb = DocxRGBColor(0, 0, 0)
        else:
            run = p_title.add_run(item.title)
            run.font.name = "Calibri"
            run.font.size = DocxPt(16)
            run.font.bold = True
            run.font.color.rgb = DocxRGBColor(0, 0, 0)

        region_text = (getattr(item, "region", None) or "Global").upper()
        importance = getattr(item, "importance_score", 0)
        p_meta = doc.add_paragraph()
        run_meta = p_meta.add_run(f"{item.date}     |     Region: {region_text}     |     Importance: {importance}/10")
        run_meta.font.name = "Calibri"
        run_meta.font.size = DocxPt(10)
        run_meta.font.color.rgb = DocxRGBColor(130, 130, 130)
        p_meta.paragraph_format.space_after = DocxPt(6)

        add_rule(doc)

        img_stream = download_image_for_export(item.top_image, add_border=True)
        if img_stream:
            p_img = doc.add_paragraph()
            p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p_img.add_run().add_picture(img_stream, width=DocxInches(5.5))
            p_img.paragraph_format.space_after = DocxPt(10)

        add_section_label(doc, "Intelligence Summary")
        sentences = [s.strip() + "." for s in item.master_summary.split(". ") if s.strip()]
        if not sentences: sentences = [item.master_summary]
        for sentence in sentences:
            p_b = doc.add_paragraph(style="List Bullet")
            run_b = p_b.add_run(sentence)
            run_b.font.name = "Calibri"
            run_b.font.size = DocxPt(11)
            run_b.font.color.rgb = DocxRGBColor(30, 30, 30)

        ai_opinion = generate_opinion(item.master_summary)
        add_section_label(doc, "Analysis")
        p_insight = doc.add_paragraph()
        p_insight.paragraph_format.left_indent = DocxInches(0.3)
        p_insight.paragraph_format.space_after = DocxPt(4)
        run_i = p_insight.add_run(ai_opinion)
        run_i.font.name = "Calibri"
        run_i.font.size = DocxPt(11)
        run_i.font.italic = True
        run_i.font.color.rgb = DocxRGBColor(80, 80, 80)

        team = determine_target_team(item.title, item.master_summary)
        add_section_label(doc, "Routing")
        p_team = doc.add_paragraph()
        run_team = p_team.add_run(team)
        run_team.font.name = "Calibri"
        run_team.font.size = DocxPt(11)
        run_team.font.color.rgb = DocxRGBColor(30, 30, 30)

        add_section_label(doc, "Sources")
        sources = item.sources if item.sources else []
        if sources:
            for source in sources:
                p_src = doc.add_paragraph()
                run_src = p_src.add_run(source.name)
                run_src.font.name = "Calibri"
                run_src.font.size = DocxPt(10)
                run_src.font.color.rgb = DocxRGBColor(60, 60, 60)
                p_src.paragraph_format.space_after = DocxPt(2)
        else:
            p_src = doc.add_paragraph()
            run_src = p_src.add_run(getattr(item, "source", "Unknown"))
            run_src.font.name = "Calibri"
            run_src.font.size = DocxPt(10)
            run_src.font.color.rgb = DocxRGBColor(60, 60, 60)

    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename={safe_filename}"},
    )


# ==========================================
# --- WORKFLOW ENDPOINTS ---
# ==========================================
@app.get("/workflow")
def get_workflow(request: Request):
    profile = get_profile_for_request(request)
    store = resolve_workflow_identities(load_workflow_store(request))
    return {
        "selected": apply_learned_regions(store.get("selected", []), profile),
        "approved": apply_learned_regions(store.get("approved", []), profile),
        "profile": profile,
    }


@app.post("/workflow/select")
def select_news(request: Request, item: dict = Body(...)):
    profile = get_profile_for_request(request)
    viewer_ip = get_client_ip(request)
    viewer_key = get_viewer_key(viewer_ip)
    viewer_name = get_viewer_profile(viewer_ip).get("display_name") or get_team_owner_for_ip(viewer_ip)
    with workflow_lock:
        store = load_workflow_store(request)
        if any(i.get("title") == item.get("title") for i in store["selected"]):
            return {"status": "exists", "message": "Already selected", "profile": profile}
        item["profile"] = profile
        item["selected_by_id"] = viewer_key
        item["selected_by"] = viewer_name or str(item.get("selected_by", "")).strip() or "Unknown"
        store["selected"].append(item)
        save_workflow_store(store, request)
    return {"status": "success", "count": len(store["selected"]), "profile": profile}


@app.post("/workflow/import")
def import_archived_news(request: Request, payload: dict = Body(...)):
    """Import selected archive signals into the shared profile review queue."""

    profile = get_profile_for_request(request)
    client_ip = get_client_ip(request)
    viewer_key = get_viewer_key(client_ip)
    viewer_name = (
        get_viewer_profile(client_ip).get("display_name")
        or get_team_owner_for_ip(client_ip)
        or "Unknown"
    )
    fingerprint = str(payload.get("_tracking_fingerprint") or "unknown")
    items = payload.get("items", [])
    if not isinstance(items, list) or not items:
        raise HTTPException(status_code=400, detail="Select at least one archived signal.")
    if len(items) > 100:
        raise HTTPException(status_code=400, detail="Import no more than 100 signals at once.")

    imported = []
    existing = []
    now = datetime.datetime.now().isoformat(timespec="minutes")
    with workflow_lock:
        store = load_workflow_store(request)
        known = {
            _article_identity(item)
            for item in (store.get("selected", []) + store.get("approved", []))
        }
        for raw_item in items:
            if not isinstance(raw_item, dict):
                continue
            identity = _article_identity(raw_item)
            if not identity:
                continue
            if identity in known:
                existing.append(raw_item.get("title") or identity)
                continue
            item = dict(raw_item)
            item["profile"] = profile
            item["selected_by_id"] = viewer_key
            item["selected_by"] = viewer_name
            item["selected_at"] = now
            item["selection_source"] = "briefing_archive"
            store["selected"].append(item)
            imported.append(item)
            known.add(identity)
        if imported:
            save_workflow_store(store, request)

    if imported:
        representative = dict(imported[0])
        representative["item_count"] = len(imported)
        representative["screen"] = "briefing_archive"
        record_usage_activity(
            client_ip,
            profile,
            fingerprint,
            "archive_import",
            json.dumps(representative, ensure_ascii=False),
        )
        for item in imported[1:]:
            PERSONALIZATION_SERVICE.record_event(
                viewer_key,
                profile,
                "archive_import",
                item,
            )
    return {
        "status": "success",
        "imported": len(imported),
        "already_present": len(existing),
        "selected_count": len(store.get("selected", [])),
        "profile": profile,
    }


@app.post("/workflow/approve")
def approve_news(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    item_title = payload.get("title")
    key = payload.get("key")
    if key != DIRECTOR_KEY:
        return {"status": "error", "message": "Invalid Director Key", "profile": profile}
    approver_ip = get_client_ip(request)
    with workflow_lock:
        store = load_workflow_store(request)
        item_to_approve = next((i for i in store["selected"] if i["title"] == item_title), None)
        if not item_to_approve:
            return {"status": "error", "message": "Item not found", "profile": profile}
        store["selected"] = [i for i in store["selected"] if i["title"] != item_title]
        item_to_approve["profile"] = profile
        item_to_approve["approved_by_id"] = get_viewer_key(approver_ip)
        item_to_approve["approved_by"] = (
            get_viewer_profile(approver_ip).get("display_name")
            or get_team_owner_for_ip(approver_ip)
            or "Authorized user"
        )
        item_to_approve["approved_at"] = datetime.datetime.now().isoformat(timespec="minutes")
        store["approved"].append(item_to_approve)
        save_workflow_store(store, request)
    return {"status": "success", "message": "Approved", "profile": profile}


@app.post("/workflow/remove")
def remove_news(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    title = payload.get("title")
    list_type = payload.get("list_type")
    if list_type not in ["selected", "approved"]:
        return {"status": "error", "message": "Invalid list type", "profile": profile}
    with workflow_lock:
        store = load_workflow_store(request)
        store[list_type] = [i for i in store[list_type] if i["title"] != title]
        save_workflow_store(store, request)
    return {"status": "success", "profile": profile}


# ==========================================
# --- REGION CORRECTION / LEARNING ENDPOINT ---
# ==========================================
@app.post("/region/correct")
def correct_region(request: Request, payload: dict = Body(...)):
    profile = get_active_profile_name(request)
    title = str(payload.get("title", "")).strip()
    region = normalize_region_label(payload.get("region"))
    keywords = normalize_region_keywords(payload.get("keywords", []))
    reason = str(payload.get("reason", "")).strip()
    if not title:
        raise HTTPException(status_code=400, detail="Article title is required")
    if not region:
        raise HTTPException(status_code=400, detail="Region must be Local or Global")
    if not keywords:
        raise HTTPException(
            status_code=400,
            detail="Add at least one keyword so future scans can learn this correction",
        )

    other_region = "Global" if region == "Local" else "Local"
    with region_learning_lock:
        learned = load_region_learning(profile)
        learned[other_region] = [
            keyword for keyword in learned[other_region] if keyword not in keywords
        ]
        for keyword in keywords:
            if keyword not in learned[region]:
                learned[region].append(keyword)
        learned["corrections"].append({
            "title": title,
            "previous_region": normalize_region_label(payload.get("previous_region")),
            "region": region,
            "keywords": keywords,
            "reason": reason,
            "created_at": datetime.datetime.now().isoformat(timespec="seconds"),
        })
        learned["corrections"] = learned["corrections"][-500:]
        save_region_learning(learned, profile)

    return {
        "status": "success",
        "profile": profile,
        "region": region,
        "keywords": keywords,
        "message": f"Saved. Future scans will use {len(keywords)} learned keyword(s) for {region} signals.",
    }


# ==========================================
# --- SITES ENDPOINTS ---
# ==========================================
@app.get("/sites")
def get_sites(request: Request):
    profile = UNIFIED_PROFILE if UNIFIED_CORPUS_ENABLED else get_profile_for_request(request)
    sites_path = DEFAULT_SITES_FILE if UNIFIED_CORPUS_ENABLED else get_sites_file_for_profile(profile)
    if os.path.exists(sites_path):
        sites = load_sites(Path(sites_path))
        sites.sort(key=lambda x: x.get("name", "").lower())
        return sites
    return []


@app.post("/sites")
def add_site(site: dict, request: Request):
    profile = UNIFIED_PROFILE if UNIFIED_CORPUS_ENABLED else get_profile_for_request(request)
    sites_path = DEFAULT_SITES_FILE if UNIFIED_CORPUS_ENABLED else get_sites_file_for_profile(profile)
    with sites_lock:
        sites = load_sites(Path(sites_path)) if os.path.exists(sites_path) else []
        requested_verticals = site.get("verticals") or site.get("vertical") or ["technology"]
        if isinstance(requested_verticals, str):
            requested_verticals = [requested_verticals]
        legacy_profile = (
            BROADCAST_PROFILE
            if "broadcast" in {str(value).strip().lower() for value in requested_verticals}
            else DEFAULT_PROFILE
        )
        normalized = normalize_site(
            {
                **site,
                "verticals": requested_verticals,
                "legacy_profile": site.get("legacy_profile") or legacy_profile,
                "audiences": site.get("audiences") or ["all"],
            },
            legacy_profile,
        )
        entrypoint = canonical_source_url(normalized.get("rss_url") or normalized.get("url"))
        if any(
            canonical_source_url(existing.get("rss_url") or existing.get("url")) == entrypoint
            for existing in sites
        ):
            raise HTTPException(status_code=409, detail="This source entrypoint already exists.")
        sites.append(normalized)
        temp_file = f"{sites_path}.{secrets.token_hex(6)}.tmp"
        with open(temp_file, "w", encoding="utf-8") as f:
            json.dump(sites, f, indent=4)
            f.flush()
            os.fsync(f.fileno())
        os.replace(temp_file, sites_path)
    return {
        "status": "success",
        "profile": "unified" if UNIFIED_CORPUS_ENABLED else profile,
        "sites_file": sites_path,
        "source": normalized,
    }


def _write_unified_sites(sites_path: str, sites: list[dict]):
    temp_file = f"{sites_path}.{secrets.token_hex(6)}.tmp"
    try:
        with open(temp_file, "w", encoding="utf-8") as handle:
            json.dump(sites, handle, indent=4, ensure_ascii=False)
            handle.flush()
            os.fsync(handle.fileno())
        os.replace(temp_file, sites_path)
    finally:
        if os.path.exists(temp_file):
            os.unlink(temp_file)


def _source_record_index(sites: list[dict], source_id: str) -> int | None:
    needle = str(source_id or "").strip().casefold()
    for index, source in enumerate(sites):
        aliases = {
            str(source.get("id") or "").strip().casefold(),
            canonical_source_url(source.get("rss_url") or source.get("url")).casefold(),
        }
        if needle in aliases:
            return index
    return None


@app.put("/sites/{source_id}")
def update_site(source_id: str, request: Request, payload: dict = Body(...)):
    """Update/toggle one record in the authoritative unified catalog."""

    sites_path = DEFAULT_SITES_FILE if UNIFIED_CORPUS_ENABLED else get_sites_file_for_profile(
        get_profile_for_request(request)
    )
    with sites_lock:
        sites = load_sites(Path(sites_path)) if os.path.exists(sites_path) else []
        index = _source_record_index(sites, source_id)
        if index is None:
            raise HTTPException(status_code=404, detail="Source not found.")
        existing = sites[index]
        protected = {
            "id": existing.get("id"),
            "legacy_profile": existing.get("legacy_profile") or DEFAULT_PROFILE,
            "verticals": existing.get("verticals") or ["technology"],
            "audiences": existing.get("audiences") or ["all"],
        }
        updated = normalize_site({**existing, **payload, **protected}, protected["legacy_profile"])
        updated_entrypoint = canonical_source_url(updated.get("rss_url") or updated.get("url"))
        if any(
            position != index
            and canonical_source_url(source.get("rss_url") or source.get("url")) == updated_entrypoint
            for position, source in enumerate(sites)
        ):
            raise HTTPException(status_code=409, detail="This source entrypoint already exists.")
        sites[index] = updated
        _write_unified_sites(sites_path, sites)
    return {"status": "success", "profile": public_profile_name(UNIFIED_PROFILE), "source": updated}


@app.delete("/sites/{source_id}")
def delete_site(source_id: str, request: Request):
    """Remove one explicitly addressed source from the unified catalog."""

    sites_path = DEFAULT_SITES_FILE if UNIFIED_CORPUS_ENABLED else get_sites_file_for_profile(
        get_profile_for_request(request)
    )
    with sites_lock:
        sites = load_sites(Path(sites_path)) if os.path.exists(sites_path) else []
        index = _source_record_index(sites, source_id)
        if index is None:
            raise HTTPException(status_code=404, detail="Source not found.")
        removed = sites.pop(index)
        _write_unified_sites(sites_path, sites)
    return {
        "status": "success",
        "profile": public_profile_name(UNIFIED_PROFILE),
        "removed": removed,
        "count": len(sites),
    }


# ==========================================
# --- STATUS ENDPOINT ---
# ==========================================
@app.get("/status")
def get_system_status(request: Request):
    active_profile = get_active_profile_name(request)
    details_allowed = get_client_ip(request) in SYSTEM_STATUS_ALLOWED_IPS
    cached_capabilities = pipeline_health_cache.get("result")
    with scheduler_lock:
        running_jobs = sum(
            1
            for job in active_jobs.values()
            if job.get("status") in {"queued", "running"}
        )
        response = {
            **SCHEDULER_STATUS,
            "active_profile": public_profile_name(active_profile),
            "active_manual_jobs": running_jobs,
            "capacity_remaining": crawl_semaphore._value,
            "details_allowed": details_allowed,
            "pipeline": {
                "mode": (
                    cached_capabilities.get("mode", "awaiting_preflight")
                    if SAMSUNG_PIPELINE_ENABLED and cached_capabilities
                    else "samsung_preflight_pending"
                    if SAMSUNG_PIPELINE_ENABLED
                    else "local_models"
                ),
                "discovery_only": (
                    cached_capabilities.get("discovery_only", False)
                    if cached_capabilities
                    else False
                ),
                "web_search_enabled": WEB_SEARCH_ENRICHMENT_ENABLED,
                "chat_summary_enabled": FINAL_CHAT_SUMMARY_ENABLED,
                "automatic_local_fallback": True,
                "credentials_ready": SAMSUNG_PIPELINE_CREDENTIALS_READY,
            },
        }
        if details_allowed:
            response["pipeline"]["last_preflight"] = cached_capabilities
            response["profiles"] = (
                {"unified": get_profile_debug_info(UNIFIED_PROFILE)}
                if UNIFIED_CORPUS_ENABLED
                else {
                    DEFAULT_PROFILE: get_profile_debug_info(DEFAULT_PROFILE),
                    BROADCAST_PROFILE: get_profile_debug_info(BROADCAST_PROFILE),
                }
            )
        return response


# ==========================================
# DOSSIER INSIGHT ENDPOINT
# ==========================================
@app.post("/insight")
def get_dossier_insight(request: Request, item: dict = Body(...)):
    profile = get_active_profile_name(request)
    existing = str(
        item.get("why_it_matters")
        or item.get("why_matters")
        or ""
    ).strip()
    if existing:
        return {
            "status": "success",
            "profile": profile,
            "why_matters": existing,
            "generated_by": item.get("summarized_by") or "stored",
        }
    insight, generated_by = generate_why_it_matters(item, profile)
    return {
        "status": "success",
        "profile": profile,
        "why_matters": insight,
        "generated_by": generated_by,
    }


# ==========================================
# --- BRIEFING META ENDPOINT ---
# ==========================================
@app.get("/briefing/meta")
def get_briefing_meta(request: Request):
    profile = get_profile_for_request(request)
    latest = get_latest_briefing_file_for_profile(profile)
    if not latest:
        return {"last_updated": None, "count": 0, "filename": None, "profile": public_profile_name(profile)}
    try:
        with open(latest, "r", encoding="utf-8") as f:
            data = json.load(f)
        return {
            "last_updated": datetime.datetime.fromtimestamp(os.path.getmtime(latest)).isoformat(),
            "filename": os.path.basename(latest),
            "count": len(data),
            "profile": public_profile_name(profile),
        }
    except Exception:
        return {"last_updated": None, "count": 0, "filename": None, "profile": public_profile_name(profile)}


# ==========================================
# --- LATEST BRIEFING ENDPOINT ---
# ==========================================
@app.get("/latest-briefing")
def get_latest_briefing(request: Request):
    profile = get_profile_for_request(request)
    client_ip = get_client_ip(request)
    viewer_key = get_viewer_key(client_ip)
    latest = get_latest_briefing_file_for_profile(profile)
    if latest:
        try:
            with open(latest, "r", encoding="utf-8") as file_obj:
                data = json.load(file_obj)
                if data and len(data) > 0:
                    visible_data = apply_learned_regions(
                        filter_viewer_hidden(data, request, profile),
                        profile,
                    )
                    ranked_data, personalization = PERSONALIZATION_SERVICE.rank_articles(
                        visible_data,
                        viewer_key,
                        profile,
                        get_viewer_saved_items(request, profile),
                    )
                    viewer = get_viewer_profile(client_ip)
                    personalization["viewer_name"] = (
                        viewer.get("display_name")
                        or get_team_owner_for_ip(client_ip)
                        or ""
                    )
                    return {
                        "status": "success",
                        "result": ranked_data,
                        "type": "scheduler",
                        "source": "shared",
                        "profile": public_profile_name(profile),
                        "filename": os.path.basename(latest),
                        "generated_at": datetime.datetime.fromtimestamp(os.path.getmtime(latest)).strftime("%d %b %Y, %I:%M %p"),
                        "personalization": personalization,
                    }
        except Exception as e:
            return {"status": "error", "result": [], "profile": public_profile_name(profile), "message": str(e)}
    return {"status": "empty", "result": [], "profile": public_profile_name(profile)}


# ==========================================
# --- BRIEFING REMOVE/RESTORE (Disk) ---
# ==========================================
def _load_briefing_items(file_path: str) -> list:
    """Read one briefing as a list or fail before any paired state mutates."""

    try:
        with open(file_path, "r", encoding="utf-8") as file_obj:
            items = json.load(file_obj)
    except (OSError, ValueError, TypeError) as error:
        raise RuntimeError(f"Could not read the shared briefing: {error}") from error
    if not isinstance(items, list):
        raise RuntimeError("The shared briefing has an invalid JSON shape.")
    return items


def _save_briefing_items(file_path: str, items: list):
    """Replace a briefing atomically so a failed write preserves the old file."""

    destination = os.path.abspath(file_path)
    os.makedirs(os.path.dirname(destination), exist_ok=True)
    temp_file = f"{destination}.{secrets.token_hex(6)}.tmp"
    try:
        with open(temp_file, "w", encoding="utf-8") as file_obj:
            json.dump(items, file_obj, indent=4, ensure_ascii=False)
            file_obj.flush()
            os.fsync(file_obj.fileno())
        os.replace(temp_file, destination)
    finally:
        if os.path.exists(temp_file):
            try:
                os.unlink(temp_file)
            except OSError:
                pass


def _paired_state_failure(operation: str, error: Exception, rollback_error: Exception = None):
    """Raise a machine-readable failure for a paired briefing/store mutation."""

    if rollback_error is not None:
        raise HTTPException(
            status_code=500,
            detail={
                "status": "partial",
                "operation": operation,
                "state": "recovery_required",
                "recoverable": True,
                "message": (
                    f"The {operation} operation was interrupted and automatic rollback "
                    "also failed. Refresh the briefing and retry; an administrator can "
                    "reconcile the Not Interested store if needed."
                ),
                "cause": str(error),
                "rollback_error": str(rollback_error),
            },
        )
    raise HTTPException(
        status_code=503,
        detail={
            "status": "error",
            "operation": operation,
            "state": "rolled_back",
            "recoverable": True,
            "message": (
                f"The {operation} operation could not be completed. No paired state "
                "change was kept; retry safely."
            ),
            "cause": str(error),
        },
    )


@app.post("/briefing/remove")
def remove_from_briefing(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    title = str(payload.get("title", "")).strip()
    if not title:
        raise HTTPException(status_code=400, detail="No title provided")

    latest = get_latest_briefing_file_for_profile(profile)
    if not latest:
        raise HTTPException(status_code=404, detail="No briefing file found")
    try:
        with briefing_lock:
            data = _load_briefing_items(latest)
            original_count = len(data)
            normalized_title = title.casefold()
            data = [
                item for item in data
                if str(item.get("title", "")).strip().casefold() != normalized_title
            ]
            removed = original_count - len(data)
            if removed > 0:
                _save_briefing_items(latest, data)
                print(f"[BRIEFING] Removed '{title[:50]}' from {os.path.basename(latest)}")
        return {"status": "success", "removed": removed, "remaining": len(data), "profile": profile}
    except Exception as e:
        if isinstance(e, HTTPException):
            raise
        raise HTTPException(status_code=500, detail=f"Could not remove article from briefing: {e}") from e


@app.post("/briefing/restore")
def restore_to_briefing(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    article = payload.get("article")
    if not isinstance(article, dict) or not str(article.get("title", "")).strip():
        raise HTTPException(status_code=400, detail="A titled article is required")

    latest = get_latest_briefing_file_for_profile(profile)
    if not latest:
        raise HTTPException(status_code=404, detail="No briefing file found")
    try:
        with briefing_lock:
            data = _load_briefing_items(latest)
            normalized_title = str(article.get("title", "")).strip().casefold()
            restored = not any(
                str(item.get("title", "")).strip().casefold() == normalized_title
                for item in data
            )
            if restored:
                data.insert(0, article)
                _save_briefing_items(latest, data)
                print(f"[BRIEFING] Restored '{article.get('title', '')[:50]}' to {os.path.basename(latest)}")
        return {"status": "success", "restored": restored, "count": len(data), "profile": profile}
    except Exception as e:
        if isinstance(e, HTTPException):
            raise
        raise HTTPException(status_code=500, detail=f"Could not restore article to briefing: {e}") from e


# ==========================================
# --- READ-ONLY EXTRACTED INTELLIGENCE SEARCH ---
# ==========================================
def _archive_file_date(file_path: str):
    """Return the briefing date encoded in an extracted archive filename."""

    match = re.search(r"briefing_(\d{4}-\d{2}-\d{2})_", os.path.basename(file_path))
    if not match:
        return None
    try:
        return datetime.datetime.strptime(match.group(1), "%Y-%m-%d").date()
    except ValueError:
        return None


def _article_archive_date(article: dict, fallback=None):
    """Read an article date without making any network or crawler request."""

    for key in ("date", "published_at", "publishedAt", "first_seen", "generated_at"):
        value = str(article.get(key) or "").strip()
        if not value:
            continue
        iso_candidate = value[:10]
        try:
            return datetime.datetime.strptime(iso_candidate, "%Y-%m-%d").date()
        except ValueError:
            continue
    return fallback


def _article_source_names(article: dict):
    names = []
    primary = str(article.get("source") or article.get("src") or "").strip()
    if primary:
        names.append(primary)
    for source in article.get("sources") or article.get("source_list") or []:
        if isinstance(source, dict):
            name = str(source.get("name") or source.get("title") or source.get("source") or "").strip()
        else:
            name = str(source).strip()
        if name:
            names.append(name)
    return list(dict.fromkeys(names))


@app.get("/internal-content/samsung-feed")
def get_samsung_internal_feed(request: Request, limit: int = Query(100, ge=1, le=100)):
    """Return Samsung channels projected from the one retained corpus.

    This endpoint is read-only.  It never starts Scrapy, enrichment, or a
    scheduler run: it scans the already-generated briefing history and applies
    the explicit source contracts used by Samsung Internal.
    """

    from news_scrapper.samsung_internal_feed import build_samsung_internal_feed

    profile = get_profile_for_request(request)
    archive_files = [
        file_path
        for file_path in get_profile_history_files(profile)
        if os.path.basename(file_path).startswith("briefing_")
    ]
    payload = build_samsung_internal_feed(archive_files, limit=limit)
    for channel in ("global", "local", "sampark"):
        visible = filter_viewer_hidden(payload[channel], request, profile)
        payload[channel] = apply_learned_regions(visible, profile)
        payload["counts"][channel] = len(payload[channel])
    payload["profile"] = public_profile_name(profile)
    return payload


def _archive_search_terms(query: str):
    """Produce useful OR-search terms from natural text or comma-separated input."""

    raw = str(query or "").strip().casefold()
    if not raw:
        return []
    phrases = [part.strip() for part in raw.split(",") if part.strip()]
    words = re.findall(r"[\w.+#-]{2,}", raw, flags=re.UNICODE)
    return list(dict.fromkeys(phrases + words))


def search_extracted_intelligence(
    profile: str,
    query: str,
    from_date: str = None,
    to_date: str = None,
    target_sites: str = None,
    limit: int = 250,
):
    """Search scheduler-produced JSON archives only.

    This function deliberately has no Scrapy, subprocess, HTTP, enrichment, or
    scheduler integration. It is the data boundary used by the UI's Scan
    screen so an interactive search can never launch an internet crawl.
    """

    terms = _archive_search_terms(query)
    if not terms:
        return {
            "status": "error",
            "message": "Enter at least one search term.",
            "results": [],
            "count": 0,
            "profile": profile,
            "search_scope": "extracted_archives_only",
            "crawler_started": False,
        }

    try:
        start = datetime.datetime.strptime(from_date, "%Y-%m-%d").date() if from_date else None
        end = datetime.datetime.strptime(to_date, "%Y-%m-%d").date() if to_date else None
    except ValueError:
        return {
            "status": "error",
            "message": "Dates must use YYYY-MM-DD format.",
            "results": [],
            "count": 0,
            "profile": profile,
            "search_scope": "extracted_archives_only",
            "crawler_started": False,
        }
    if start and end and start > end:
        start, end = end, start

    requested_sources = {
        source.strip().casefold()
        for source in str(target_sites or "").split(",")
        if source.strip() and source.strip().casefold() != "all"
    }

    archive_files = [
        file_path
        for file_path in get_profile_history_files(profile)
        if os.path.basename(file_path).startswith("briefing_")
    ]
    archive_files.sort(key=os.path.getmtime, reverse=True)

    matches = {}
    scanned_articles = 0
    searchable_files = 0
    for file_path in archive_files:
        file_date = _archive_file_date(file_path)
        try:
            with open(file_path, "r", encoding="utf-8") as file_obj:
                payload = json.load(file_obj)
        except (OSError, ValueError, TypeError):
            continue
        if not isinstance(payload, list):
            continue
        searchable_files += 1
        for raw_article in payload:
            if not isinstance(raw_article, dict):
                continue
            scanned_articles += 1
            article_date = _article_archive_date(raw_article, file_date)
            if start and (article_date is None or article_date < start):
                continue
            if end and (article_date is None or article_date > end):
                continue

            source_names = _article_source_names(raw_article)
            normalized_sources = {name.casefold() for name in source_names}
            if requested_sources and not requested_sources.intersection(normalized_sources):
                continue

            title = str(raw_article.get("title") or "")
            keywords = raw_article.get("keywords_found") or raw_article.get("keywords") or []
            if not isinstance(keywords, list):
                keywords = [keywords]
            summary = " ".join(
                str(raw_article.get(key) or "")
                for key in ("master_summary", "summary", "ppt_summary", "snippet")
            )
            full_text = " ".join(
                [
                    title,
                    summary,
                    str(raw_article.get("full_contents") or raw_article.get("full_content") or ""),
                    " ".join(str(value) for value in keywords),
                    " ".join(source_names),
                    str(raw_article.get("category") or ""),
                    str(raw_article.get("region") or ""),
                ]
            ).casefold()
            title_text = title.casefold()
            keyword_text = " ".join(str(value) for value in keywords).casefold()
            summary_text = summary.casefold()
            matched_terms = [term for term in terms if term in full_text]
            if not matched_terms:
                continue

            score = 0
            whole_query = str(query or "").strip().casefold()
            if whole_query and whole_query in title_text:
                score += 18
            for term in matched_terms:
                if term in title_text:
                    score += 9
                if term in keyword_text:
                    score += 7
                if term in summary_text:
                    score += 3
                else:
                    score += 1
            score += min(int(raw_article.get("source_count") or len(source_names) or 1), 5)
            score += min(int(raw_article.get("importance_score") or 0) // 20, 5)

            article = apply_learned_region(dict(raw_article), profile)
            article["search_score"] = score
            article["matched_terms"] = matched_terms
            article["archive_file"] = os.path.basename(file_path)
            article["archive_date"] = article_date.isoformat() if article_date else None
            article["search_scope"] = "extracted_archives_only"
            stable_key = str(article.get("link") or article.get("url") or title).strip().casefold()
            existing = matches.get(stable_key)
            if existing is None or score > existing.get("search_score", 0):
                matches[stable_key] = article

    results = list(matches.values())
    results.sort(
        key=lambda article: (
            int(article.get("search_score") or 0),
            str(article.get("archive_date") or article.get("date") or ""),
            int(article.get("importance_score") or 0),
        ),
        reverse=True,
    )
    results = results[:limit]
    return {
        "status": "success",
        "results": results,
        "count": len(results),
        "profile": profile,
        "query": query,
        "from_date": start.isoformat() if start else None,
        "to_date": end.isoformat() if end else None,
        "archive_files_searched": searchable_files,
        "articles_searched": scanned_articles,
        "search_scope": "extracted_archives_only",
        "crawler_started": False,
    }


@app.get("/archive/search")
def search_archive(
    request: Request,
    query: str = Query(..., min_length=1),
    from_date: str = Query(None),
    to_date: str = Query(None),
    target_sites: str = Query(None),
    limit: int = Query(250, ge=1, le=500),
):
    profile = get_profile_for_request(request)
    result = search_extracted_intelligence(
        profile=profile,
        query=query,
        from_date=from_date,
        to_date=to_date,
        target_sites=target_sites,
        limit=limit,
    )
    visible_results = filter_viewer_hidden(result.get("results", []), request, profile)
    result["results"] = visible_results
    result["count"] = len(visible_results)
    return result


# ==========================================
# --- HISTORY ENDPOINTS ---
# ==========================================
@app.get("/history/list")
def get_history_list(request: Request, session_id: str = Query(None)):
    profile = get_active_profile_name(request)
    files = get_profile_history_files(profile)
    files.sort(key=os.path.getmtime, reverse=True)
    file_list = []
    for f in files:
        filename = os.path.basename(f)
        try:
            if filename.startswith("briefing_"):
                ts = filename.replace("briefing_", "").replace(".json", "")
                display_date = datetime.datetime.strptime(ts, "%Y-%m-%d_%H-%M-%S").strftime("%b %d, %Y - %I:%M %p")
                file_list.append({"filename": filename, "display": display_date, "type": "scheduler", "profile": profile})
            elif filename.startswith("manual_"):
                if not session_id: continue
                expected_prefix = f"manual_{session_id}_"
                if not filename.startswith(expected_prefix): continue
                ts = filename.replace(expected_prefix, "").replace(".json", "")
                display_date = datetime.datetime.strptime(ts, "%Y-%m-%d_%H-%M-%S").strftime("%b %d, %Y - %I:%M %p")
                file_list.append({"filename": filename, "display": display_date, "type": "manual", "profile": profile})
        except:
            pass
    return file_list


@app.get("/history/range")
def get_history_by_range(request: Request, from_date: str, to_date: str, session_id: str = Query(None)):
    profile = get_active_profile_name(request)
    try:
        start_date = datetime.datetime.strptime(from_date, "%Y-%m-%d").date()
        end_date = datetime.datetime.strptime(to_date, "%Y-%m-%d").date()
    except ValueError:
        return {"status": "error", "message": "Invalid date format.", "profile": profile}

    merged_results = []
    seen_titles = set()
    files = get_profile_history_files(profile)
    for f in files:
        filename = os.path.basename(f)
        try:
            if filename.startswith("briefing_"):
                ts = filename.replace("briefing_", "").replace(".json", "")
                date_part = ts.split("_")[0]
            elif filename.startswith("manual_"):
                if not session_id: continue
                expected_prefix = f"manual_{session_id}_"
                if not filename.startswith(expected_prefix): continue
                ts = filename.replace(expected_prefix, "").replace(".json", "")
                date_part = ts.split("_")[0]
            else:
                continue

            file_date = datetime.datetime.strptime(date_part, "%Y-%m-%d").date()
            if start_date <= file_date <= end_date:
                with open(f, "r", encoding="utf-8") as file_obj:
                    data = json.load(file_obj)
                    for item in data:
                        item = apply_learned_region(item, profile)
                        title = item.get("title", "")
                        if title not in seen_titles:
                            seen_titles.add(title)
                            merged_results.append(item)
        except:
            continue

    merged_results.sort(key=lambda x: x.get("date", ""), reverse=True)
    return {"status": "success", "count": len(merged_results), "results": merged_results, "profile": profile}


@app.get("/history/{filename}")
def get_history_file(request: Request, filename: str):
    profile = get_active_profile_name(request)
    file_path = resolve_profile_history_file(filename, profile)
    if not file_path:
        raise HTTPException(status_code=404, detail="File not found")
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return {"status": "success", "results": apply_learned_regions(data, profile), "profile": profile}
    except Exception as e:
        return {"status": "error", "message": str(e)}


# ==========================================
# VOC FEEDBACK STORE
# ==========================================
VOC_FEEDBACK_FILE = os.path.join(ROOT_DIR, "voc_feedback.json")


@app.post("/voc")
async def submit_voc_feedback(request: Request):
    import uuid

    payload = await request.json()
    message = str(payload.get("message", "")).strip()
    if not message:
        return {"status": "error", "message": "Feedback message is required"}
    profile = get_active_profile_name(request)
    feedback_item = {
        "id": str(uuid.uuid4()),
        "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
        "name": str(payload.get("name", "anonymous")).strip() or "anonymous",
        "type": str(payload.get("type", "ui_feedback")).strip(),
        "message": message,
        "page": str(payload.get("page", "unknown")).strip(),
        "profile": profile,
        "viewer_ip_hash": get_viewer_key(get_client_ip(request)),
        "viewer_name": get_viewer_profile(get_client_ip(request)).get("display_name", "anonymous"),
    }
    with voc_lock:
        items = []
        if os.path.exists(VOC_FEEDBACK_FILE):
            try:
                with open(VOC_FEEDBACK_FILE, "r", encoding="utf-8") as f:
                    items = json.load(f)
            except Exception:
                items = []
        items.insert(0, feedback_item)
        temp_file = f"{VOC_FEEDBACK_FILE}.{secrets.token_hex(6)}.tmp"
        with open(temp_file, "w", encoding="utf-8") as f:
            json.dump(items[:500], f, indent=2, ensure_ascii=False)
            f.flush()
            os.fsync(f.fileno())
        os.replace(temp_file, VOC_FEEDBACK_FILE)
    return {"status": "success", "message": "Feedback saved", "item": feedback_item}


# ==========================================
# GATEKEEPER REVIEW
# ==========================================
GATEKEEPER_QUEUE_FILE = os.path.join(ROOT_DIR, "gatekeeper_restore_queue.json")


def require_gatekeeper_ip_access(request: Request):
    """Authorize destructive shared-feed actions by resolved client IP only."""
    ip = get_client_ip(request)
    if ip not in GATEKEEPER_ALLOWED_IPS:
        raise HTTPException(status_code=403, detail="Global article removal is not enabled for this network.")
    return ip


def require_gatekeeper_access(request: Request):
    ip = require_gatekeeper_ip_access(request)
    provided = request.headers.get("x-gatekeeper-key", "")
    if not provided or not secrets.compare_digest(str(provided), str(GATEKEEPER_KEY)):
        raise HTTPException(status_code=403, detail="Invalid Gatekeeper key.")
    return ip


def load_dropped_articles():
    path = os.path.join(ROOT_DIR, "dropped_articles.json")
    try:
        with open(path, "r", encoding="utf-8") as file_obj:
            data = json.load(file_obj)
    except (OSError, ValueError):
        data = []
    output = []
    for index, raw in enumerate(data if isinstance(data, list) else []):
        item = dict(raw)
        if not item.get("id"):
            stable = f"{item.get('profile')}|{item.get('timestamp')}|{item.get('title')}|{index}"
            item["id"] = hashlib.sha256(stable.encode("utf-8")).hexdigest()[:24]
            item["legacy"] = True
        item.setdefault("status", "dropped")
        item.setdefault("updated_at", item.get("timestamp", ""))
        item.setdefault("keywords_found", item.get("keyword", []))
        item.setdefault("summary", item.get("master_summary", ""))
        item.setdefault("restore_eligible", bool(item.get("link")) and not item.get("legacy", False))
        output.append(item)
    return output


def save_dropped_articles(items):
    path = os.path.join(ROOT_DIR, "dropped_articles.json")
    temp = f"{path}.tmp"
    with open(temp, "w", encoding="utf-8") as file_obj:
        json.dump(items[-500:], file_obj, indent=2, ensure_ascii=False)
    os.replace(temp, path)


def load_gatekeeper_jobs():
    try:
        with open(GATEKEEPER_QUEUE_FILE, "r", encoding="utf-8") as file_obj:
            data = json.load(file_obj)
            return data if isinstance(data, list) else []
    except (OSError, ValueError):
        return []


def save_gatekeeper_jobs(jobs):
    temp = f"{GATEKEEPER_QUEUE_FILE}.tmp"
    with open(temp, "w", encoding="utf-8") as file_obj:
        json.dump(jobs[-200:], file_obj, indent=2, ensure_ascii=False)
    os.replace(temp, GATEKEEPER_QUEUE_FILE)


@app.get("/gatekeeper/access")
def gatekeeper_access(request: Request):
    ip = get_client_ip(request)
    profile = get_active_profile_name(request)
    return {
        "allowed": ip in GATEKEEPER_ALLOWED_IPS,
        "ip": ip,
        "owner": get_viewer_profile(ip).get("display_name") or get_team_owner_for_ip(ip) or "Authorized user",
        "active_profile": profile,
    }


@app.get("/gatekeeper/dropped")
def gatekeeper_dropped(
    request: Request,
    profile: str = Query("all"),
    status: str = Query("all"),
    search: str = Query(""),
    offset: int = Query(0, ge=0),
    limit: int = Query(100, ge=1, le=500),
):
    require_gatekeeper_access(request)
    items = load_dropped_articles()
    if profile in PROFILE_CONFIGS:
        items = [item for item in items if item.get("profile", DEFAULT_PROFILE) == profile]
    scoped_items = list(items)
    counts = {"all": len(scoped_items)}
    for state in ("dropped", "queued", "processing", "restored", "failed"):
        counts[state] = sum(1 for item in scoped_items if item.get("status", "dropped") == state)
    counts["eligible"] = sum(1 for item in scoped_items if item.get("restore_eligible"))
    if status != "all":
        items = [item for item in items if item.get("status", "dropped") == status]
    query = search.strip().casefold()
    if query:
        items = [
            item for item in items
            if query in " ".join(
                str(item.get(key, ""))
                for key in ("title", "summary", "source", "keywords_found", "bouncer_reason")
            ).casefold()
        ]
    items.sort(key=lambda item: item.get("updated_at") or item.get("timestamp") or "", reverse=True)
    return {
        "status": "success",
        "items": items[offset : offset + limit],
        "counts": counts,
        "has_more": offset + limit < len(items),
    }


@app.get("/gatekeeper/queue")
def gatekeeper_queue(request: Request, profile: str = Query("all")):
    require_gatekeeper_access(request)
    jobs = load_gatekeeper_jobs()
    if profile in PROFILE_CONFIGS:
        jobs = [job for job in jobs if job.get("profile") == profile]
    counts = {
        "queued": sum(1 for job in jobs if job.get("status") == "queued"),
        "processing": sum(1 for job in jobs if job.get("status") == "processing"),
        "completed": sum(1 for job in jobs if job.get("status") == "completed"),
        "failed": sum(1 for job in jobs if job.get("status") == "failed"),
    }
    return {
        "status": "success",
        "jobs": list(reversed(jobs[-50:])),
        "counts": counts,
        "worker": {"running": counts["processing"] > 0, "current_article": next((job.get("title") for job in jobs if job.get("status") == "processing"), None)},
    }


def update_gatekeeper_job(job_id: str, **changes):
    with gatekeeper_queue_lock:
        jobs = load_gatekeeper_jobs()
        job = next((entry for entry in jobs if entry.get("id") == job_id), None)
        if job:
            job.update(changes)
            save_gatekeeper_jobs(jobs)
        return job


def restore_gatekeeper_article(article_id: str, job_id: str):
    """Process one queued restore on the single background worker."""

    with dropped_lock:
        items = load_dropped_articles()
        item = next((entry for entry in items if entry.get("id") == article_id), None)
        if not item:
            raise HTTPException(status_code=404, detail="Dropped article not found.")
        if not item.get("restore_eligible"):
            raise HTTPException(status_code=409, detail="This legacy record lacks a restorable source URL.")
        profile = item.get("profile", DEFAULT_PROFILE)
        now = datetime.datetime.now().isoformat(timespec="seconds")
        item.update({"status": "processing", "updated_at": now, "restore_error": None})
        save_dropped_articles(items)
    update_gatekeeper_job(job_id, status="processing", updated_at=now)
    try:
        restored = enrich_raw_articles([item], get_profile_config(profile)["keywords"], profile)[0]
        restored = enrich_final_articles([restored], profile)[0]
        restored["profile"] = profile
        restored.pop("status", None)
        latest_file = get_latest_briefing_file_for_profile(profile)
        briefing = []
        if latest_file:
            try:
                with open(latest_file, "r", encoding="utf-8") as file_obj:
                    briefing = json.load(file_obj)
            except (OSError, ValueError):
                briefing = []
        identity = restored.get("link") or restored.get("title", "").casefold()
        if not any((entry.get("link") or entry.get("title", "").casefold()) == identity for entry in briefing):
            briefing.insert(0, restored)
        if not latest_file:
            timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
            latest_file = os.path.join(get_profile_history_dir(profile), f"briefing_{timestamp}.json")
        with open(latest_file, "w", encoding="utf-8") as file_obj:
            json.dump(briefing, file_obj, indent=2, ensure_ascii=False)
        with dropped_lock:
            items = load_dropped_articles()
            item = next(entry for entry in items if entry.get("id") == article_id)
            item.update({
                "status": "restored", "updated_at": datetime.datetime.now().isoformat(timespec="seconds"),
                "briefing_filename": os.path.basename(latest_file),
            })
            save_dropped_articles(items)
        update_gatekeeper_job(job_id, status="completed", updated_at=item["updated_at"], error=None)
        return item
    except Exception as error:
        with dropped_lock:
            items = load_dropped_articles()
            failed = next(entry for entry in items if entry.get("id") == article_id)
            failed.update({"status": "failed", "restore_error": str(error)[:500], "updated_at": datetime.datetime.now().isoformat(timespec="seconds")})
            save_dropped_articles(items)
        update_gatekeeper_job(job_id, status="failed", error=str(error)[:500], updated_at=failed["updated_at"])
        print(f"[GATEKEEPER:{profile}] Restore failed for {item.get('title', '')[:70]}: {error}", flush=True)


def queue_gatekeeper_article(article_id: str, retry: bool = False):
    with dropped_lock:
        items = load_dropped_articles()
        item = next((entry for entry in items if entry.get("id") == article_id), None)
        if not item:
            raise HTTPException(status_code=404, detail="Dropped article not found.")
        if not item.get("restore_eligible"):
            raise HTTPException(status_code=409, detail="This legacy record lacks a restorable source URL.")
        if item.get("status") in {"queued", "processing"}:
            return {"status": "exists", "item": item}
        if retry and item.get("status") != "failed":
            raise HTTPException(status_code=409, detail="Only failed restorations can be retried.")
        now = datetime.datetime.now().isoformat(timespec="seconds")
        item.update({"status": "queued", "updated_at": now, "restore_error": None})
        save_dropped_articles(items)

    job = {
        "id": secrets.token_hex(10),
        "article_id": article_id,
        "title": item.get("title", ""),
        "profile": item.get("profile", DEFAULT_PROFILE),
        "status": "queued",
        "attempts": 1,
        "created_at": now,
        "updated_at": now,
    }
    with gatekeeper_queue_lock:
        jobs = load_gatekeeper_jobs()
        job["attempts"] = 1 + sum(1 for entry in jobs if entry.get("article_id") == article_id)
        jobs.append(job)
        save_gatekeeper_jobs(jobs)
    gatekeeper_executor.submit(restore_gatekeeper_article, article_id, job["id"])
    return {"status": "queued", "item": item, "job": job}


@app.post("/gatekeeper/restore")
def gatekeeper_restore(request: Request, payload: dict = Body(...)):
    require_gatekeeper_access(request)
    article_id = str(payload.get("id", "")).strip()
    if not article_id:
        raise HTTPException(status_code=400, detail="Dropped article id is required.")
    return queue_gatekeeper_article(article_id)


@app.post("/gatekeeper/retry")
def gatekeeper_retry(request: Request, payload: dict = Body(...)):
    require_gatekeeper_access(request)
    article_id = str(payload.get("id", "")).strip()
    if not article_id:
        raise HTTPException(status_code=400, detail="Dropped article id is required.")
    return queue_gatekeeper_article(article_id, retry=True)


# ==========================================
# --- TRAINING ENDPOINTS ---
# ==========================================
@app.post("/train")
def save_training_data(request: Request, data: VotePayload, background_tasks: BackgroundTasks):
    profile = get_active_profile_name(request)
    total = save_training_vote(data.keywords, data.summary, data.vote, data.title or "", profile)
    background_tasks.add_task(enqueue_bouncer_retrain, profile)
    return {"status": "success", "total_samples": total, "profile": profile, "retrain_scheduled": True}


# ==========================================
# --- NOT INTERESTED ENDPOINTS ---
# ==========================================
@app.post("/not-interested")
def add_not_interested(request: Request, background_tasks: BackgroundTasks, payload: dict = Body(...)):
    authorized_ip = require_gatekeeper_ip_access(request)
    profile = get_active_profile_name(request)
    title = str(payload.get("title", "")).strip()
    if not title:
        raise HTTPException(status_code=400, detail="An article title is required.")
    summary = (
        payload.get("master_summary")
        or payload.get("summary")
        or payload.get("snippet")
        or payload.get("full_content")
        or payload.get("full_contents")
        or ""
    )
    keywords = payload.get("keywords_found", [])
    entry = {
        "title": title,
        "master_summary": summary,
        "ppt_summary": payload.get("ppt_summary", ""),
        "snippet": payload.get("snippet", ""),
        "date": payload.get("date", ""),
        "link": payload.get("link", ""),
        "top_image": payload.get("top_image", ""),
        "sources": payload.get("sources", []),
        "importance_score": payload.get("importance_score", 0),
        "keywords_found": keywords,
        "region": payload.get("region", "Global"),
        "full_contents": payload.get("full_contents", ""),
        "category": payload.get("category", "Tech News"),
        "source": payload.get("source", "Unknown"),
        "sentiment": payload.get("sentiment", "neutral"),
        "is_fresh": payload.get("is_fresh", True),
        "first_seen": payload.get("first_seen", ""),
        "source_count": payload.get("source_count", 1),
        "entities": payload.get("entities", {}),
        "profile": profile,
        "rejected_at": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "rejected_by": payload.get("rejected_by", "authorized kill switch"),
        "rejected_by_ip": authorized_ip,
    }
    store_before = []
    store_changed = False
    try:
        with not_interested_lock, briefing_lock:
            store_before = load_not_interested_store(request)
            already_rejected = is_already_rejected(title, store_before)
            latest = get_latest_briefing_file_for_profile(profile)
            briefing_before = _load_briefing_items(latest) if latest else []
            normalized_title = title.casefold()
            briefing_after = [
                item for item in briefing_before
                if str(item.get("title", "")).strip().casefold() != normalized_title
            ]
            store_after = list(store_before)
            if not already_rejected:
                store_after.append(entry)
                save_not_interested_store(store_after, request)
                store_changed = True
            try:
                if latest and len(briefing_after) != len(briefing_before):
                    _save_briefing_items(latest, briefing_after)
            except Exception as error:
                rollback_error = None
                if store_changed:
                    try:
                        save_not_interested_store(store_before, request)
                    except Exception as rollback:
                        rollback_error = rollback
                _paired_state_failure("reject", error, rollback_error)
    except HTTPException:
        raise
    except Exception as error:
        _paired_state_failure("reject", error)

    if already_rejected:
        return {
            "status": "exists",
            "message": "Already in Not Interested; shared briefing state reconciled.",
            "count": len(store_before),
            "briefing_removed": len(briefing_before) - len(briefing_after),
            "profile": profile,
        }

    try:
        save_training_vote(keywords, summary, "not_interested", title, profile)
    except Exception as error:
        raise HTTPException(
            status_code=500,
            detail={
                "status": "partial",
                "operation": "reject",
                "state": "content_committed_training_pending",
                "recoverable": True,
                "message": "The article was removed, but its learning vote could not be saved. Refresh before retrying.",
                "cause": str(error),
            },
        ) from error
    background_tasks.add_task(enqueue_bouncer_retrain, profile)

    return {
        "status": "success",
        "message": "Moved to Not Interested and removed from the shared briefing",
        "count": len(store_before) + 1,
        "briefing_removed": len(briefing_before) - len(briefing_after),
        "profile": profile,
        "retrain_scheduled": True,
    }


# ==========================================
# --- PERSONAL HIDDEN SIGNALS ---
# ==========================================
@app.get("/viewer/hidden")
def get_personal_hidden(request: Request):
    profile = get_profile_for_request(request)
    items = get_viewer_hidden_items(request, profile)
    return {
        "status": "success",
        "items": apply_learned_regions(items, profile),
        "count": len(items),
        "profile": profile,
        "scope": "current_viewer_only",
        "trains_bouncer": False,
    }


@app.post("/viewer/hidden")
def hide_for_current_viewer(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    viewer_key = get_viewer_key(get_client_ip(request))
    article_key = _article_identity(payload)
    if not article_key:
        raise HTTPException(status_code=400, detail="An article title or link is required.")

    with viewer_hidden_lock:
        store = load_viewer_hidden_store()
        viewer_store = store.setdefault(viewer_key, {})
        items = viewer_store.setdefault(profile, [])
        if not any(
            str(item.get("article_key") or _article_identity(item)) == article_key
            for item in items
        ):
            entry = dict(payload)
            entry["article_key"] = article_key
            entry["hidden_at"] = datetime.datetime.now().isoformat(timespec="seconds")
            entry["hidden_scope"] = "current_viewer_only"
            items.insert(0, entry)
            save_viewer_hidden_store(store)

    return {
        "status": "success",
        "message": "Hidden from your feed only.",
        "count": len(items),
        "profile": profile,
        "scope": "current_viewer_only",
        "trains_bouncer": False,
    }


@app.post("/viewer/hidden/restore")
def restore_for_current_viewer(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    viewer_key = get_viewer_key(get_client_ip(request))
    target_key = str(payload.get("article_key") or _article_identity(payload))
    if not target_key:
        raise HTTPException(status_code=400, detail="An article title or link is required.")

    with viewer_hidden_lock:
        store = load_viewer_hidden_store()
        viewer_store = store.setdefault(viewer_key, {})
        items = viewer_store.get(profile, [])
        remaining = [
            item for item in items
            if str(item.get("article_key") or _article_identity(item)) != target_key
        ]
        viewer_store[profile] = remaining
        save_viewer_hidden_store(store)

    return {
        "status": "success",
        "message": "Signal restored to your feed.",
        "count": len(remaining),
        "profile": profile,
        "scope": "current_viewer_only",
        "trains_bouncer": False,
    }


# ==========================================
# --- PERSONAL URL BRIEFINGS ---
# ==========================================
@app.get("/viewer/briefings")
def get_personal_url_briefings(request: Request):
    profile = get_profile_for_request(request)
    with viewer_briefing_lock:
        store = PERSONAL_BRIEFING_STORE.read()
        viewer_key, migrated = claim_legacy_private_bucket(store, request)
        if migrated:
            PERSONAL_BRIEFING_STORE.write(store)
        jobs = store.get(viewer_key, {}).get(profile, [])
    jobs = jobs if isinstance(jobs, list) else []
    return {
        "status": "success",
        "jobs": list(reversed(jobs)),
        "items": [
            job["article"]
            for job in reversed(jobs)
            if job.get("status") == "complete"
            and isinstance(job.get("article"), dict)
        ],
        "profile": profile,
        "scope": "current_viewer_only",
        "count": len(jobs),
        "active_count": sum(
            job.get("status") in {"queued", "processing"} for job in jobs
        ),
    }


@app.post("/viewer/briefings")
def create_personal_url_briefings(request: Request, payload: dict = Body(...)):
    profile = get_profile_for_request(request)
    client_ip = get_client_ip(request)
    viewer_key = get_private_viewer_key(request)
    fingerprint = str(payload.get("_tracking_fingerprint") or "unknown")
    raw_value = payload.get("urls", [])
    if isinstance(raw_value, list):
        candidates = [str(value).strip() for value in raw_value]
    else:
        candidates = re.findall(r"https?://[^\s,]+", str(raw_value or ""))
    candidates = [value.rstrip(".,;") for value in candidates if value.strip()]
    if not candidates:
        raise HTTPException(
            status_code=400,
            detail="Paste at least one HTTP or HTTPS news article URL.",
        )
    if len(candidates) > PERSONAL_BRIEFING_MAX_URLS:
        raise HTTPException(
            status_code=400,
            detail=(
                f"Submit no more than {PERSONAL_BRIEFING_MAX_URLS} URLs "
                "at one time."
            ),
        )

    valid_urls, invalid = [], []
    for candidate in candidates:
        try:
            canonical = assert_public_article_url(candidate)
            if canonical not in valid_urls:
                valid_urls.append(canonical)
        except ValueError as error:
            invalid.append({"url": candidate, "error": str(error)})

    accepted, duplicates = [], []
    with viewer_briefing_lock:
        store = PERSONAL_BRIEFING_STORE.read()
        viewer_key, _ = claim_legacy_private_bucket(store, request)
        viewer = store.setdefault(viewer_key, {})
        jobs = viewer.setdefault(profile, [])
        existing_by_url = {
            str(job.get("url")): job
            for job in jobs
            if job.get("url")
        }
        now = datetime.datetime.now().isoformat(timespec="seconds")
        for url in valid_urls:
            existing = existing_by_url.get(url)
            if existing:
                duplicates.append(
                    {
                        "url": url,
                        "job_id": existing.get("id"),
                        "status": existing.get("status"),
                        "article_id": (
                            existing.get("article", {}).get("id")
                            if isinstance(existing.get("article"), dict)
                            else None
                        ),
                    }
                )
                continue
            job = {
                "id": secrets.token_hex(10),
                "url": url,
                "status": "queued",
                "stage": "queued",
                "progress": 5,
                "message": "Waiting for a private briefing worker.",
                "created_at": now,
                "updated_at": now,
                "profile": profile,
                "scope": "current_viewer_only",
                "article": None,
                "error": None,
            }
            jobs.append(job)
            existing_by_url[url] = job
            accepted.append(job)
        PERSONAL_BRIEFING_STORE.write(store)

    dispatch_failures = []
    for job in accepted:
        try:
            personal_briefing_executor.submit(
                process_personal_briefing_job,
                viewer_key,
                profile,
                job["id"],
                job["url"],
            )
        except Exception as error:
            failure = {
                "status": "failed",
                "stage": "dispatch_failed",
                "progress": 100,
                "message": (
                    "The private briefing worker could not start. "
                    "Use Retry when the service is ready."
                ),
                "error": f"{type(error).__name__}: {error}"[:500],
            }
            update_personal_briefing_job(
                viewer_key,
                profile,
                job["id"],
                **failure,
            )
            # Keep the response consistent with the committed job record so
            # the UI never claims work started when dispatch actually failed.
            job.update(failure)
            dispatch_failures.append(
                {
                    "job_id": job["id"],
                    "url": job["url"],
                    "error": failure["error"],
                    "retryable": True,
                }
            )
            print(
                f"[PERSONAL:{profile}] Could not dispatch job {job['id']}: "
                f"{failure['error']}",
                flush=True,
            )
    activity_tracked = record_usage_best_effort(
        client_ip,
        profile,
        fingerprint,
        "personal_briefing_submit",
        json.dumps(
            {
                "accepted": len(accepted),
                "duplicates": len(duplicates),
                "invalid": len(invalid),
            }
        ),
    )
    return {
        "status": "partial" if dispatch_failures else "accepted",
        "accepted": accepted,
        "duplicates": duplicates,
        "invalid": invalid,
        "dispatch_failures": dispatch_failures,
        "activity_tracked": activity_tracked,
        "profile": profile,
        "scope": "current_viewer_only",
    }


@app.post("/viewer/briefings/{job_id}/retry")
def retry_personal_url_briefing(
    job_id: str,
    request: Request,
):
    profile = get_profile_for_request(request)
    viewer_key = get_private_viewer_key(request)
    transition = {"job": None, "queued": False}

    def updater(store):
        nonlocal viewer_key
        viewer_key, _ = claim_legacy_private_bucket(store, request)
        jobs = store.setdefault(viewer_key, {}).setdefault(profile, [])
        job = next(
            (candidate for candidate in jobs if candidate.get("id") == job_id),
            None,
        )
        if not job:
            return store
        transition["job"] = dict(job)
        if job.get("status") in {"queued", "processing"}:
            return store
        job.update(
            {
                "status": "queued",
                "stage": "queued",
                "progress": 5,
                "message": "Retry queued.",
                "error": None,
                "updated_at": datetime.datetime.now().isoformat(timespec="seconds"),
            }
        )
        transition["job"] = dict(job)
        transition["queued"] = True
        return store

    # The read/check/update is one critical section. Two simultaneous Retry
    # clicks can no longer submit the same URL to the worker twice.
    with viewer_briefing_lock:
        PERSONAL_BRIEFING_STORE.update(updater)

    job = transition["job"]
    if not job:
        raise HTTPException(status_code=404, detail="Private briefing job not found.")
    if not transition["queued"]:
        return {"status": "already_running", "job": job}
    try:
        personal_briefing_executor.submit(
            process_personal_briefing_job,
            viewer_key,
            profile,
            job_id,
            job.get("url"),
        )
    except Exception as error:
        failure_message = (
            "The private briefing worker could not start. "
            "Use Retry when the service is ready."
        )
        update_personal_briefing_job(
            viewer_key,
            profile,
            job_id,
            status="failed",
            stage="dispatch_failed",
            progress=100,
            message=failure_message,
            error=f"{type(error).__name__}: {error}"[:500],
        )
        raise HTTPException(
            status_code=503,
            detail={
                "status": "error",
                "operation": "personal_briefing_retry",
                "state": "dispatch_failed",
                "recoverable": True,
                "message": failure_message,
                "job_id": job_id,
            },
        ) from error
    return {"status": "queued", "job_id": job_id}


@app.post("/viewer/briefings/clear")
def clear_personal_url_briefings(request: Request, payload: dict = Body(default={})):
    """Clear finished private URL jobs without interrupting active work."""

    profile = get_profile_for_request(request)
    client_ip = get_client_ip(request)
    viewer_key = get_private_viewer_key(request)
    fingerprint = str(payload.get("_tracking_fingerprint") or "unknown")
    requested_scope = str(payload.get("scope") or "finished").strip().lower()
    if requested_scope not in {"finished", "failed", "complete"}:
        raise HTTPException(
            status_code=400,
            detail="Clear scope must be finished, failed, or complete.",
        )
    statuses = {
        "finished": {"complete", "failed"},
        "failed": {"failed"},
        "complete": {"complete"},
    }[requested_scope]
    removed = 0
    remaining = 0

    with viewer_briefing_lock:
        def updater(store):
            nonlocal removed, remaining, viewer_key
            viewer_key, _ = claim_legacy_private_bucket(store, request)
            viewer = store.setdefault(viewer_key, {})
            jobs = viewer.get(profile, [])
            jobs = jobs if isinstance(jobs, list) else []
            kept = [job for job in jobs if job.get("status") not in statuses]
            removed = len(jobs) - len(kept)
            remaining = len(kept)
            viewer[profile] = kept
            return store

        PERSONAL_BRIEFING_STORE.update(updater)

    activity_tracked = False
    if removed:
        activity_tracked = record_usage_best_effort(
            client_ip,
            profile,
            fingerprint,
            "personal_briefing_clear",
            json.dumps(
                {"removed": removed, "scope": requested_scope, "screen": "my_briefing"}
            ),
        )
    return {
        "status": "success",
        "removed": removed,
        "remaining": remaining,
        "active_jobs_preserved": True,
        "activity_tracked": activity_tracked,
        "profile": profile,
        "scope": "current_viewer_only",
    }


# ==========================================
# --- PERSONAL SAVED-FOR-LATER SIGNALS ---
# ==========================================
@app.get("/viewer/saved")
def get_personal_saved(request: Request):
    profile = get_profile_for_request(request)
    items = get_viewer_saved_items(request, profile)
    return {
        "status": "success",
        "items": apply_learned_regions(items, profile),
        "count": len(items),
        "profile": profile,
        "scope": "current_viewer_only",
        "affects_ranking": True,
        "follow_window_days": 30,
    }


@app.post("/viewer/saved")
def save_for_current_viewer(
    request: Request,
    payload: dict = Body(...),
    response: Response = None,
):
    profile = get_profile_for_request(request)
    client_ip = get_client_ip(request)
    viewer_key = get_private_viewer_key(request)
    fingerprint = str(payload.get("_tracking_fingerprint") or "unknown")
    article_payload = {
        key: value
        for key, value in payload.items()
        if key != "_tracking_fingerprint"
    }
    article_key = _article_identity(article_payload)
    if not article_key:
        raise HTTPException(status_code=400, detail="An article title or link is required.")

    changed = False
    with viewer_saved_lock:
        store = load_viewer_saved_store()
        viewer_key, migrated = claim_legacy_private_bucket(store, request)
        viewer_store = store.setdefault(viewer_key, {})
        items = viewer_store.setdefault(profile, [])
        existing = next(
            (
                item
                for item in items
                if str(item.get("article_key") or _article_identity(item))
                == article_key
            ),
            None,
        )
        if existing is None:
            entry = dict(article_payload)
            entry["article_key"] = article_key
            entry["saved_at"] = datetime.datetime.now().isoformat(timespec="seconds")
            entry["saved_scope"] = "current_viewer_only"
            items.insert(0, entry)
            changed = True
        if changed or migrated:
            save_viewer_saved_store(store)

    activity_tracked = False
    if changed:
        activity_detail = {
            key: article_payload.get(key)
            for key in (
                "title", "link", "url", "canonical_link", "source", "src",
                "category", "region", "keywords", "keywords_found",
                "matched_keywords", "matched_terms", "entities", "article_intent",
                "intent", "cluster_id",
            )
            if article_payload.get(key) not in (None, "", [])
        }
        activity_detail["screen"] = "saved_endpoint"
        activity_tracked = record_usage_best_effort(
            client_ip,
            profile,
            fingerprint,
            "save_for_later",
            json.dumps(activity_detail, ensure_ascii=False),
        )
        record_recommendation_best_effort(
            request,
            response or Response(),
            "save_for_later",
            activity_detail,
            event_id=str(payload.get("recommendation_event_id") or ""),
            occurred_at=str(payload.get("occurred_at") or ""),
        )

    return {
        "status": "success",
        "saved": True,
        "changed": changed,
        "activity_tracked": activity_tracked,
        "count": len(items),
        "profile": profile,
        "scope": "current_viewer_only",
        "affects_ranking": True,
        "follow_window_days": 30,
    }


@app.post("/viewer/saved/remove")
def remove_saved_for_current_viewer(
    request: Request,
    payload: dict = Body(...),
    response: Response = None,
):
    profile = get_profile_for_request(request)
    client_ip = get_client_ip(request)
    viewer_key = get_private_viewer_key(request)
    fingerprint = str(payload.get("_tracking_fingerprint") or "unknown")
    target_key = str(payload.get("article_key") or _article_identity(payload))
    if not target_key:
        raise HTTPException(status_code=400, detail="An article title or link is required.")

    with viewer_saved_lock:
        store = load_viewer_saved_store()
        viewer_key, migrated = claim_legacy_private_bucket(store, request)
        viewer_store = store.setdefault(viewer_key, {})
        items = viewer_store.get(profile, [])
        remaining = [
            item
            for item in items
            if str(item.get("article_key") or _article_identity(item))
            != target_key
        ]
        viewer_store[profile] = remaining
        changed = len(remaining) != len(items)
        if changed or migrated:
            save_viewer_saved_store(store)

    activity_tracked = False
    if changed:
        removal_detail = {
            "title": payload.get("title", ""),
            "link": payload.get("link") or payload.get("url", ""),
            "source": payload.get("source", ""),
            "screen": "saved_endpoint",
        }
        activity_tracked = record_usage_best_effort(
            client_ip,
            profile,
            fingerprint,
            "save_for_later_remove",
            json.dumps(
                removal_detail,
                ensure_ascii=False,
            ),
        )
        record_recommendation_best_effort(
            request,
            response or Response(),
            "save_for_later_remove",
            removal_detail,
            event_id=str(payload.get("recommendation_event_id") or ""),
            occurred_at=str(payload.get("occurred_at") or ""),
        )

    return {
        "status": "success",
        "saved": False,
        "changed": changed,
        "activity_tracked": activity_tracked,
        "count": len(remaining),
        "profile": profile,
        "scope": "current_viewer_only",
        "affects_ranking": True,
        "follow_window_days": 30,
    }


@app.get("/viewer/personalization")
def get_viewer_personalization(request: Request):
    profile = get_profile_for_request(request)
    viewer_key = get_viewer_key(get_client_ip(request))
    return {
        "status": "success",
        "profile": profile,
        "scope": "current_viewer_only",
        **PERSONALIZATION_SERVICE.summary(viewer_key, profile),
    }


@app.post("/viewer/personalization/reset")
def reset_viewer_personalization(request: Request):
    profile = get_profile_for_request(request)
    viewer_key = get_viewer_key(get_client_ip(request))
    removed = PERSONALIZATION_SERVICE.reset(viewer_key, profile)
    return {
        "status": "success",
        "removed_events": removed,
        "profile": profile,
        "scope": "current_viewer_only",
        "saved_signals_preserved": True,
    }


@app.get("/not-interested")
def get_not_interested(request: Request):
    profile = get_active_profile_name(request)
    with not_interested_lock:
        store = load_not_interested_store(request)
    return {"status": "success", "items": apply_learned_regions(store, profile), "count": len(store), "expiry_hours": NOT_INTERESTED_EXPIRY_HOURS, "profile": profile}


@app.post("/not-interested/restore")
def restore_from_not_interested(request: Request, background_tasks: BackgroundTasks, payload: dict = Body(...)):
    profile = get_active_profile_name(request)
    title = str(payload.get("title", "")).strip()
    if not title:
        raise HTTPException(status_code=400, detail="An article title is required.")

    store_before = []
    try:
        with not_interested_lock, briefing_lock:
            store_before = load_not_interested_store(request)
            article_to_restore = None
            remaining = []
            normalized_title = title.casefold()
            for item in store_before:
                if str(item.get("title", "")).strip().casefold() == normalized_title:
                    article_to_restore = dict(item)
                else:
                    remaining.append(item)

            if not article_to_restore:
                raise HTTPException(status_code=404, detail="Article not found in Not Interested")

            article_to_restore.pop("rejected_at", None)
            article_to_restore.pop("rejected_by", None)
            article_to_restore = apply_learned_region(article_to_restore, profile)
            latest = get_latest_briefing_file_for_profile(profile)
            if latest:
                briefing_before = _load_briefing_items(latest)
            else:
                history_dir = get_profile_history_dir(profile)
                timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
                latest = os.path.join(history_dir, f"briefing_{timestamp}.json")
                briefing_before = []
            briefing_after = list(briefing_before)
            briefing_restored = not any(
                str(item.get("title", "")).strip().casefold() == normalized_title
                for item in briefing_after
            )
            if briefing_restored:
                briefing_after.insert(0, article_to_restore)

            save_not_interested_store(remaining, request)
            try:
                if briefing_restored:
                    _save_briefing_items(latest, briefing_after)
            except Exception as error:
                rollback_error = None
                try:
                    save_not_interested_store(store_before, request)
                except Exception as rollback:
                    rollback_error = rollback
                _paired_state_failure("restore", error, rollback_error)
    except HTTPException:
        raise
    except Exception as error:
        _paired_state_failure("restore", error)

    summary = article_to_restore.get("master_summary", "")
    keywords = article_to_restore.get("keywords_found", [])
    try:
        save_training_vote(keywords, summary, "interested", article_to_restore.get("title", ""), profile)
    except Exception as error:
        raise HTTPException(
            status_code=500,
            detail={
                "status": "partial",
                "operation": "restore",
                "state": "content_committed_training_pending",
                "recoverable": True,
                "message": "The article was restored, but its learning vote could not be saved. Refresh before retrying.",
                "cause": str(error),
            },
        ) from error

    print(f"Restored: {title[:60]}. Counter-vote saved. Triggering retrain...")
    background_tasks.add_task(enqueue_bouncer_retrain, profile)
    return {
        "status": "success",
        "message": "Restored to the shared briefing",
        "article": article_to_restore,
        "count": len(remaining),
        "briefing_restored": briefing_restored,
        "profile": profile,
        "retrain_scheduled": True,
    }


# ==========================================
# --- USAGE TRACKING ENDPOINTS ---
# ==========================================
def record_usage_activity(ip, profile, fingerprint, action, detail=""):
    """Persist one authoritative activity event.

    Save/remove endpoints call this only after their JSON state actually
    changes. The public /track endpoint also uses it for ordinary UI events.
    Raw IP addresses are never written to the tracker.
    """

    if not action:
        return False

    team_owner = get_team_owner_for_ip(ip)
    device_id = get_device_id(ip, fingerprint or "unknown")
    today = get_today()
    try:
        event_detail = json.loads(detail) if detail else ""
    except (TypeError, ValueError):
        event_detail = detail

    with tracker_lock:
        tracker = load_tracker()

        if device_id not in tracker:
            tracker[device_id] = {
                "ip_hash": get_viewer_key(ip),
                "fingerprint": fingerprint or "unknown",
                "profile": profile,
                "owner": team_owner or "Unknown",
                "known_team_member": bool(team_owner),
                "display_name": get_viewer_profile(ip).get(
                    "display_name", team_owner or "Unknown"
                ),
                "first_seen": today,
                "last_seen": today,
                "activity": {},
            }

        device = tracker[device_id]
        device["last_seen"] = today
        device["ip_hash"] = get_viewer_key(ip)
        device.pop("ip", None)
        device["profile"] = profile
        device["owner"] = team_owner or "Unknown"
        device["known_team_member"] = bool(team_owner)
        device["display_name"] = get_viewer_profile(ip).get(
            "display_name",
            team_owner or device.get("display_name", "Unknown"),
        )
        device.setdefault("activity", {})
        if today not in device["activity"]:
            device["activity"][today] = get_empty_day()

        day = device["activity"][today]
        day.setdefault("action_counts", {})
        day.setdefault("events", [])
        supported_actions = {
            "page_load",
            "search",
            "article_click",
            "dossier_open",
            "dossier_dwell",
            "source_open",
            "why_this_story_open",
            "vote",
            "vote_interested",
            "vote_not_interested",
            "save_for_later",
            "save_for_later_remove",
            "export",
            "draft_export",
            "briefing_view",
            "heartbeat",
            "voc_feedback",
            "select",
            "batch_select",
            "approve",
            "hide_personal",
            "restore_personal_hidden",
            "remove_selected",
            "remove_approved",
            "add_source",
            "personal_briefing_submit",
            "personal_briefing_open",
            "personal_briefing_export",
            "personal_briefing_clear",
            "archive_import",
        }
        if action not in supported_actions:
            return False
        day["action_counts"][action] = (
            int(day["action_counts"].get(action, 0)) + 1
        )
        day["events"].append(
            {
                "timestamp": datetime.datetime.now().isoformat(
                    timespec="seconds"
                ),
                "action": action,
                "detail": event_detail,
            }
        )
        day["events"] = day["events"][-500:]
        if action == "page_load":
            day["page_loads"] = day.get("page_loads", 0) + 1
        elif action == "search":
            day.setdefault("searches", [])
            if detail and detail not in day["searches"]:
                day["searches"].append(detail)
        elif action in {"article_click", "dossier_open"}:
            day["articles_clicked"] = day.get("articles_clicked", 0) + 1
        elif action == "vote_interested":
            day["votes_interested"] = day.get("votes_interested", 0) + 1
        elif action == "vote_not_interested":
            day["votes_not_interested"] = day.get("votes_not_interested", 0) + 1
        elif action == "save_for_later":
            day["saved_for_later"] = day.get("saved_for_later", 0) + 1
        elif action == "save_for_later_remove":
            day["removed_from_saved"] = day.get("removed_from_saved", 0) + 1
        elif action in {"export", "draft_export"}:
            day.setdefault("exports", [])
            if detail and detail not in day["exports"]:
                day["exports"].append(detail)
        elif action == "briefing_view":
            day["briefing_views"] = day.get("briefing_views", 0) + 1
        elif action == "heartbeat":
            day["heartbeats"] = day.get("heartbeats", 0) + 1
        elif action == "voc_feedback":
            day.setdefault("voc_feedback", []).append(detail)
        elif action in {"select", "batch_select", "archive_import"}:
            increment = 1
            if action in {"batch_select", "archive_import"}:
                try:
                    payload = (
                        event_detail
                        if isinstance(event_detail, dict)
                        else {}
                    )
                    increment = max(1, int(payload.get("item_count", 1)))
                except (TypeError, ValueError):
                    increment = 1
            day["selections"] = day.get("selections", 0) + increment
        elif action == "approve":
            day["approvals"] = day.get("approvals", 0) + 1
        elif action == "hide_personal":
            day["personal_hides"] = day.get("personal_hides", 0) + 1
        elif action in {
            "restore_personal_hidden",
            "remove_selected",
            "remove_approved",
        }:
            day["workflow_removals"] = (
                day.get("workflow_removals", 0) + 1
            )
        elif action == "add_source":
            day["sources_added"] = day.get("sources_added", 0) + 1
        elif action == "vote":
            vote_value = str(detail or "").split(":", 1)[0].strip().lower()
            counter = (
                "votes_not_interested"
                if vote_value in {"down", "not_interested"}
                else "votes_interested"
            )
            day[counter] = day.get(counter, 0) + 1

        purge_old_entries(device)
        save_tracker(tracker)
    try:
        if action == "batch_select" and isinstance(event_detail, dict):
            selected_items = event_detail.get("items", [])
            selected_items = selected_items if isinstance(selected_items, list) else []
            for selected_item in selected_items[:100]:
                PERSONALIZATION_SERVICE.record_event(
                    get_viewer_key(ip),
                    profile,
                    "select",
                    selected_item if isinstance(selected_item, dict) else None,
                )
        else:
            PERSONALIZATION_SERVICE.record_event(
                get_viewer_key(ip),
                profile,
                action,
                event_detail if isinstance(event_detail, dict) else None,
            )
    except Exception as personalization_error:
        print(
            f"[PERSONALIZATION] Could not record {action}: "
            f"{personalization_error}",
            flush=True,
        )
    return True


@app.post("/track")
def track_activity(request: Request, response: Response, payload: dict = Body(...)):
    ip = get_client_ip(request)
    profile = get_active_profile_name(request)
    fingerprint = payload.get("fingerprint", "unknown")
    action = payload.get("action", "")
    detail = payload.get("detail", "")
    tracked = record_usage_activity(ip, profile, fingerprint, action, detail)
    try:
        parsed_detail = json.loads(detail) if isinstance(detail, str) else detail
    except (TypeError, ValueError):
        parsed_detail = {}
    parsed_detail = parsed_detail if isinstance(parsed_detail, dict) else {}
    recommendation = record_recommendation_best_effort(
        request,
        response,
        action,
        parsed_detail,
        event_id=str(
            payload.get("recommendation_event_id")
            or payload.get("event_id")
            or parsed_detail.get("event_id")
            or ""
        ),
        occurred_at=str(payload.get("occurred_at") or parsed_detail.get("occurred_at") or ""),
        active_ms=payload.get("active_ms") or parsed_detail.get("active_ms") or 0,
        visible_ratio=payload.get("visible_ratio") or parsed_detail.get("visible_ratio") or 0.0,
    )
    return {"status": "ok", "tracked": tracked, "recommendation": recommendation}


@app.get("/analytics/access")
def get_analytics_access(request: Request):
    ip = get_client_ip(request)
    allowed = is_analytics_allowed_ip(ip)
    viewer = get_viewer_profile(ip)

    return {
        "allowed": allowed,
        "ip": ip,
        "owner": get_team_owner_for_ip(ip) or "Unknown",
        "known_team_member": bool(get_team_owner_for_ip(ip)),
        "viewer": {
            "display_name": viewer.get("display_name", get_team_owner_for_ip(ip) or ""),
            "email": viewer.get("email", ""),
            "ip": ip,
            "ip_hash": get_viewer_key(ip),
        },
    }


@app.get("/viewer/profile")
def read_viewer_profile(request: Request):
    ip = get_client_ip(request)
    viewer = get_viewer_profile(ip)
    return {
        "status": "success",
        "display_name": viewer.get("display_name", get_team_owner_for_ip(ip) or ""),
        "email": viewer.get("email", ""),
        "ip": ip,
        "ip_hash": get_viewer_key(ip),
    }


@app.post("/viewer/profile")
def update_viewer_profile(request: Request, payload: dict = Body(...)):
    display_name = str(payload.get("display_name", "")).strip()
    email = str(payload.get("email", "")).strip()
    if len(display_name) < 2:
        raise HTTPException(status_code=400, detail="Display name must contain at least two characters.")
    ip = get_client_ip(request)
    viewer_key = get_viewer_key(ip)
    with tracker_lock:
        profiles = load_viewer_profiles()
        previous_name = str(profiles.get(viewer_key, {}).get("display_name", "")).strip()
        duplicate = next((profile for key, profile in profiles.items() if key != viewer_key and str(profile.get("display_name", "")).casefold() == display_name.casefold()), None)
        if duplicate:
            raise HTTPException(status_code=409, detail="That display name is already in use. Please choose another one.")
        profiles[viewer_key] = {
            "display_name": display_name,
            "email": email,
            "updated_at": datetime.datetime.now().isoformat(timespec="seconds"),
        }
        save_viewer_profiles(profiles)
        tracker = load_tracker()
        for device in tracker.values():
            if device.get("ip_hash") == viewer_key or device.get("ip") == ip:
                device["display_name"] = display_name
                device["owner"] = display_name
                device["ip_hash"] = viewer_key
                device.pop("ip", None)
        save_tracker(tracker)
        refresh_workflow_identity(viewer_key, previous_name, display_name)
    return {
        "status": "success",
        "display_name": display_name,
        "email": email,
        "ip": ip,
        "ip_hash": viewer_key,
    }


@app.get("/analytics")
def get_analytics(request: Request, key: str = Query(None)):
    ip = require_analytics_access(request, key)

    tracker = load_tracker()
    today = get_today()
    summary = []

    for device_id, device in tracker.items():
        activity = device.get("activity", {})
        total_loads = sum(d.get("page_loads", 0) for d in activity.values())
        total_searches = sum(len(d.get("searches", [])) for d in activity.values())
        total_clicks = sum(d.get("articles_clicked", 0) for d in activity.values())
        total_votes = sum(d.get("votes_interested", 0) + d.get("votes_not_interested", 0) for d in activity.values())
        total_exports = sum(len(d.get("exports", [])) for d in activity.values())
        total_heartbeats = sum(d.get("heartbeats", 0) for d in activity.values())
        total_voc = sum(len(d.get("voc_feedback", [])) for d in activity.values())
        action_counts = {}
        for day_data in activity.values():
            for action, count in day_data.get("action_counts", {}).items():
                action_counts[action] = (
                    action_counts.get(action, 0) + int(count)
                )
        today_data = activity.get(today, get_empty_day())
        active_days = len([d for d in activity.values() if d.get("page_loads", 0) > 0])

        engagement = (total_clicks * 3 + total_votes * 5 + total_searches * 4 + total_exports * 10 + total_heartbeats * 1)

        summary.append({
            "device_id": device_id,
            "ip": f"hashed:{str(device.get('ip_hash', 'unknown'))[:12]}",
            "ip_hash": device.get("ip_hash", "unknown"),
            "owner": device.get("display_name") or device.get("owner", "Unknown"),
            "display_name": device.get("display_name") or device.get("owner", "Unknown"),
            "known_team_member": bool(device.get("known_team_member")),
            "profile": device.get("profile", DEFAULT_PROFILE),
            "first_seen": device.get("first_seen", ""),
            "last_seen": device.get("last_seen", ""),
            "active_days": active_days,
            "today": today_data,
            "totals": {
                "page_loads": total_loads,
                "searches": total_searches,
                "articles_clicked": total_clicks,
                "votes": total_votes,
                "exports": total_exports,
                "minutes_approx": total_heartbeats,
                "voc_feedback": total_voc,
                "actions": action_counts,
            },
            "engagement_score": engagement,
            "daily": activity,
        })

    summary.sort(key=lambda x: x["engagement_score"])
    known_count = sum(1 for device in summary if device.get("known_team_member"))
    return {
        "status": "success",
        "device_count": len(summary),
        "known_team_member_count": known_count,
        "unknown_device_count": len(summary) - known_count,
        "date": today,
        "viewer": {
            "ip": ip,
            "owner": get_team_owner_for_ip(ip) or "Unknown",
        },
        "team_ip_map": TEAM_IP_MAP,
        "devices": summary,
    }


# ==========================================
# --- MANUAL CRAWL ENDPOINT (STREAMING) ---
# ==========================================
@app.get("/crawl")
async def crawl(
    request: Request,
    keywords: str = Query(None),
    from_date: str = Query(None),
    to_date: str = Query(None),
    target_sites: str = Query("All"),
    session_id: str = Query(None),
):
    profile = get_profile_for_request(request)
    sites_file = get_sites_file_for_profile(profile)
    effective_keywords = keywords or get_profile_config(profile)["keywords"]
    job_id = secrets.token_hex(8)
    output_file = os.path.join(ROOT_DIR, f"ui_results_{job_id}.json")
    cluster_file = os.path.join(ROOT_DIR, f"clustered_results_{job_id}.json")

    with scheduler_lock:
        blocked_by_scheduler = SCHEDULER_STATUS["is_active"]
        active_jobs[job_id] = {
            "status": "blocked" if blocked_by_scheduler else "queued",
            "keywords": effective_keywords,
            "started_at": datetime.datetime.now().isoformat(),
            "profile": profile,
        }

    def event_stream():
        live_item_prefix = "SENSE_STREAM_ITEM:"

        def sse_event(event_name, payload):
            return (
                f"event: {event_name}\n"
                f"data: {json.dumps(payload, ensure_ascii=False)}\n\n"
            )

        def live_article_key(article):
            return (
                str(article.get("link") or article.get("url") or "").strip()
                or str(article.get("title") or "").strip().lower()
            )

        def prepare_live_card(raw_article):
            article = dict(raw_article or {})
            title = str(article.get("title", "")).strip()
            summary = get_bouncer_summary_from_item(article)
            keywords_found = article.get("keywords_found", [])

            decision = bouncer_decision(
                title,
                summary,
                keywords_found,
                profile=profile,
            )

            if not decision["keep"]:
                log_dropped_article(
                    title,
                    summary,
                    keywords_found,
                    decision,
                    profile=profile,
                    article=raw_article,
                )
                return None, decision

            source_name = str(article.get("source") or "Unknown").strip() or "Unknown"
            article = attach_bouncer_metadata(article, decision)
            article["profile"] = profile
            article["bouncer_stage"] = "manual_live"
            article["master_summary"] = (
                article.get("master_summary")
                or article.get("summary")
                or article.get("snippet")
                or summary
                or title
            )
            article["ppt_summary"] = article.get("ppt_summary") or article["master_summary"]
            article["full_contents"] = (
                article.get("full_contents")
                or article.get("full_content")
                or article.get("master_summary")
                or ""
            )
            article["source_count"] = int(article.get("source_count") or 1)
            article["importance_score"] = int(article.get("importance_score") or 50)
            article["category"] = assign_category(
                title,
                article.get("master_summary", "") or article.get("snippet", ""),
            )

            if not article.get("sources"):
                article["sources"] = [
                    {
                        "name": source_name,
                        "link": article.get("link") or article.get("url") or "#",
                        "date": article.get("date", ""),
                    }
                ]

            article.update(apply_learned_region(article, profile))

            return article, decision

        if blocked_by_scheduler:
            yield f"data: {json.dumps({'type': 'error', 'message': 'The scheduled briefing is running now. Please start the deep scan again when it completes.'})}\n\n"
            return

        if not crawl_semaphore.acquire(blocking=False):
            yield f"data: {json.dumps({'type': 'error', 'message': 'Server is at capacity. Please wait a moment and try again.'})}\n\n"
            with scheduler_lock:
                active_jobs[job_id]["status"] = "error"
            return

        process = None
        streamed_count = 0
        live_articles = []
        live_seen_keys = set()
        live_dropped_count = 0
        live_low_priority_count = 0
        discovered_candidate_count = 0

        try:
            with scheduler_lock:
                active_jobs[job_id]["status"] = "running"
            capabilities = resolve_pipeline_capabilities()
            yield f"data: {json.dumps({'type': 'job_started', 'job_id': job_id, 'profile': profile})}\n\n"
            yield f"data: {json.dumps({'type': 'status', 'message': f'Using {profile} profile'})}\n\n"
            yield f"data: {json.dumps({'type': 'status', 'message': 'Pipeline preflight selected ' + capabilities['mode'] + '.'})}\n\n"
            yield f"data: {json.dumps({'type': 'status', 'message': 'Deploying Spider...'})}\n\n"

            cmd = [
                sys.executable, "-m", "scrapy", "crawl", "news_spider",
                "-a", f"keyword={effective_keywords}",
                "-a", f"from_date={from_date}",
                "-a", f"to_date={to_date}",
                "-a", f"target_sites={target_sites}",
                "-a", f"sites_file={sites_file}",
                "-a", f"discovery_only={'true' if capabilities['discovery_only'] else 'false'}",
                "-s", f"ROBOTSTXT_OBEY={'True' if SCRAPY_ROBOTSTXT_OBEY else 'False'}",
                "-s", "USER_AGENT=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
                "-s", "TWISTED_REACTOR=twisted.internet.asyncioreactor.AsyncioSelectorReactor",
                "-O", output_file,
            ]
            spider_cwd = str(NEWS_CRAWLER_DIR)

            process_env = os.environ.copy()
            process_env["SENSE_STREAM_ITEMS"] = "1"
            process_env["PYTHONIOENCODING"] = "utf-8"

            process = subprocess.Popen(
                cmd, cwd=spider_cwd,
                stdout=subprocess.PIPE, stderr=subprocess.STDOUT,
                text=True, bufsize=1, universal_newlines=True,
                encoding="utf-8", errors="replace",
                env=process_env,
                close_fds=CLOSE_FDS,
            )

            try:
                for line in process.stdout:
                    line = line.strip()
                    if not line: continue

                    if line.startswith(live_item_prefix):
                        try:
                            raw_article = json.loads(line[len(live_item_prefix):])
                            if capabilities["web_search"]:
                                discovered_candidate_count += 1
                                if discovered_candidate_count == 1 or discovered_candidate_count % 10 == 0:
                                    yield f"data: {json.dumps({'type': 'status', 'message': f'Discovered {discovered_candidate_count} matching article URL(s); Samsung extraction starts after discovery.'})}\n\n"
                                continue
                            key = live_article_key(raw_article)

                            if key and key in live_seen_keys:
                                continue

                            if key:
                                live_seen_keys.add(key)

                            live_card, decision = prepare_live_card(raw_article)

                            if live_card:
                                live_articles.append(live_card)
                                streamed_count += 1

                                if decision.get("decision") == "low_priority":
                                    live_low_priority_count += 1

                                if streamed_count == 1:
                                    yield f"data: {json.dumps({'type': 'status', 'message': 'First bouncer-approved card streamed while crawler is still running.'})}\n\n"
                                elif streamed_count % 10 == 0:
                                    yield f"data: {json.dumps({'type': 'status', 'message': f'Live stream: {streamed_count} cards approved so far.'})}\n\n"

                                yield sse_event("card", {"card": live_card})
                            else:
                                live_dropped_count += 1
                                print(
                                    f"[LIVE-BOUNCER:{profile}] Dropped: "
                                    f"{str(raw_article.get('title', ''))[:90]} "
                                    f"(score={decision.get('score')})",
                                    flush=True,
                                )

                        except Exception as e:
                            print(f"[LIVE-STREAM] Item parse/filter error: {e}", flush=True)

                        continue

                    sys.stdout.write(f"{line}\n")
                    sys.stdout.flush()
                    if "LOG:" in line:
                        yield f"data: {json.dumps({'type': 'status', 'message': line.split('LOG:', 1)[1].strip()})}\n\n"
                    elif "item_scraped_count" in line:
                        yield f"data: {json.dumps({'type': 'status', 'message': 'Gathering Intelligence...'})}\n\n"
                process.wait()
            finally:
                if process and process.stdout:
                    process.stdout.close()

            if capabilities["web_search"] and os.path.exists(output_file):
                try:
                    with open(output_file, "r", encoding="utf-8") as f:
                        raw_data = json.load(f)
                    yield f"data: {json.dumps({'type': 'status', 'message': f'Extracting {len(raw_data)} discovered article(s) through Samsung Web Search at the configured safe rate.'})}\n\n"
                    enriched_data = enrich_raw_articles(
                        raw_data,
                        effective_keywords,
                        profile,
                        use_web_search=True,
                        raise_on_service_failure=True,
                    )
                    yield f"data: {json.dumps({'type': 'status', 'message': 'Running AI Gatekeeper on extracted article text...'})}\n\n"
                    filtered_data, dropped_count, low_priority_count = run_bouncer_filter_on_items(
                        enriched_data, profile, "manual_enriched"
                    )
                    with open(output_file, "w", encoding="utf-8") as f:
                        json.dump(filtered_data, f, indent=4, ensure_ascii=False)
                    yield f"data: {json.dumps({'type': 'status', 'message': f'Samsung extraction and Gatekeeper complete. Kept {len(filtered_data)}; removed {dropped_count}; low priority kept {low_priority_count}.'})}\n\n"
                except WebSearchRuntimeFailure as e:
                    print(
                        f"[PIPELINE:{profile}] Samsung extraction failed after "
                        f"preflight: {e}. Restarting with full Scrapy extraction.",
                        flush=True,
                    )
                    yield f"data: {json.dumps({'type': 'status', 'message': 'Samsung Web Search became unavailable. Restarting this scan with full local Scrapy extraction; no discovered articles will be silently discarded.'})}\n\n"
                    fallback_cmd = [
                        (
                            "discovery_only=false"
                            if str(argument).startswith("discovery_only=")
                            else argument
                        )
                        for argument in cmd
                    ]
                    fallback_result = subprocess.run(
                        fallback_cmd,
                        cwd=spider_cwd,
                        env=process_env,
                        timeout=3600,
                        check=False,
                    )
                    if fallback_result.returncode != 0:
                        raise RuntimeError(
                            "Full Scrapy fallback exited with code "
                            f"{fallback_result.returncode}"
                        )
                    if not os.path.exists(output_file):
                        raise RuntimeError(
                            "Full Scrapy fallback completed without an output file"
                        )
                    with open(output_file, "r", encoding="utf-8") as f:
                        raw_data = json.load(f)
                    filtered_data, dropped_count, low_priority_count = (
                        run_bouncer_filter_on_items(
                            raw_data,
                            profile,
                            "manual_scrapy_fallback",
                        )
                    )
                    with open(output_file, "w", encoding="utf-8") as f:
                        json.dump(
                            filtered_data,
                            f,
                            indent=4,
                            ensure_ascii=False,
                        )
                    yield f"data: {json.dumps({'type': 'status', 'message': f'Full Scrapy fallback and Gatekeeper complete. Kept {len(filtered_data)}; removed {dropped_count}; low priority kept {low_priority_count}.'})}\n\n"
                except Exception as e:
                    print(f"[PIPELINE:{profile}] Samsung extraction failed: {e}", flush=True)
                    raise
            else:
                # Local legacy flow: bouncer the crawler-extracted text first,
                # then apply optional Web Search enrichment.
                yield f"data: {json.dumps({'type': 'status', 'message': 'Running AI Gatekeeper...'})}\n\n"

                if live_articles:
                    with open(output_file, "w", encoding="utf-8") as f:
                        json.dump(live_articles, f, indent=4, ensure_ascii=False)

                    yield f"data: {json.dumps({'type': 'status', 'message': f'Live gatekeeper complete. Streamed {len(live_articles)} cards while crawling. Removed {live_dropped_count}. Low priority kept: {live_low_priority_count}.'})}\n\n"

                elif os.path.exists(output_file):
                    try:
                        with open(output_file, "r", encoding="utf-8") as f:
                            raw_data = json.load(f)

                        filtered_data, dropped_count, low_priority_count = run_bouncer_filter_on_items(
                            raw_data, profile, "manual_raw"
                        )

                        with open(output_file, "w", encoding="utf-8") as f:
                            json.dump(filtered_data, f, indent=4, ensure_ascii=False)

                        yield f"data: {json.dumps({'type': 'status', 'message': f'Gatekeeper done. Removed {dropped_count} articles. Low priority kept: {low_priority_count}.'})}\n\n"
                        print(f"Bouncer complete [{profile}]. Dropped {dropped_count} articles.", flush=True)
                    except Exception as e:
                        print(f"Bouncer error, skipping filter: {e}", flush=True)

                if os.path.exists(output_file):
                    try:
                        with open(output_file, "r", encoding="utf-8") as f:
                            enrichable = json.load(f)
                        enrichable = enrich_raw_articles(
                            enrichable,
                            effective_keywords,
                            profile,
                            use_web_search=False,
                        )
                        with open(output_file, "w", encoding="utf-8") as f:
                            json.dump(enrichable, f, indent=4, ensure_ascii=False)
                    except Exception as e:
                        print(f"[PIPELINE:{profile}] Raw enrichment failed: {e}", flush=True)
                        if WEB_SEARCH_REQUIRE_SUCCESS:
                            raise

            # ==========================================
            # PHASE 1: STREAM CARDS IN REAL-TIME
            # ==========================================
            yield f"data: {json.dumps({'type': 'status', 'message': 'Activating Fusion Engine...'})}\n\n"

            try:
                from news_scrapper.semantic_clustering import MinimalSemanticEngine
                engine = MinimalSemanticEngine(
                    load_summarizer=not capabilities["chat"]
                )

                if capabilities["chat"]:
                    yield f"data: {json.dumps({'type': 'status', 'message': 'Clustering extracted articles without local BART summarization...'})}\n\n"
                    engine.fuse(job_id=job_id, fast_mode=True)
                    yield f"data: {json.dumps({'type': 'status', 'message': 'Semantic clustering complete. Samsung Chat summarization starts next.'})}\n\n"
                elif live_articles:
                    yield f"data: {json.dumps({'type': 'status', 'message': 'Optimizing streamed cards into clustered events...'})}\n\n"
                    engine.fuse_cluster(job_id=job_id, fast_mode=False)
                    yield f"data: {json.dumps({'type': 'status', 'message': 'Optimization complete.'})}\n\n"
                else:
                    for event in engine.fuse_stream(job_id=job_id, fast_mode=False):
                        event["category"] = assign_category(
                            event.get("title", ""),
                            event.get("master_summary", "") or event.get("snippet", ""),
                        )
                        event["profile"] = profile
                        event.update(apply_learned_region(event, profile))
                        yield sse_event("card", {"card": event})
                        streamed_count += 1

                    if streamed_count == 0:
                        yield f"data: {json.dumps({'type': 'status', 'message': 'No articles to process.'})}\n\n"
                    else:
                        yield f"data: {json.dumps({'type': 'status', 'message': f'All {streamed_count} cards streamed. Optimizing...'})}\n\n"

                    if streamed_count > 1:
                        yield f"data: {json.dumps({'type': 'status', 'message': 'Clustering duplicate stories...'})}\n\n"
                        engine.fuse_cluster(job_id=job_id, fast_mode=False)
                        yield f"data: {json.dumps({'type': 'status', 'message': 'Optimization complete.'})}\n\n"
                    else:
                        if os.path.exists(output_file):
                            with open(output_file, "r", encoding="utf-8") as f:
                                raw = json.load(f)
                            with open(cluster_file, "w", encoding="utf-8") as f:
                                json.dump(raw, f, indent=4, ensure_ascii=False)

            except Exception as e:
                print(f"Fusion streaming error: {e}", flush=True)
                yield f"data: {json.dumps({'type': 'status', 'message': 'Fusion Engine Error. Falling back...'})}\n\n"
                try:
                    fusion_fallback = subprocess.Popen(
                        [sys.executable, "-u", "-m", "news_scrapper.semantic_clustering", "--job-id", job_id],
                        cwd=str(PROJECT_ROOT),
                        stdout=subprocess.PIPE, stderr=subprocess.STDOUT,
                        text=True, encoding="utf-8", errors="replace",
                        close_fds=CLOSE_FDS,
                    )
                    try:
                        for line in fusion_fallback.stdout:
                            line = line.strip()
                            sys.stdout.write(f"{line}\n")
                            sys.stdout.flush()
                            if "FUSION ENGINE:" in line:
                                yield f"data: {json.dumps({'type': 'status', 'message': line.split('FUSION ENGINE:', 1)[1].strip()})}\n\n"
                        fusion_fallback.wait()
                    finally:
                        if fusion_fallback and fusion_fallback.stdout:
                            fusion_fallback.stdout.close()
                except Exception as fallback_err:
                    print(f"Fallback fusion also failed: {fallback_err}", flush=True)

            # Load final results
            results = []
            if os.path.exists(cluster_file):
                with open(cluster_file, "r", encoding="utf-8") as f:
                    results = json.load(f)
            elif os.path.exists(output_file):
                with open(output_file, "r", encoding="utf-8") as f:
                    results = json.load(f)

            if results:
                results, final_dropped_count, final_low_priority_count = run_bouncer_filter_on_items(
                    results, profile, "manual_final"
                )
                for r in results:
                    r["category"] = assign_category(
                        r.get("title", ""),
                        r.get("master_summary", "") or r.get("snippet", ""),
                    )
                    r["profile"] = profile
                    r.update(apply_learned_region(r, profile))
                results = enrich_final_articles(
                    results,
                    profile,
                    use_chat=capabilities["chat"],
                )

            # Archive
            if results:
                yield f"data: {json.dumps({'type': 'status', 'message': 'Archiving Intelligence...'})}\n\n"
                learner.log_search_data(effective_keywords, results)
                timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
                sid = session_id or "unknown"
                manual_path = os.path.join(get_profile_history_dir(profile), f"manual_{sid}_{timestamp}.json")
                with open(manual_path, "w", encoding="utf-8") as f:
                    json.dump(results, f, indent=4, ensure_ascii=False)
                print(f"Archived [{profile}]: manual_{sid}_{timestamp}.json", flush=True)

            with scheduler_lock:
                active_jobs[job_id]["status"] = "complete"

            if results and streamed_count == 0:
                for item in results:
                    yield sse_event("card", {"card": item})

            yield sse_event(
                "data",
                {
                    "results": results,
                    "job_id": job_id,
                    "reclustered": True,
                    "profile": profile,
                },
            )

        finally:
            if process and process.poll() is None:
                process.terminate()
            crawl_semaphore.release()
            with scheduler_lock:
                active_jobs[job_id]["status"] = "complete"
            threading.Timer(300, cleanup_job_files, args=[output_file, cluster_file]).start()

    return StreamingResponse(event_stream(), media_type="text/event-stream")
